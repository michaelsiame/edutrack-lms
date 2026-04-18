#!/usr/bin/env python3
"""
Module 1: Foundations of Computing and Mathematics
Five separate detailed PPTX presentations with custom diagrams and thorough notes.
Cybersecurity Certificate Program | EduTrack LMS | Zambia
"""

import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont
import math

# ---------------------------------------------------------------------------
# CONFIGURATION
# ---------------------------------------------------------------------------
BASE_DIR = "/Users/michael/Documents/Programming/edutrack-lms"
ASSETS_DIR = os.path.join(BASE_DIR, "module1_pptx_assets")
os.makedirs(ASSETS_DIR, exist_ok=True)

DARK_BLUE = RGBColor(0x1A, 0x23, 0x7E)
ACCENT_ORANGE = RGBColor(0xF4, 0x4C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1F, 0x29, 0x33)
LIGHT_GRAY = RGBColor(0xF5, 0xF7, 0xFA)

W, H = 1280, 720

# ---------------------------------------------------------------------------
# FONT HELPERS
# ---------------------------------------------------------------------------
def get_font(size):
    candidates = [
        "/System/Library/Fonts/Helvetica.ttc",
        "/System/Library/Fonts/HelveticaNeue.ttc",
        "/Library/Fonts/Arial.ttf",
        "/System/Library/Fonts/SFNS.ttf",
    ]
    for c in candidates:
        if os.path.exists(c):
            return ImageFont.truetype(c, size)
    return ImageFont.load_default()

FONT_TITLE = get_font(36)
FONT_BIG = get_font(28)
FONT_BODY = get_font(22)
FONT_NOTE = get_font(18)
FONT_SMALL = get_font(16)

def save_img(img, name):
    img.save(os.path.join(ASSETS_DIR, name))

# ---------------------------------------------------------------------------
# DIAGRAM GENERATION
# ---------------------------------------------------------------------------
def generate_all_diagrams():
    print("Generating educational diagrams...")
    
    # 1. IPOS Cycle
    img = Image.new("RGB", (W, H), "#F5F7FA")
    d = ImageDraw.Draw(img)
    d.text((W//2, 30), "The IPOS Model and Security Relevance", font=FONT_TITLE, fill="#1A237E", anchor="mm")
    cx, cy = 640, 360
    r_main = 180
    stages = [("INPUT", 270, "#1A237E"), ("PROCESS", 0, "#F44C00"), ("OUTPUT", 90, "#1A237E"), ("STORAGE", 180, "#F44C00")]
    for i, (label, angle, col) in enumerate(stages):
        rad = math.radians(angle)
        x = cx + int(r_main * math.cos(rad))
        y = cy + int(r_main * math.sin(rad))
        r = 70
        d.ellipse([x-r, y-r, x+r, y+r], fill=col, outline="white", width=5)
        d.text((x, y), label, font=FONT_BODY, fill="white", anchor="mm")
        next_angle = stages[(i+1)%4][1]
        rad2 = math.radians(next_angle)
        x2 = cx + int(r_main * math.cos(rad2))
        y2 = cy + int(r_main * math.sin(rad2))
        d.arc([cx-r_main-30, cy-r_main-30, cx+r_main+30, cy+r_main+30], angle-10, next_angle+10, fill="#333", width=4)
        d.polygon([(x2+r-5, y2), (x2+r-15, y2-8), (x2+r-15, y2+8)], fill="#333")
    d.text((W//2, 620), "Each stage is an attack surface that cybersecurity professionals must protect", font=FONT_NOTE, fill="#333", anchor="mm")
    save_img(img, "ipos_cycle.png")
    
    # 2. Hardware Components
    img = Image.new("RGB", (W, H), "#F5F7FA")
    d = ImageDraw.Draw(img)
    d.text((W//2, 25), "Computer Hardware Components and Security Relevance", font=FONT_TITLE, fill="#1A237E", anchor="mm")
    boxes = [
        (80, 90, "CENTRAL PROCESSING UNIT", "The brain of the computer", "Executes instructions\nCryptographic operations\nVulnerable to side-channel attacks", "#1A237E"),
        (660, 90, "RANDOM ACCESS MEMORY", "Temporary working memory", "Fast but loses data when off\nFileless malware hides here\nCold boot attacks possible", "#F44C00"),
        (80, 390, "STORAGE DEVICES", "HDD, SSD, USB Flash", "Permanent data retention\nEncryption at rest protects\nForensic evidence lives here", "#1A237E"),
        (660, 390, "NETWORK INTERFACE CARD", "Connects to network", "Unique MAC address\nMAC spoofing bypasses controls\nIdentifies devices on network", "#F44C00"),
    ]
    for x, y, title, subtitle, body, col in boxes:
        d.rounded_rectangle([x, y, x+540, y+280], radius=12, fill=col, outline="white", width=4)
        d.text((x+270, y+30), title, font=FONT_BODY, fill="white", anchor="mm")
        d.text((x+270, y+65), subtitle, font=FONT_NOTE, fill="#FFD180", anchor="mm")
        lines = body.split("\n")
        for li, line in enumerate(lines):
            d.text((x+20, y+100+li*28), "- " + line, font=FONT_SMALL, fill="white", anchor="lm")
    save_img(img, "hardware_components.png")
    
    # 3. Von Neumann
    img = Image.new("RGB", (W, H), "#F5F7FA")
    d = ImageDraw.Draw(img)
    d.text((W//2, 25), "Von Neumann Architecture", font=FONT_TITLE, fill="#1A237E", anchor="mm")
    def draw_block(x, y, w, h, title, body, col):
        d.rounded_rectangle([x, y, x+w, y+h], radius=8, fill=col, outline="white", width=3)
        d.text((x+w//2, y+25), title, font=FONT_BODY, fill="white", anchor="mm")
        for i, line in enumerate(body.split("\n")):
            d.text((x+w//2, y+55+i*20), line, font=FONT_SMALL, fill="#EEE", anchor="mm")
    def arrow(d, x1, y1, x2, y2):
        d.line([(x1, y1), (x2, y2)], fill="#333", width=3)
        angle = math.atan2(y2-y1, x2-x1)
        arr_len = 12
        d.polygon([
            (x2, y2),
            (x2 - arr_len*math.cos(angle - 0.5), y2 - arr_len*math.sin(angle - 0.5)),
            (x2 - arr_len*math.cos(angle + 0.5), y2 - arr_len*math.sin(angle + 0.5))
        ], fill="#333")
    draw_block(540, 120, 200, 90, "CPU", "Arithmetic Logic Unit\nControl Unit", "#1A237E")
    draw_block(540, 260, 200, 70, "MEMORY", "Stores data and instructions", "#F44C00")
    draw_block(540, 400, 200, 70, "STORAGE", "Long-term retention", "#1A237E")
    draw_block(120, 260, 140, 70, "INPUT", "Keyboard, Mouse", "#555")
    draw_block(1020, 260, 140, 70, "OUTPUT", "Monitor, Printer", "#555")
    arrow(d, 260, 295, 540, 165)
    arrow(d, 640, 210, 640, 260)
    arrow(d, 540, 295, 260, 295)
    arrow(d, 740, 295, 1020, 295)
    arrow(d, 640, 330, 640, 400)
    arrow(d, 640, 400, 640, 330)
    d.text((W//2, 530), "Critical Security Implication:", font=FONT_BIG, fill="#1A237E", anchor="mm")
    d.text((W//2, 570), "Because data and instructions share memory, buffer overflow attacks can overwrite code.", font=FONT_BODY, fill="#333", anchor="mm")
    d.text((W//2, 650), "Defenses: DEP / NX bit and Address Space Layout Randomization (ASLR)", font=FONT_NOTE, fill="#F44C00", anchor="mm")
    save_img(img, "von_neumann.png")
    
    # 4. Windows FS
    img = Image.new("RGB", (W, H), "#F5F7FA")
    d = ImageDraw.Draw(img)
    d.text((W//2, 20), "Windows File System Hierarchy", font=FONT_TITLE, fill="#1A237E", anchor="mm")
    tree_lines = [
        ("C:\\", 80, 80, "#1A237E", 0),
        ("├── Windows\\", 120, 120, "#555", 0),
        ("│   └── System32", 160, 150, "#F44C00", 1),
        ("├── Program Files\\", 120, 180, "#555", 0),
        ("├── Users\\", 120, 210, "#555", 0),
        ("│   └── Chanda\\", 160, 240, "#1A237E", 0),
        ("│       ├── Desktop\\", 200, 270, "#555", 0),
        ("│       ├── Documents\\", 200, 300, "#555", 0),
        ("│       ├── Downloads\\", 200, 330, "#F44C00", 0),
        ("│       └── AppData\\", 200, 360, "#F44C00", 1),
        ("└── Temp\\", 120, 400, "#F44C00", 0),
    ]
    for text, x, y, col, warn in tree_lines:
        d.text((x, y), text, font=FONT_BODY, fill=col)
        if warn:
            d.text((x+250, y), "HIGH RISK", font=FONT_SMALL, fill="#F44C00")
    notes = [
        ("System32: Core OS files. Malware targets for DLL hijacking.", 700, 120),
        ("AppData: Hidden user data. Top location for malware persistence.", 700, 240),
        ("Downloads: Common malware delivery point.", 700, 330),
        ("Temp: Temporary files. Check during incident response.", 700, 400),
    ]
    for text, x, y in notes:
        d.text((x, y), text, font=FONT_SMALL, fill="#333")
    save_img(img, "windows_fs.png")
    
    # 5. Linux FS
    img = Image.new("RGB", (W, H), "#F5F7FA")
    d = ImageDraw.Draw(img)
    d.text((W//2, 20), "Linux File System Hierarchy", font=FONT_TITLE, fill="#1A237E", anchor="mm")
    tree_lines = [
        ("/", 80, 80, "#1A237E"),
        ("├── /bin/       Essential system commands", 120, 115, "#555"),
        ("├── /etc/       Configuration files", 120, 145, "#F44C00"),
        ("├── /home/      User home directories", 120, 175, "#555"),
        ("│   └── /home/mulenga/", 160, 205, "#1A237E"),
        ("├── /root/      Root admin home", 120, 235, "#1A237E"),
        ("├── /tmp/       Temporary files", 120, 265, "#F44C00"),
        ("├── /usr/       User programs", 120, 295, "#555"),
        ("└── /var/log/   System log files", 120, 325, "#F44C00"),
    ]
    for text, x, y, col in tree_lines:
        d.text((x, y), text, font=FONT_BODY, fill=col)
    notes = [
        ("/etc/ contains passwd, shadow, and hosts files.", 650, 145),
        ("/var/log/ is the security goldmine for investigations.", 650, 265),
        ("/tmp/ is world-writable: malware favorite drop zone.", 650, 295),
    ]
    for text, x, y in notes:
        d.text((x, y), text, font=FONT_SMALL, fill="#333")
    save_img(img, "linux_fs.png")
    
    # 6. Linux Permissions
    img = Image.new("RGB", (W, H), "#F5F7FA")
    d = ImageDraw.Draw(img)
    d.text((W//2, 20), "Linux File Permissions Explained", font=FONT_TITLE, fill="#1A237E", anchor="mm")
    d.text((60, 70), "Permission string example:", font=FONT_BODY, fill="#333")
    d.text((60, 105), "-rwxr-xr--  1  mulenga  staff  1024  Mar 2026  script.sh", font=FONT_BIG, fill="#1A237E")
    parts = [("-", 60, 145, "File type"), ("rwx", 100, 145, "Owner"), ("r-x", 175, 145, "Group"), ("r--", 245, 145, "Others")]
    for text, x, y, label in parts:
        w = d.textlength(text, font=FONT_BIG)
        d.text((x, y-5), text, font=FONT_BIG, fill="#F44C00")
        d.line([(x, y+30), (x+w, y+30)], fill="#F44C00", width=3)
        d.text((x+w//2, y+40), label, font=FONT_SMALL, fill="#333", anchor="mt")
    table_y = 210
    row_h = 42
    cols = [60, 280, 430, 540, 850]
    headers = ["Permission", "Symbol", "Value", "Meaning", "Security Use"]
    rows = [
        ["Read", "r", "4", "View contents", "List directory / read file"],
        ["Write", "w", "2", "Modify contents", "Edit file / add delete files"],
        ["Execute", "x", "1", "Run as program", "Run scripts / enter directory"],
        ["No permission", "-", "0", "Denied", "Restrict sensitive data"],
    ]
    for i, h in enumerate(headers):
        d.rectangle([cols[i], table_y, cols[i+1]-5 if i+1 < len(cols) else W-60, table_y+row_h], fill="#1A237E", outline="white")
        d.text((cols[i]+10, table_y+row_h//2), h, font=FONT_SMALL, fill="white", anchor="lm")
    for ri, row in enumerate(rows):
        y = table_y + (ri+1)*row_h
        bg = "#FFFFFF" if ri % 2 == 0 else "#E8EAF6"
        for ci, cell in enumerate(row):
            w_end = cols[ci+1]-5 if ci+1 < len(cols) else W-60
            d.rectangle([cols[ci], y, w_end, y+row_h], fill=bg, outline="#CCC")
            d.text((cols[ci]+10, y+row_h//2), cell, font=FONT_SMALL, fill="#333", anchor="lm")
    d.text((60, 470), "Example: chmod 755 script.sh", font=FONT_BODY, fill="#1A237E")
    d.text((60, 510), "Owner=7 (rwx) | Group=5 (r-x) | Others=5 (r-x)", font=FONT_NOTE, fill="#333")
    d.text((60, 560), "Best practice: chmod 600 ~/.ssh/id_rsa (owner only)", font=FONT_SMALL, fill="#F44C00")
    save_img(img, "linux_permissions.png")
    
    # 7. Binary Conversion
    img = Image.new("RGB", (W, H), "#F5F7FA")
    d = ImageDraw.Draw(img)
    d.text((W//2, 20), "Decimal to Binary: Convert 45 to Binary", font=FONT_TITLE, fill="#1A237E", anchor="mm")
    steps = [
        "Step 1:  45 / 2 = 22  remainder 1",
        "Step 2:  22 / 2 = 11  remainder 0",
        "Step 3:  11 / 2 =  5  remainder 1",
        "Step 4:   5 / 2 =  2  remainder 1",
        "Step 5:   2 / 2 =  1  remainder 0",
        "Step 6:   1 / 2 =  0  remainder 1",
        "",
        "Read remainders bottom to top:  1 0 1 1 0 1",
        "Therefore: 45 decimal = 101101 binary",
    ]
    for i, step in enumerate(steps):
        color = "#1A237E" if "Therefore" in step else "#333"
        size = FONT_BODY if "Therefore" in step else FONT_NOTE
        d.text((80, 70+i*32), step, font=size, fill=color)
    d.text((80, 360), "Why this matters:", font=FONT_BIG, fill="#F44C00")
    d.text((80, 400), "IPv4 addresses use 32 bits divided into four 8-bit octets.", font=FONT_NOTE, fill="#333")
    d.text((80, 435), "Each octet ranges 0-255 decimal, which is 00000000 to 11111111 binary.", font=FONT_NOTE, fill="#333")
    d.text((80, 470), "Binary knowledge is essential for subnetting and network security design.", font=FONT_NOTE, fill="#333")
    save_img(img, "07_binary_conversion.png")
    
    # 8. Binary to Decimal
    img = Image.new("RGB", (W, H), "#F5F7FA")
    d = ImageDraw.Draw(img)
    d.text((W//2, 20), "Binary to Decimal: Convert 10101100 to Decimal", font=FONT_TITLE, fill="#1A237E", anchor="mm")
    headers = ["128", "64", "32", "16", "8", "4", "2", "1"]
    bits = ["1", "0", "1", "0", "1", "1", "0", "0"]
    start_x = 100
    cell_w = 130
    y1, y2 = 100, 160
    for i, (h, b) in enumerate(zip(headers, bits)):
        x = start_x + i*cell_w
        d.rectangle([x, y1, x+cell_w-5, y2], fill="#1A237E", outline="white")
        d.text((x+cell_w//2, y1+25), h, font=FONT_BODY, fill="white", anchor="mm")
        d.rectangle([x, y2, x+cell_w-5, y2+60], fill="#E8EAF6" if i%2==0 else "#FFFFFF", outline="#CCC")
        d.text((x+cell_w//2, y2+30), b, font=FONT_BIG, fill="#1A237E", anchor="mm")
    d.text((80, 260), "= (1 x 128) + (0 x 64) + (1 x 32) + (0 x 16) + (1 x 8) + (1 x 4) + (0 x 2) + (0 x 1)", font=FONT_BODY, fill="#333")
    d.text((80, 300), "= 128 + 0 + 32 + 0 + 8 + 4 + 0 + 0", font=FONT_BODY, fill="#333")
    d.text((80, 340), "= 172", font=FONT_BIG, fill="#1A237E")
    d.text((80, 390), "Therefore: 10101100 binary = 172 decimal", font=FONT_BODY, fill="#1A237E")
    d.text((80, 460), "This octet appears in IP addresses like 172.16.0.1", font=FONT_NOTE, fill="#333")
    save_img(img, "08_binary_to_decimal.png")
    
    # 9. Hex Conversion
    img = Image.new("RGB", (W, H), "#F5F7FA")
    d = ImageDraw.Draw(img)
    d.text((W//2, 20), "Binary to Hexadecimal: The 4-Bit Grouping Shortcut", font=FONT_TITLE, fill="#1A237E", anchor="mm")
    d.text((80, 80), "Binary:    1010    1100", font=FONT_BIG, fill="#333")
    d.text((80, 120), "Groups:    |       |", font=FONT_BODY, fill="#555")
    d.text((80, 160), "Hex:         A       C", font=FONT_BIG, fill="#1A237E")
    d.text((80, 210), "Therefore: 10101100 binary = AC hex = 172 decimal", font=FONT_BODY, fill="#1A237E")
    d.text((80, 280), "Reference Table:", font=FONT_BIG, fill="#333")
    refs = [
        ("Decimal", "0", "5", "9", "10", "15", "16", "255"),
        ("Binary", "0000", "0101", "1001", "1010", "1111", "00010000", "11111111"),
        ("Hex", "0", "5", "9", "A", "F", "10", "FF"),
    ]
    tx, ty = 80, 330
    cw = 140
    for ri, row in enumerate(refs):
        for ci, cell in enumerate(row):
            bg = "#1A237E" if ci == 0 else ("#E8EAF6" if ri%2==0 else "#FFFFFF")
            d.rectangle([tx+ci*cw, ty+ri*45, tx+(ci+1)*cw-5, ty+(ri+1)*45], fill=bg, outline="#CCC")
            d.text((tx+ci*cw+10, ty+ri*45+22), cell, font=FONT_SMALL, fill="white" if ci==0 else "#333", anchor="lm")
    d.text((80, 500), "Why hex matters:", font=FONT_BIG, fill="#F44C00")
    d.text((80, 540), "- MAC addresses use hex pairs", font=FONT_NOTE, fill="#333")
    d.text((80, 575), "- SHA-256 hashes are displayed as hex strings", font=FONT_NOTE, fill="#333")
    d.text((80, 610), "- Memory addresses in forensics use hex notation", font=FONT_NOTE, fill="#333")
    save_img(img, "09_hex_conversion.png")
    
    # 10. Boolean Algebra
    img = Image.new("RGB", (W, H), "#F5F7FA")
    d = ImageDraw.Draw(img)
    d.text((W//2, 20), "Boolean Algebra and Logic Operations", font=FONT_TITLE, fill="#1A237E", anchor="mm")
    ops = [
        ("AND", "Both inputs must be true", "1 AND 1 = 1", "#1A237E"),
        ("OR", "At least one input true", "0 OR 1 = 1", "#F44C00"),
        ("NOT", "Inverts the input value", "NOT 1 = 0", "#1A237E"),
        ("XOR", "One or other, not both", "1 XOR 1 = 0", "#F44C00"),
    ]
    for i, (name, desc, example, col) in enumerate(ops):
        x = 60 + (i % 2) * 620
        y = 70 + (i // 2) * 220
        d.rounded_rectangle([x, y, x+580, y+190], radius=10, fill=col, outline="white", width=3)
        d.text((x+290, y+30), name, font=FONT_BIG, fill="white", anchor="mm")
        d.text((x+290, y+70), desc, font=FONT_NOTE, fill="#FFD180", anchor="mm")
        d.text((x+290, y+110), "Example: " + example, font=FONT_BODY, fill="white", anchor="mm")
    d.text((W//2, 510), "Security Application:", font=FONT_BIG, fill="#1A237E", anchor="mm")
    d.text((W//2, 550), "Multi-factor authentication uses AND logic:", font=FONT_BODY, fill="#333", anchor="mm")
    d.text((W//2, 590), "Access granted ONLY IF password correct AND one-time password valid.", font=FONT_NOTE, fill="#333", anchor="mm")
    save_img(img, "10_boolean_algebra.png")
    
    # 11. XOR Encryption
    img = Image.new("RGB", (W, H), "#F5F7FA")
    d = ImageDraw.Draw(img)
    d.text((W//2, 20), "XOR: The Mathematical Foundation of Encryption", font=FONT_TITLE, fill="#1A237E", anchor="mm")
    d.text((80, 80), "The special property of XOR:", font=FONT_BIG, fill="#333")
    d.text((80, 120), "If you XOR a value with a key, then XOR the result with the SAME key,", font=FONT_BODY, fill="#333")
    d.text((80, 155), "you get back the original value.", font=FONT_BODY, fill="#333")
    d.text((80, 210), "Example:", font=FONT_BIG, fill="#1A237E")
    d.text((80, 255), "  Plaintext:     1010", font=FONT_BODY, fill="#333")
    d.text((80, 290), "  Key:           1100", font=FONT_BODY, fill="#333")
    d.text((80, 325), "  Ciphertext:    0110   (1010 XOR 1100)", font=FONT_BODY, fill="#F44C00")
    d.text((80, 360), "  Decrypt:       1010   (0110 XOR 1100)  -> Original back!", font=FONT_BODY, fill="#1A237E")
    d.text((80, 430), "Real-world applications:", font=FONT_BIG, fill="#F44C00")
    d.text((80, 475), "- One-Time Pad: only mathematically unbreakable encryption", font=FONT_NOTE, fill="#333")
    d.text((80, 510), "- Stream ciphers like ChaCha20 protect TLS web traffic", font=FONT_NOTE, fill="#333")
    d.text((80, 545), "- AES uses XOR in every encryption round", font=FONT_NOTE, fill="#333")
    d.text((80, 580), "- Malware authors use simple XOR to hide strings from antivirus", font=FONT_NOTE, fill="#333")
    save_img(img, "11_xor_encryption.png")
    
    # 12. Network Types
    img = Image.new("RGB", (W, H), "#F5F7FA")
    d = ImageDraw.Draw(img)
    d.text((W//2, 20), "Network Types and Security Requirements", font=FONT_TITLE, fill="#1A237E", anchor="mm")
    nets = [
        ("LAN", "Local Area Network", "Office, School lab", "Segmentation limits breach spread", "#1A237E"),
        ("WAN", "Wide Area Network", "National backbone", "Encryption prevents interception", "#F44C00"),
        ("WLAN", "Wireless LAN", "Home WiFi, Cafes", "Vulnerable to eavesdropping", "#1A237E"),
        ("VPN", "Virtual Private Network", "Remote work tunnel", "Misconfigurations cause breaches", "#F44C00"),
    ]
    for i, (abbr, full, scope, risk, col) in enumerate(nets):
        x = 60 + (i % 2) * 620
        y = 70 + (i // 2) * 220
        d.rounded_rectangle([x, y, x+580, y+190], radius=10, fill=col, outline="white", width=3)
        d.text((x+290, y+30), abbr + " - " + full, font=FONT_BODY, fill="white", anchor="mm")
        d.text((x+290, y+70), "Scope: " + scope, font=FONT_NOTE, fill="#FFD180", anchor="mm")
        d.text((x+290, y+120), "Security: " + risk, font=FONT_SMALL, fill="white", anchor="mm")
    save_img(img, "12_network_types.png")
    
    # 13. OSI Stack
    img = Image.new("RGB", (W, H), "#F5F7FA")
    d = ImageDraw.Draw(img)
    d.text((W//2, 15), "The OSI Model - 7 Layers", font=FONT_TITLE, fill="#1A237E", anchor="mm")
    layers = [
        ("7. APPLICATION", "HTTP, HTTPS, DNS, FTP, SMTP", "#1A237E", "User interaction"),
        ("6. PRESENTATION", "SSL, TLS, JPEG, ASCII", "#3949AB", "Encryption and format"),
        ("5. SESSION", "NetBIOS, RPC", "#1A237E", "Connection management"),
        ("4. TRANSPORT", "TCP, UDP", "#3949AB", "Delivery and ports"),
        ("3. NETWORK", "IP, ICMP, ARP", "#1A237E", "Routing and addressing"),
        ("2. DATA LINK", "Ethernet, WiFi", "#3949AB", "MAC addressing"),
        ("1. PHYSICAL", "Cables, radio, fiber", "#1A237E", "Bit transmission"),
    ]
    bw, bh = 700, 60
    bx = 80
    by_start = 80
    for i, (name, protos, col, desc) in enumerate(layers):
        y = by_start + i*(bh+8)
        d.rounded_rectangle([bx, y, bx+bw, y+bh], radius=6, fill=col, outline="white", width=2)
        d.text((bx+20, y+bh//2), name, font=FONT_BODY, fill="white", anchor="lm")
        d.text((bx+bw//2, y+bh//2), protos, font=FONT_SMALL, fill="#FFD180", anchor="mm")
        d.text((bx+bw-20, y+bh//2), desc, font=FONT_SMALL, fill="white", anchor="rm")
    d.text((bx+bw//2, by_start+7*(bh+8)+10), "Memory aid: All People Seem To Need Data Processing", font=FONT_NOTE, fill="#333", anchor="mm")
    d.text((bx+bw+60, 80), "Data Unit:", font=FONT_BIG, fill="#1A237E")
    pdus = ["Data", "Data", "Data", "Segment", "Packet", "Frame", "Bit"]
    for i, pdu in enumerate(pdus):
        y = by_start + i*(bh+8)
        d.text((bx+bw+60, y+bh//2), pdu, font=FONT_BODY, fill="#333", anchor="lm")
    save_img(img, "13_osi_stack.png")
    
    # 14. OSI vs TCP/IP
    img = Image.new("RGB", (W, H), "#F5F7FA")
    d = ImageDraw.Draw(img)
    d.text((W//2, 15), "OSI Model versus TCP/IP Model", font=FONT_TITLE, fill="#1A237E", anchor="mm")
    osi_layers = [
        ("Application", "HTTP, DNS, FTP"),
        ("Presentation", "SSL/TLS"),
        ("Session", "NetBIOS, RPC"),
        ("Transport", "TCP, UDP"),
        ("Network", "IP, ICMP, ARP"),
        ("Data Link", "Ethernet, WiFi"),
        ("Physical", "Cables, Radio"),
    ]
    tcp_map = [
        (0, 0, 3),
        (3, 1, 1),
        (4, 2, 1),
        (5, 3, 2),
    ]
    bx_osi, bx_tcp = 100, 760
    by = 70
    h = 55
    w = 420
    d.text((bx_osi+w//2, by-30), "OSI Model (7 Layers)", font=FONT_BIG, fill="#1A237E", anchor="mm")
    d.text((bx_tcp+w//2, by-30), "TCP/IP Model (4 Layers)", font=FONT_BIG, fill="#F44C00", anchor="mm")
    for i, (name, proto) in enumerate(osi_layers):
        y = by + i*(h+5)
        col = "#1A237E" if i%2==0 else "#3949AB"
        d.rounded_rectangle([bx_osi, y, bx_osi+w, y+h], radius=5, fill=col, outline="white", width=2)
        d.text((bx_osi+15, y+h//2), f"L{7-i}. {name}", font=FONT_SMALL, fill="white", anchor="lm")
        d.text((bx_osi+w-15, y+h//2), proto, font=FONT_SMALL, fill="#FFD180", anchor="rm")
    tcp_layers = [
        ("Application", "HTTP, HTTPS, DNS, SSH"),
        ("Transport", "TCP, UDP"),
        ("Internet", "IP, ICMP, ARP"),
        ("Network Access", "Ethernet, WiFi"),
    ]
    for osi_idx, tcp_idx, span in tcp_map:
        y = by + osi_idx*(h+5)
        height = span*(h+5) - 5
        col = "#F44C00" if tcp_idx%2==0 else "#FF7043"
        d.rounded_rectangle([bx_tcp, y, bx_tcp+w, y+height], radius=5, fill=col, outline="white", width=2)
        d.text((bx_tcp+15, y+height//2), tcp_layers[tcp_idx][0], font=FONT_SMALL, fill="white", anchor="lm")
        d.text((bx_tcp+w-15, y+height//2), tcp_layers[tcp_idx][1], font=FONT_SMALL, fill="#FFD180", anchor="rm")
    for osi_idx, tcp_idx, span in tcp_map:
        y_mid = by + osi_idx*(h+5) + (span*(h+5) - 5)//2
        d.line([(bx_osi+w, y_mid), (bx_tcp, y_mid)], fill="#333", width=2)
    d.text((W//2, 470), "OSI = conceptual teaching model", font=FONT_BODY, fill="#1A237E", anchor="mm")
    d.text((W//2, 505), "TCP/IP = practical implementation on all internet devices", font=FONT_BODY, fill="#F44C00", anchor="mm")
    d.text((W//2, 550), "Cybersecurity professionals must understand both", font=FONT_BIG, fill="#333", anchor="mm")
    save_img(img, "14_osi_tcpip.png")
    
    # 15. TCP Handshake
    img = Image.new("RGB", (W, H), "#F5F7FA")
    d = ImageDraw.Draw(img)
    d.text((W//2, 15), "TCP Three-Way Handshake", font=FONT_TITLE, fill="#1A237E", anchor="mm")
    d.rounded_rectangle([150, 120, 350, 220], radius=10, fill="#1A237E", outline="white", width=3)
    d.text((250, 170), "CLIENT", font=FONT_BIG, fill="white", anchor="mm")
    d.rounded_rectangle([930, 120, 1130, 220], radius=10, fill="#F44C00", outline="white", width=3)
    d.text((1030, 170), "SERVER", font=FONT_BIG, fill="white", anchor="mm")
    d.line([(250, 250), (250, 650)], fill="#CCC", width=2)
    d.line([(1030, 250), (1030, 650)], fill="#CCC", width=2)
    d.line([(250, 300), (1030, 300)], fill="#333", width=3)
    d.polygon([(1030, 300), (1020, 295), (1020, 305)], fill="#333")
    d.text((640, 285), "SYN (Synchronize)", font=FONT_BODY, fill="#1A237E", anchor="mm")
    d.text((250, 310), "Step 1", font=FONT_NOTE, fill="#333", anchor="mm")
    d.line([(1030, 400), (250, 400)], fill="#333", width=3)
    d.polygon([(250, 400), (260, 395), (260, 405)], fill="#333")
    d.text((640, 385), "SYN-ACK (Synchronize-Acknowledge)", font=FONT_BODY, fill="#F44C00", anchor="mm")
    d.text((1030, 410), "Step 2", font=FONT_NOTE, fill="#333", anchor="mm")
    d.line([(250, 500), (1030, 500)], fill="#333", width=3)
    d.polygon([(1030, 500), (1020, 495), (1020, 505)], fill="#333")
    d.text((640, 485), "ACK (Acknowledge)", font=FONT_BODY, fill="#1A237E", anchor="mm")
    d.text((250, 510), "Step 3", font=FONT_NOTE, fill="#333", anchor="mm")
    d.text((W//2, 570), "After the handshake, the connection is established.", font=FONT_BODY, fill="#333", anchor="mm")
    d.text((W//2, 620), "Attack vector: SYN Flood - attackers send SYN but never complete handshake,", font=FONT_NOTE, fill="#F44C00", anchor="mm")
    d.text((W//2, 655), "exhausting server resources and denying service to legitimate users.", font=FONT_NOTE, fill="#F44C00", anchor="mm")
    save_img(img, "15_tcp_handshake.png")
    
    # 16. Network Devices
    img = Image.new("RGB", (W, H), "#F5F7FA")
    d = ImageDraw.Draw(img)
    d.text((W//2, 20), "Common Network Security Devices", font=FONT_TITLE, fill="#1A237E", anchor="mm")
    devices = [
        ("FIREWALL", "Filters traffic based on rules", "Network and Transport layers (L3 and L4)", "#1A237E"),
        ("ROUTER", "Directs traffic between networks", "Network layer (L3). Uses IP addresses", "#F44C00"),
        ("SWITCH", "Connects devices in same network", "Data Link layer (L2). Uses MAC addresses", "#1A237E"),
        ("IDS / IPS", "Monitors for suspicious activity", "Transport to Application layers (L4-L7)", "#F44C00"),
    ]
    for i, (name, desc, layer, col) in enumerate(devices):
        x = 60 + (i % 2) * 620
        y = 70 + (i // 2) * 210
        d.rounded_rectangle([x, y, x+580, y+180], radius=10, fill=col, outline="white", width=3)
        d.text((x+290, y+30), name, font=FONT_BODY, fill="white", anchor="mm")
        d.text((x+290, y+70), desc, font=FONT_NOTE, fill="#FFD180", anchor="mm")
        d.text((x+290, y+120), layer, font=FONT_SMALL, fill="white", anchor="mm")
    save_img(img, "16_network_devices.png")
    
    print("All diagrams generated successfully!")

# ---------------------------------------------------------------------------
# PPTX BUILDER HELPERS
# ---------------------------------------------------------------------------
def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def add_title_slide(prs, title, subtitle, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.333), Inches(1.0))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(26)
    p2.font.color.rgb = RGBColor(0xFF, 0xB7, 0x4D)
    p2.alignment = PP_ALIGN.CENTER
    slide.notes_slide.notes_text_frame.text = notes
    return slide

def add_section_slide(prs, title, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = ACCENT_ORANGE
    bg.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    slide.notes_slide.notes_text_frame.text = notes
    return slide

def add_bullet_slide(prs, title, bullets, notes="", img=None, img_x=Inches(6.9), img_y=Inches(1.35), img_w=Inches(6.0)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    tb_title = slide.shapes.add_textbox(Inches(0.4), Inches(0.22), Inches(12.5), Inches(0.7))
    p = tb_title.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    has_img = img and os.path.exists(img)
    text_width = Inches(6.6) if has_img else Inches(12.5)
    tb_body = slide.shapes.add_textbox(Inches(0.4), Inches(1.12), text_width, Inches(6.38))
    tf = tb_body.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = b
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT
        p.line_spacing = 1.0
        p.space_after = Pt(1 if i < len(bullets) - 1 else 0)
    if has_img:
        try:
            slide.shapes.add_picture(img, img_x, img_y, width=img_w)
        except Exception as e:
            print(f"Image error {img}: {e}")
    slide.notes_slide.notes_text_frame.text = notes
    return slide

def add_two_col_slide(prs, title, left_bullets, right_bullets, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    tb_title = slide.shapes.add_textbox(Inches(0.4), Inches(0.22), Inches(12.5), Inches(0.7))
    p = tb_title.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    for col_idx, bullets in enumerate([left_bullets, right_bullets]):
        x = Inches(0.4 + col_idx*6.4)
        tb = slide.shapes.add_textbox(x, Inches(1.1), Inches(6.2), Inches(6.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        for i, b in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = b
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_TEXT
            p.line_spacing = 1.0
            p.space_after = Pt(1 if i < len(bullets) - 1 else 0)
    slide.notes_slide.notes_text_frame.text = notes
    return slide


# ---------------------------------------------------------------------------
# TOPIC 1: COMPUTER FUNDAMENTALS
# ---------------------------------------------------------------------------
def build_topic_1():
    prs = new_prs()
    
    add_title_slide(prs,
        "Topic 1: Computer Fundamentals",
        "Foundations of Computing and Mathematics | Cybersecurity Certificate Program",
        "Welcome to Topic 1: Computer Fundamentals. This topic establishes the physical and logical building blocks of computing. Every cybersecurity professional must understand how computers process information because every cyberattack ultimately targets hardware, software, or data at some stage. We will examine the Input-Processing-Output-Storage (IPOS) cycle, core hardware components, system architecture, and different types of computing devices. All abbreviations will be defined on first use."
    )
    
    add_bullet_slide(prs,
        "Why Computer Fundamentals Matter in Cybersecurity",
        [
            "- You cannot protect what you do not understand.",
            "- Attackers exploit weaknesses at every level: physical hardware, operating systems, applications, and user behavior.",
            "- Understanding hardware helps incident responders find where malware hides and how data is stolen.",
            "- In Zambia, banks like Zambia National Commercial Bank (Zanaco), mobile network operators like Mobile Telecommunications Network (MTN), and utility companies like Zambia Electricity Supply Corporation (ZESCO) all rely on these same computing principles.",
            "- Job roles that require this knowledge: Security Operations Center (SOC) Analyst, Network Security Technician, Digital Forensics Investigator, and Incident Responder."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "This slide sets the context for the entire topic. Emphasize that cybersecurity is not just about firewalls and antivirus software. It begins with understanding the machines themselves.\n\n"
        "Key point: The 2021 Colonial Pipeline ransomware attack in the United States began with a compromised Virtual Private Network (VPN) account, but the attackers then moved through the organization's computer systems, exploiting basic understanding of how Windows servers, user accounts, and network connections work. If the defenders had stronger foundational knowledge, the attack might have been contained earlier.\n\n"
        "In Zambia, the Information and Communication Technology (ICT) sector is growing rapidly. ZICTA (Zambia Information and Communications Technology Authority) regulates this space, and organizations across the country need people who understand both the technology and how to secure it.\n\n"
        "Abbreviations defined:\n"
        "- IPOS = Input, Processing, Output, Storage\n"
        "- VPN = Virtual Private Network\n"
        "- SOC = Security Operations Center\n"
        "- MTN = Mobile Telecommunications Network\n"
        "- ZESCO = Zambia Electricity Supply Corporation\n"
        "- ZICTA = Zambia Information and Communications Technology Authority"
    )
    
    add_bullet_slide(prs,
        "What is a Computer? The IPOS Model",
        [
            "- A computer is an electronic device that accepts data, processes it according to instructions, stores the results, and produces output.",
            "- The IPOS Model stands for Input, Processing, Output, and Storage.",
            "- Input: Data enters the system through devices like a keyboard, mouse, touchscreen, biometric scanner, or card reader.",
            "- Processing: The Central Processing Unit (CPU) performs calculations and logical operations on the input data.",
            "- Output: The results are displayed to the user through a monitor, printer, or speakers.",
            "- Storage: Data and instructions are saved for future use on devices like Hard Disk Drives (HDD), Solid State Drives (SSD), or Universal Serial Bus (USB) flash drives."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "The IPOS model is fundamental because it breaks down every computing action into four stages. Each stage represents a potential attack surface.\n\n"
        "Input Stage: Attackers can intercept input before it reaches the processor. Examples include keyloggers (software or hardware devices that record every keystroke), fake biometric scanners that steal fingerprint data, and card skimmers placed on Automated Teller Machines (ATMs).\n\n"
        "Processing Stage: The CPU is where all computation happens. Modern CPUs use speculative execution (performing calculations before they are confirmed as needed) to speed up performance. The Spectre and Meltdown vulnerabilities, discovered in 2018, showed that this feature could leak sensitive information between programs. Over 90 percent of processors worldwide were affected.\n\n"
        "Output Stage: Shoulder surfing is when an attacker physically watches a user's screen or keyboard in public places like internet cafes in Lusaka. Screen privacy filters are a physical security control that limits viewing angles.\n\n"
        "Storage Stage: Data at rest (data stored on a device) must be protected. If a laptop is stolen from a government office or non-governmental organization (NGO) in Zambia and the drive is not encrypted, the thief can access all files. Full Disk Encryption (FDE) using BitLocker on Windows or Linux Unified Key Setup (LUKS) on Linux prevents this.\n\n"
        "Abbreviations defined:\n"
        "- CPU = Central Processing Unit\n"
        "- HDD = Hard Disk Drive\n"
        "- SSD = Solid State Drive\n"
        "- USB = Universal Serial Bus\n"
        "- ATM = Automated Teller Machine\n"
        "- NGO = Non-Governmental Organization\n"
        "- FDE = Full Disk Encryption\n"
        "- LUKS = Linux Unified Key Setup",
        img=os.path.join(ASSETS_DIR, "ipos_cycle.png")
    )
    
    add_bullet_slide(prs,
        "Hardware Components and Their Security Relevance",
        [
            "- Central Processing Unit (CPU): Often called the brain of the computer. It executes program instructions and performs arithmetic and logical operations. Speed is measured in gigahertz (GHz). Modern CPUs have multiple cores, allowing them to perform multiple tasks simultaneously.",
            "- Security relevance: CPUs handle cryptographic operations like encrypting and decrypting data. Vulnerabilities such as Spectre and Meltdown are side-channel attacks that exploit how the CPU processes data to leak secrets across programs.",
            "- Random Access Memory (RAM): Temporary, high-speed memory that the CPU uses while the computer is running. It is volatile, meaning all data is lost when power is turned off.",
            "- Security relevance: Fileless malware operates entirely in RAM to avoid being detected by traditional antivirus software that scans files on the hard drive. Incident responders use memory forensics tools like Volatility to detect such threats."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Central Processing Unit (CPU): The CPU contains two main parts: the Arithmetic Logic Unit (ALU), which performs mathematical calculations and logical comparisons, and the Control Unit (CU), which fetches instructions from memory and directs operations. Multi-core processors (dual-core, quad-core, etc.) allow parallel processing. In cybersecurity, the CPU is responsible for encrypting data using algorithms like the Advanced Encryption Standard (AES). However, researchers discovered that speculative execution - a performance optimization where the CPU guesses which instructions it will need next - could be exploited. The Spectre and Meltdown attacks allowed malicious programs to read memory that should have been protected, including passwords and encryption keys. Over 90 percent of processors worldwide were affected.\n\n"
        "Random Access Memory (RAM): RAM is much faster than storage but temporary. When you open a document, it is loaded from storage into RAM so the CPU can work with it quickly. Because RAM does not retain data after power loss, some attackers perform cold boot attacks: they quickly restart a computer from a special device and dump the RAM contents before the data fades. Fileless malware (also called living-off-the-land attacks) resides only in RAM, using legitimate system tools like PowerShell to avoid writing any files to disk. This makes it invisible to signature-based antivirus. Memory forensics is the practice of analyzing RAM dumps to find such threats. Tools like Volatility and Rekall are industry standards for this work.\n\n"
        "Abbreviations defined:\n"
        "- ALU = Arithmetic Logic Unit\n"
        "- CU = Control Unit\n"
        "- AES = Advanced Encryption Standard\n"
        "- RAM = Random Access Memory",
        img=os.path.join(ASSETS_DIR, "hardware_components.png")
    )
    
    add_bullet_slide(prs,
        "More Hardware Components: Storage and Network Interface",
        [
            "- Storage Devices: Hard Disk Drives (HDD) use spinning magnetic platters. Solid State Drives (SSD) use flash memory with no moving parts, making them faster and more durable. USB flash drives provide portable storage.",
            "- Security relevance: Stolen storage devices are a major data breach vector. Encryption at rest ensures that even if a device is physically stolen, the data remains unreadable without the decryption key. Forensic investigators also examine storage for deleted files and hidden data.",
            "- Network Interface Card (NIC): This hardware component connects the computer to a network, either through an Ethernet cable or wirelessly via WiFi.",
            "- Every NIC has a unique Media Access Control (MAC) address burned into it by the manufacturer, such as 00:1A:2B:3C:4D:5E.",
            "- Security relevance: MAC addresses identify devices on a local network. Attackers can use MAC spoofing (changing their MAC address to impersonate another device) to bypass Network Access Control (NAC) systems or frame innocent users."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Storage Devices: Hard Disk Drives (HDD) store data magnetically on spinning platters. They are slower but offer large capacities at low cost. Solid State Drives (SSD) store data in flash memory chips. They are faster, quieter, and more shock-resistant but typically more expensive per gigabyte. USB (Universal Serial Bus) flash drives are convenient for transferring files but are easily lost or stolen. In 2022, the United Kingdom's Information Commissioner's Office (ICO) fined multiple organizations millions of pounds after unencrypted USB drives containing personal data were lost.\n\n"
        "Encryption at rest means that data is encrypted while stored on the device. Even if the physical device is stolen, the thief cannot read the data without the encryption key. Windows includes BitLocker for this purpose. Linux includes LUKS (Linux Unified Key Setup). Both require a password or recovery key to decrypt the drive.\n\n"
        "Network Interface Card (NIC): Every device that connects to a network has a NIC. Wired NICs use Ethernet cables. Wireless NICs use WiFi. Each NIC has a unique 48-bit identifier called the Media Access Control (MAC) address. The first 24 bits (three pairs of hex digits) identify the manufacturer. The last 24 bits identify the specific device.\n\n"
        "MAC spoofing is when an attacker changes their device's MAC address to match another device's address. This can be used to bypass Network Access Control (NAC), which restricts network access to known devices. It can also be used to impersonate another user on a network. Tools like macchanger on Linux make MAC spoofing easy. However, MAC addresses only work on the local network; they are not sent over the internet. Internet routing uses Internet Protocol (IP) addresses instead.\n\n"
        "Abbreviations defined:\n"
        "- HDD = Hard Disk Drive\n"
        "- SSD = Solid State Drive\n"
        "- USB = Universal Serial Bus\n"
        "- ICO = Information Commissioner's Office\n"
        "- NIC = Network Interface Card\n"
        "- MAC = Media Access Control\n"
        "- NAC = Network Access Control\n"
        "- IP = Internet Protocol"
    )
    
    add_bullet_slide(prs,
        "System Architecture: 32-bit versus 64-bit",
        [
            "- A 32-bit system can address a maximum of 4 gigabytes (GB) of Random Access Memory (RAM).",
            "- A 64-bit system can theoretically address 16 exabytes of RAM - far more than any current computer needs.",
            "- Most modern computers, smartphones, and servers use 64-bit architecture.",
            "- Security relevance: 64-bit operating systems support stronger security features including Address Space Layout Randomization (ASLR) and Data Execution Prevention (DEP), also known as the No-eXecute (NX) bit.",
            "- ASLR randomizes where programs load in memory, making it harder for attackers to predict memory addresses during exploits.",
            "- DEP marks certain memory regions as non-executable, preventing buffer overflow attacks from running injected malicious code."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "The terms 32-bit and 64-bit refer to how much data the processor can handle in one operation and how much memory it can access. A 32-bit processor uses 32 binary digits (bits) to represent memory addresses. Because 2 raised to the power of 32 equals approximately 4.29 billion, a 32-bit system can only use about 4 gigabytes (GB) of RAM. This was sufficient in the 1990s but is a serious limitation today.\n\n"
        "A 64-bit processor uses 64 bits for memory addressing. While no computer today uses anywhere near the theoretical maximum of 16 exabytes (1 exabyte = 1 billion gigabytes), the larger address space allows systems to use hundreds of gigabytes of RAM, which is essential for running modern security tools like Security Information and Event Management (SIEM) systems.\n\n"
        "From a cybersecurity perspective, the move to 64-bit is about more than just memory capacity. 64-bit operating systems can implement stronger versions of Address Space Layout Randomization (ASLR). ASLR is a security technique that randomly arranges the memory positions of key data areas, including the base of the executable and the positions of the stack, heap, and libraries. This makes it much harder for an attacker to predict where to redirect execution during a buffer overflow or other memory corruption attack.\n\n"
        "Data Execution Prevention (DEP) is another critical defense. Also called the NX (No-eXecute) bit on AMD processors or XD (eXecute Disable) bit on Intel processors, DEP marks certain areas of memory as non-executable. In a buffer overflow attack, the attacker typically injects malicious code into a memory area that should only contain data (like the stack). With DEP enabled, the processor will refuse to execute code from that region, causing the attack to fail.\n\n"
        "Abbreviations defined:\n"
        "- GB = Gigabyte\n"
        "- RAM = Random Access Memory\n"
        "- ASLR = Address Space Layout Randomization\n"
        "- DEP = Data Execution Prevention\n"
        "- NX = No-eXecute\n"
        "- SIEM = Security Information and Event Management"
    )
    
    add_bullet_slide(prs,
        "Von Neumann Architecture and Security Implications",
        [
            "- The Von Neumann architecture, developed by mathematician John von Neumann in 1945, is the design model used by almost all modern computers.",
            "- It consists of five main parts: Input devices, the Central Processing Unit (CPU) containing the Arithmetic Logic Unit (ALU) and Control Unit (CU), Memory (RAM), Storage, and Output devices.",
            "- The defining characteristic is that both data and program instructions are stored in the same memory space.",
            "- Security implication: Because data and instructions share memory, a buffer overflow attack can overwrite program instructions with malicious code, causing the computer to execute attacker-controlled commands.",
            "- Modern defenses include DEP/NX bit, ASLR, and stack canaries (special values placed before return addresses to detect overflows)."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "John von Neumann was a Hungarian-American mathematician who, in 1945, published a paper describing a computer design where the program (instructions) and data are both stored in the same memory. Before this, computers were programmed by physically rewiring circuits or using punch cards that were separate from the data. The stored-program concept revolutionized computing and is still used today.\n\n"
        "The five components are:\n"
        "1. Input: Devices that bring data into the computer (keyboard, mouse, scanner, microphone).\n"
        "2. CPU (Central Processing Unit): The brain containing the ALU (Arithmetic Logic Unit) for math and the Control Unit (CU) for directing operations.\n"
        "3. Memory (RAM): Fast temporary storage for both data and instructions currently in use.\n"
        "4. Storage: Long-term retention of programs and files when the computer is off.\n"
        "5. Output: Devices that present results to the user (monitor, printer, speakers).\n\n"
        "The security problem: Because instructions and data live in the same memory, there is no hard boundary preventing a program from accidentally or maliciously treating data as instructions. This is the root cause of buffer overflow attacks. When a program receives more input than expected and fails to check the size, the excess input can spill over into adjacent memory that contains return addresses or executable instructions. An attacker can craft input that includes malicious machine code and overwrite a return address to point to that code. When the function returns, the CPU executes the attacker's code.\n\n"
        "Defenses evolved to mitigate this architectural weakness:\n"
        "- Stack canaries: A secret value placed on the stack before the return address. If a buffer overflow occurs, the canary is overwritten first. Before returning, the program checks if the canary has changed. If so, it terminates immediately.\n"
        "- DEP/NX bit: Marks data regions (stack, heap) as non-executable.\n"
        "- ASLR: Randomizes memory locations so attackers cannot predict where to jump.\n\n"
        "The 2017 WannaCry ransomware exploited EternalBlue, a buffer overflow in the Windows Server Message Block (SMB) protocol, to spread automatically across networks. Understanding Von Neumann architecture helps students understand WHY such vulnerabilities are possible.",
        img=os.path.join(ASSETS_DIR, "von_neumann.png")
    )
    
    add_bullet_slide(prs,
        "Types of Computers and Their Security Posture",
        [
            "- Personal Computer (PC): Includes desktop and laptop computers used by individuals. Security focus is on endpoint protection, including antivirus software, Endpoint Detection and Response (EDR), and user awareness training.",
            "- Server: A powerful computer that provides services to many users simultaneously, such as hosting websites, managing databases, or handling email. Servers are high-value targets and require hardening (securing the configuration), continuous monitoring, and regular patching.",
            "- Embedded System: A computer built into another device, such as an ATM, traffic light, or medical equipment. These often run outdated software that cannot be easily patched, making them vulnerable.",
            "- Internet of Things (IoT) Device: Smart devices connected to the internet, such as security cameras, smart meters, and home assistants. These devices often ship with default passwords and lack security updates, creating a massive attack surface."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Different categories of computers face different threats and require different security strategies.\n\n"
        "Personal Computers (PCs): In Zambia, most office workers at banks, schools, and government offices use Windows PCs. The primary security controls for PCs are endpoint security solutions. Antivirus software detects known malware signatures. Endpoint Detection and Response (EDR) goes further by monitoring behavior in real time, detecting fileless malware, and allowing remote investigation. User awareness training teaches employees to recognize phishing emails, which are the most common attack vector against endpoints.\n\n"
        "Servers: A single compromised web server can expose thousands or millions of customer records. In 2017, Equifax - one of the largest credit reporting agencies in the world - suffered a breach exposing 147 million records. The cause was an unpatched Apache Struts web server. In Zambia, servers at MTN, Airtel, Zanaco, and government data centers predominantly run Linux. This makes Linux security skills highly employable. Server hardening includes disabling unnecessary services, applying security patches promptly, using strong authentication, and maintaining detailed logs.\n\n"
        "Embedded Systems: Unlike PCs and servers, embedded systems are often physically difficult to access and are designed to run for years without intervention. This means they frequently run old operating systems with known vulnerabilities that were never patched. An ATM running Windows XP (which Microsoft stopped supporting in 2014) is vulnerable to attacks that have been public knowledge for years.\n\n"
        "Internet of Things (IoT): The number of IoT devices worldwide is in the tens of billions. In Zambia, smart electricity meters, cheap Closed-Circuit Television (CCTV) cameras, and smart home devices are increasingly common. Many of these devices ship with default usernames and passwords like admin/admin or root/root. The 2016 Mirai botnet infected over 600,000 IoT devices with default credentials and used them to launch record-breaking Distributed Denial of Service (DDoS) attacks. The first security step for any IoT device is changing the default password and isolating it on a separate network segment.\n\n"
        "Abbreviations defined:\n"
        "- PC = Personal Computer\n"
        "- EDR = Endpoint Detection and Response\n"
        "- IoT = Internet of Things\n"
        "- CCTV = Closed-Circuit Television\n"
        "- DDoS = Distributed Denial of Service"
    )
    
    add_bullet_slide(prs,
        "Topic 1 Summary: Key Takeaways",
        [
            "- The IPOS Model (Input, Processing, Output, Storage) describes how computers handle data. Every stage represents an attack surface.",
            "- The Central Processing Unit (CPU) executes instructions and handles encryption. Side-channel attacks like Spectre and Meltdown target the CPU.",
            "- Random Access Memory (RAM) is temporary and volatile. Fileless malware hides in RAM to avoid disk-based antivirus detection.",
            "- Storage devices (HDD, SSD, USB) hold data permanently. Encryption at rest protects against theft and unauthorized access.",
            "- The Network Interface Card (NIC) connects devices to networks using a unique Media Access Control (MAC) address. MAC spoofing is a common attack technique.",
            "- The Von Neumann architecture stores data and instructions in the same memory, which enables buffer overflow attacks. Modern defenses include ASLR, DEP/NX bit, and stack canaries.",
            "- Different computer types (PC, server, embedded, IoT) require different security strategies."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "This summary slide recaps the foundational concepts of Topic 1. Emphasize the connections between hardware and security.\n\n"
        "The IPOS model is not just theory - it maps to real attack vectors. Input = keyloggers and skimmers. Processing = CPU side-channel attacks. Output = shoulder surfing and screen capture malware. Storage = data theft from unencrypted devices.\n\n"
        "The CPU and RAM are primary targets for advanced attackers. While traditional antivirus scans files on the hard drive, sophisticated threats live in memory. This is why memory forensics and behavior-based detection (EDR) are essential skills for modern cybersecurity professionals.\n\n"
        "The Von Neumann architecture is a historical design decision with lasting security consequences. Buffer overflows have existed for decades and remain one of the most dangerous vulnerability classes. Every major operating system invests heavily in mitigations like ASLR and DEP.\n\n"
        "Finally, the growth of IoT devices has massively expanded the attack surface. In Zambia and globally, organizations struggle to inventory and secure these devices. Students who understand how to segment networks and harden embedded systems will be valuable assets."
    )
    
    out = os.path.join(BASE_DIR, "Topic_1_Computer_Fundamentals.pptx")
    prs.save(out)
    print(f"Saved: {out}")


# ---------------------------------------------------------------------------
# TOPIC 2: OPERATING SYSTEMS
# ---------------------------------------------------------------------------
def build_topic_2():
    prs = new_prs()
    
    add_title_slide(prs,
        "Topic 2: Operating Systems",
        "Windows and Linux Fundamentals for Cybersecurity",
        "Welcome to Topic 2: Operating Systems. An Operating System (OS) is the software that manages computer hardware and provides services for applications. It sits between the user and the hardware, making it the ultimate security gatekeeper. Every file you open, every network connection you make, and every password you type passes through the OS. If the OS is compromised, nothing running on it can be trusted. This topic covers Windows and Linux, the two most important operating systems for cybersecurity professionals."
    )
    
    add_bullet_slide(prs,
        "What is an Operating System?",
        [
            "- An Operating System (OS) is system software that manages hardware resources and provides common services for computer programs.",
            "- It acts as an intermediary between the user, applications, and the physical hardware.",
            "- Key functions of an OS include:",
            "   + Process management: Creating, scheduling, and terminating running programs.",
            "   + Memory management: Allocating and protecting Random Access Memory (RAM).",
            "   + File system management: Organizing, reading, and writing files on storage devices.",
            "   + Security and access control: Managing user accounts, passwords, and permissions.",
            "   + Device management: Communicating with hardware through drivers."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "The Operating System (OS) is the most critical piece of software on any computer. Without an OS, the hardware cannot be used. The OS is responsible for making sure that multiple programs can run simultaneously without interfering with each other.\n\n"
        "Process Management: Every running program is called a process. The OS decides which process gets to use the CPU and for how long. Attackers often try to hide malicious processes by giving them names that look like legitimate system processes. For example, malware might name itself svch0st.exe (with a zero) instead of svchost.exe (with the letter o), hoping that a busy administrator will not notice the difference.\n\n"
        "Memory Management: The OS ensures that each process only accesses its own allocated memory. If a program tries to read or write memory belonging to another program, the OS should stop it. However, vulnerabilities in the OS or applications can bypass these protections, leading to attacks like buffer overflows and heap spraying.\n\n"
        "File System Management: The OS controls how files are created, named, stored, and retrieved. It also enforces permissions - who can read, write, or delete files. Poorly configured file permissions are a common cause of data breaches.\n\n"
        "Security and Access Control: The OS verifies who you are (authentication) and what you are allowed to do (authorization). Windows uses Active Directory in corporate environments. Linux uses the Pluggable Authentication Modules (PAM) system.\n\n"
        "Device Management: The OS communicates with hardware devices through drivers. Because drivers run with kernel-level privileges (the highest level of system access), a vulnerability in a printer driver, graphics card driver, or network driver can give an attacker complete control of the computer. The 2021 PrintNightmare vulnerability in Windows was a driver flaw that allowed remote code execution.\n\n"
        "Abbreviations defined:\n"
        "- OS = Operating System\n"
        "- RAM = Random Access Memory\n"
        "- CPU = Central Processing Unit\n"
        "- PAM = Pluggable Authentication Modules"
    )
    
    add_bullet_slide(prs,
        "Windows Operating System: File System Basics",
        [
            "- Windows is the most widely used desktop operating system globally and in Zambia.",
            "- The file system is organized in a tree structure starting from drive letters like C:\\.",
            "- C:\\Windows\\System32 contains core operating system files, Dynamic Link Libraries (DLLs), and executables. This is a frequent target for malware.",
            "- C:\\Users\\[Username]\\AppData is a hidden directory that stores application data. Malware frequently uses this location for persistence because it is writable by standard users.",
            "- C:\\Temp and C:\\Windows\\Temp store temporary files. Incident responders check these directories for recently dropped malware.",
            "- C:\\Program Files and C:\\Program Files (x86) contain installed applications. Unauthorized software here may indicate a security breach."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Microsoft Windows dominates the personal computer market, with approximately 70 percent market share worldwide. In Zambian businesses, schools, and government offices, Windows is the standard desktop operating system. Because of this dominance, cybersecurity professionals spend a large portion of their time investigating, securing, and managing Windows systems.\n\n"
        "C:\\Windows\\System32: This directory is the heart of the Windows operating system. It contains hundreds of Dynamic Link Libraries (DLLs) - shared code libraries used by multiple programs - and essential executables like cmd.exe (the Command Prompt) and svchost.exe (which hosts Windows services). Malware frequently targets System32 through a technique called DLL hijacking. In this attack, a malicious DLL with the same name as a legitimate one is placed in a directory where Windows will load it before the real version. When a program runs, it unknowingly loads the malicious DLL and executes attacker code.\n\n"
        "AppData: This directory is hidden by default, meaning users do not see it unless they enable 'Show hidden files' in File Explorer. It has three subfolders: Local, LocalLow, and Roaming. Malware loves AppData because standard users can write to it without administrator privileges, and some antivirus programs exclude it from scanning to avoid breaking legitimate applications. Ransomware families like Ryuk and TrickBot commonly store configuration files and secondary payloads in AppData. To view it, open File Explorer, click the View tab, and check 'Hidden items.' Alternatively, in Command Prompt, use: dir /a C:\\Users\\[username]\\AppData\n\n"
        "Temp Directories: The %TEMP% environment variable points to a temporary folder, usually C:\\Users\\[username]\\AppData\\Local\\Temp. Any user can write here, making it an ideal place for malware to extract and execute files. During an incident response investigation, one of the first commands an analyst runs is dir /o-d %TEMP% (the /o-d flag sorts by date, newest first). Recently created executable files in Temp are highly suspicious.\n\n"
        "Program Files: This is where most installed applications live. During security audits, analysts check for unauthorized or outdated software. Old versions of web browsers, PDF readers, or media players often contain known vulnerabilities that attackers can exploit.\n\n"
        "Abbreviations defined:\n"
        "- DLL = Dynamic Link Library\n"
        "- cmd.exe = Command Prompt executable",
        img=os.path.join(ASSETS_DIR, "windows_fs.png")
    )
    
    add_bullet_slide(prs,
        "Windows Command Prompt: Essential Security Commands",
        [
            "- ipconfig: Displays Internet Protocol (IP) configuration, including the computer's IP address, subnet mask, and default gateway. Using ipconfig /all shows additional details including the Media Access Control (MAC) address and Domain Name System (DNS) servers.",
            "- netstat -an: Displays all active network connections and listening ports. Analysts use this to spot suspicious connections to unknown remote addresses.",
            "- tasklist: Lists all running processes. Using tasklist /svc shows which Windows services each process is hosting, helping identify injected or rogue processes.",
            "- systeminfo: Displays detailed system information including the operating system version, installed hotfixes (security patches), and hardware details.",
            "- net user: Lists local user accounts. Attackers often create hidden backdoor accounts with innocent-sounding names.",
            "- dir /a: Lists all files in a directory, including hidden and system files."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "The Windows Command Prompt (cmd.exe) and its more powerful successor, PowerShell, are essential tools for system administration and incident response. While PowerShell has largely replaced Command Prompt for advanced tasks, these basic commands are guaranteed to work on every Windows system and are excellent starting points for beginners.\n\n"
        "ipconfig: This command reveals the computer's network identity. The output includes the IPv4 address (the local network address), subnet mask (which defines the network boundary), default gateway (the router that connects to other networks), and MAC address. The command ipconfig /displaydns shows the DNS resolver cache - a record of recently resolved domain names. If you see entries for suspicious domains, it may indicate malware communication. Unexpected DNS server addresses can also be a sign of DNS hijacking.\n\n"
        "netstat -an: This is one of the most important commands for network investigation. The -a flag shows all connections and listening ports. The -n flag displays addresses and port numbers in numerical form, which is faster than resolving names. If you see a connection to an unusual foreign IP address on port 4444 (the default port for the Metasploit penetration testing framework) or port 3389 (Remote Desktop Protocol), it may indicate that the computer has been compromised. For more detail, the command netstat -anb (requires administrator privileges) shows which executable file owns each connection.\n\n"
        "tasklist: This lists all running processes. A process is a running instance of a program. The basic tasklist command shows the process name and Process Identifier (PID). Adding the /svc flag shows which Windows services are running inside each svchost.exe process. Malware frequently injects itself into legitimate processes or creates processes with names that mimic legitimate ones. Comparing tasklist output against a known-good baseline helps spot anomalies.\n\n"
        "systeminfo: This command outputs a comprehensive system profile, including the operating system name and version, installation date, registered owner, processor type, and total physical memory. Most importantly for security, it lists installed hotfixes (security patches). An unpatched system is a vulnerable system. For example, if the hotfix for EternalBlue (MS17-010) is missing, the computer is vulnerable to WannaCry and similar ransomware.\n\n"
        "net user: This command lists all local user accounts on the computer. Attackers who gain access often create backdoor accounts with names like 'helpdesk,' 'backup,' or 'support' to maintain access. The command net localgroup administrators shows which accounts have administrator privileges. Any unexpected account in this list is a critical red flag.\n\n"
        "dir /a: The /a flag shows all files, including hidden files (which have the Hidden attribute) and system files. Malware often marks itself as hidden to avoid detection.\n\n"
        "Abbreviations defined:\n"
        "- IP = Internet Protocol\n"
        "- MAC = Media Access Control\n"
        "- DNS = Domain Name System\n"
        "- PID = Process Identifier\n"
        "- MS17-010 = Microsoft Security Bulletin 2017-010 (EternalBlue patch)"
    )
    
    add_bullet_slide(prs,
        "Windows User Accounts and the Principle of Least Privilege",
        [
            "- Administrator Account: Has full control over the system. Can install software, modify all files, and change security settings. Should be used only when necessary.",
            "- Standard User Account: Has limited privileges. Cannot install software system-wide or modify protected operating system files. This is the recommended account type for daily use.",
            "- Guest Account: Has very limited access and cannot make persistent changes to the system.",
            "- The Principle of Least Privilege means giving users only the minimum level of access they need to perform their job.",
            "- Security impact: If a standard user account is compromised by malware, the malware's ability to damage the system is limited. If an administrator account is compromised, the attacker gains full control."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "User account management is a fundamental security control in Windows. Microsoft designed the User Account Control (UAC) system to reduce the damage caused by malware running under administrator accounts. When UAC is enabled, even administrators run with standard privileges by default, and Windows prompts for confirmation before performing administrative tasks.\n\n"
        "However, many users and even some organizations disable UAC or give everyone administrator accounts for convenience. This is a dangerous practice. If an attacker tricks an administrator into running malware - for example, through a phishing email with a malicious attachment - the malware immediately has full system access. It can disable antivirus, create hidden accounts, install persistent backdoors, and encrypt files (ransomware).\n\n"
        "In contrast, if the same malware runs on a standard user account, its capabilities are restricted. It cannot modify system directories, install drivers, or change security settings. While it can still steal the user's personal files, it cannot easily spread across the network or deeply compromise the system. This is why the Principle of Least Privilege is one of the most effective and least expensive security controls.\n\n"
        "In enterprise environments, Windows uses Active Directory (AD) to centrally manage thousands of user accounts. Group Policy allows administrators to enforce least privilege across the entire organization. For example, Group Policy can prevent standard users from running executable files from USB drives, disable guest accounts, or enforce password complexity requirements.\n\n"
        "Abbreviations defined:\n"
        "- UAC = User Account Control\n"
        "- AD = Active Directory"
    )
    
    add_bullet_slide(prs,
        "Linux Operating System: Why It Matters in Cybersecurity",
        [
            "- Linux is an open-source operating system based on Unix. The source code is freely available, and anyone can inspect, modify, and distribute it.",
            "- Linux powers the majority of web servers, cloud infrastructure, and supercomputers worldwide. Over 80 percent of web servers run Linux.",
            "- In Zambia, core infrastructure at banks, mobile network operators, and government institutions runs on Linux servers.",
            "- Kali Linux is a specialized Linux distribution designed for cybersecurity professionals. It comes pre-installed with hundreds of penetration testing and forensics tools.",
            "- Because Linux is free, highly configurable, and stable, it is the platform of choice for most security tools and research."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Linux was created in 1991 by Linus Torvalds, a Finnish computer science student. Unlike Windows, which is developed by Microsoft as proprietary software, Linux is open-source. This means its source code is publicly available, and a global community of developers contributes to it. There are many different versions of Linux called distributions (or distros), including Ubuntu, Debian, Fedora, CentOS, and Red Hat Enterprise Linux.\n\n"
        "The dominance of Linux in the server market cannot be overstated. According to W3Techs, over 80 percent of all web servers run Linux. Major websites like Google, Facebook, Amazon, and Wikipedia all rely heavily on Linux infrastructure. In the cloud computing market, Amazon Web Services (AWS), Microsoft Azure, and Google Cloud Platform all run Linux virtual machines as their primary compute offering.\n\n"
        "For cybersecurity professionals, Linux proficiency is non-negotiable. The vast majority of security tools - including Nmap (network scanner), Wireshark (packet analyzer), Metasploit (penetration testing framework), and Volatility (memory forensics) - were originally developed for Linux. Kali Linux, maintained by Offensive Security, is a Debian-based distribution that includes over 600 pre-installed security tools. It is the standard operating system used in ethical hacking courses and professional penetration testing.\n\n"
        "In Zambia, Linux skills are directly employable. MTN Zambia, Airtel Zambia, Zanaco, ZESCO, and ZICTA all maintain Linux-based servers for web hosting, databases, and network infrastructure. Students who learn Linux command-line navigation, file permissions, and basic scripting will have a significant advantage when applying for technical roles.\n\n"
        "Abbreviations defined:\n"
        "- AWS = Amazon Web Services\n"
        "- ZICTA = Zambia Information and Communications Technology Authority"
    )
    
    add_bullet_slide(prs,
        "Linux File System Hierarchy",
        [
            "- The Linux file system is organized as a single tree starting from the root directory, represented by a forward slash (/).",
            "- /bin/ contains essential system commands available to all users, such as ls (list files), cat (display file contents), and chmod (change file permissions).",
            "- /etc/ contains system-wide configuration files. Critical security files here include /etc/passwd (user account information), /etc/shadow (encrypted password hashes), and /etc/hosts (local hostname mappings).",
            "- /home/ contains personal directories for each user. For example, /home/mulenga/ belongs to the user named mulenga.",
            "- /var/log/ contains system log files. This is where authentication attempts, system errors, and application activity are recorded.",
            "- /tmp/ is a temporary directory that any user can write to. Malware frequently drops files here.",
            "- /root/ is the home directory of the root user - the Linux equivalent of the Windows Administrator account."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Unlike Windows, which uses drive letters like C:\\ and D:\\, Linux treats all storage as part of a single directory tree starting at the root (/). This includes hard drives, USB drives, network shares, and even special system information.\n\n"
        "/bin/ (binaries): This directory contains essential command-line programs that all users need. The name 'bin' comes from 'binary,' which is the machine-readable format of compiled programs. Important commands here include ls (list directory contents), cp (copy files), mv (move files), rm (remove files), and chmod (change permissions). Because these programs are essential for system recovery, /bin/ is kept separate from other directories and is available even when the system is in rescue mode.\n\n"
        "/etc/ (et cetera): This directory contains system-wide configuration files. It is one of the most important directories for security professionals. Key files include:\n"
        "- /etc/passwd: Contains a list of all user accounts on the system, along with their user ID numbers, home directories, and default shells. By design, this file is readable by all users, but it does NOT contain the actual passwords.\n"
        "- /etc/shadow: Contains the encrypted (hashed) passwords for all users. This file should be readable ONLY by the root user. If an attacker gains read access to /etc/shadow, they can attempt to crack the passwords offline using tools like John the Ripper or Hashcat.\n"
        "- /etc/hosts: Maps hostnames to IP addresses locally. Attackers sometimes modify this file to redirect traffic from legitimate websites (like banking sites) to malicious servers controlled by the attacker. This is called hosts file poisoning.\n"
        "- /etc/crontab and /etc/cron.d/: Define scheduled tasks. Attackers often add malicious cron jobs to maintain persistence - ensuring their malware runs again after a reboot.\n\n"
        "/var/log/ (variable logs): This is the security goldmine. Linux logs almost everything that happens on the system. The most important log files for security include:\n"
        "- /var/log/auth.log (on Debian/Ubuntu systems) or /var/log/secure (on Red Hat/CentOS systems): Records every authentication attempt, including successful logins, failed logins, and sudo usage. Multiple failed login attempts from the same IP address indicate a brute-force attack.\n"
        "- /var/log/syslog: General system messages.\n"
        "- /var/log/apache2/access.log or /var/log/httpd/access_log: Web server access logs. These show every request made to a website, including the source IP address, requested file, and response code.\n"
        "- /var/log/audit/audit.log: Detailed audit logs, especially on systems running Security-Enhanced Linux (SELinux).\n\n"
        "/tmp/ (temporary): This directory is world-writable, meaning any user can create, modify, or delete files here. Because of this, it is a favorite drop zone for malware. Many modern Linux distributions mitigate this risk by mounting /tmp with special flags like noexec (which prevents execution of files in /tmp) and nodev (which prevents device files).\n\n"
        "/root/: This is the home directory of the root user, who has unlimited privileges on the system. Unlike regular users whose home directories are under /home/, root's home is separate for security reasons. Regular users should never be using the root account for daily tasks. Instead, they should use the sudo (superuser do) command to temporarily elevate privileges when needed.\n\n"
        "Abbreviations defined:\n"
        "- IP = Internet Protocol\n"
        "- SELinux = Security-Enhanced Linux\n"
        "- sudo = superuser do",
        img=os.path.join(ASSETS_DIR, "linux_fs.png")
    )
    
    add_bullet_slide(prs,
        "Essential Linux Commands for Security Professionals",
        [
            "- ls -la: Lists files in long format, showing permissions, owner, size, and modification date. The -a flag includes hidden files (files starting with a dot).",
            "- cat /etc/passwd: Displays the contents of the /etc/passwd file, which lists all user accounts on the system.",
            "- grep 'Failed' /var/log/auth.log: Searches the authentication log for lines containing the word 'Failed.' This helps identify brute-force login attempts.",
            "- ps aux: Displays all running processes on the system, including those started by other users and system services.",
            "- sudo chmod 600 file.txt: Changes file permissions so that only the owner can read and write the file. The group and others have no permissions.",
            "- whoami: Displays the username of the currently logged-in user.",
            "- history: Shows the command history for the current session. Attackers may delete or modify this to hide their tracks."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Linux command-line proficiency is essential for anyone pursuing a career in cybersecurity. The command line provides direct, low-level access to the system that is often faster and more powerful than graphical tools.\n\n"
        "ls -la: The ls command lists directory contents. The -l flag uses the 'long' format, showing permissions, number of links, owner name, group name, file size, modification date, and filename. The -a flag shows all files, including hidden files that start with a dot (.). Hidden files are not displayed by default. Attackers frequently hide configuration files, SSH keys, and backdoors as hidden files. For example, .bashrc is a hidden file in the user's home directory that executes every time a terminal is opened. An attacker who modifies .bashrc can achieve persistence - their malicious code will run automatically whenever the user opens a terminal.\n\n"
        "cat: Short for 'concatenate,' the cat command displays the contents of a file. It is commonly used to quickly view configuration files, log files, or scripts. For large files, commands like less or more are better because they allow scrolling.\n\n"
        "grep: This is one of the most powerful text-searching tools in Linux. It searches files for lines matching a given pattern. Examples for security analysis:\n"
        "- grep 'Failed password' /var/log/auth.log | wc -l  (counts failed login attempts)\n"
        "- grep -i 'error|warning' /var/log/syslog  (finds lines containing either 'error' or 'warning,' case-insensitive)\n"
        "- grep '192.168.1.50' /var/log/apache2/access.log  (finds all web requests from a specific IP address)\n\n"
        "ps aux: The ps command shows running processes. The a flag shows processes for all users. The u flag displays the user who owns each process. The x flag includes processes not attached to a terminal (background services). Security analysts use ps aux to identify suspicious processes, unusually high CPU usage, or processes running from unexpected locations like /tmp/.\n\n"
        "sudo chmod 600 file.txt: The chmod command changes file permissions. The numeric mode 600 means: owner has read and write permissions (4+2=6), group has no permissions (0), and others have no permissions (0). This is the recommended permission for sensitive files like SSH private keys (~/.ssh/id_rsa). If an SSH private key has permissions like 644 (readable by everyone), anyone with access to the system can copy it and use it to impersonate the owner on remote servers.\n\n"
        "whoami: A simple but useful command that prints the current username. When performing security tasks, it is important to verify whether you are running as a regular user or as root.\n\n"
        "history: Bash (the default Linux shell) keeps a history of commands entered by the user. The history command displays this list. Attackers often try to clear the history using commands like history -c or rm ~/.bash_history to cover their tracks. A missing or unusually short history file can itself be a sign of compromise.\n\n"
        "Abbreviations defined:\n"
        "- SSH = Secure Shell\n"
        "- Bash = Bourne Again Shell\n"
        "- CPU = Central Processing Unit"
    )
    
    add_bullet_slide(prs,
        "Understanding Linux File Permissions",
        [
            "- Every file and directory in Linux has three sets of permissions: one for the owner, one for the group, and one for all other users.",
            "- The three permission types are: Read (r), Write (w), and Execute (x).",
            "- Read (r) allows viewing the contents of a file or listing the contents of a directory.",
            "- Write (w) allows modifying a file or adding and removing files from a directory.",
            "- Execute (x) allows running a file as a program or script. For directories, it allows entering the directory.",
            "- Permissions can be represented numerically: Read = 4, Write = 2, Execute = 1. The values are added together for each category.",
            "- Example: chmod 755 script.sh means the owner has read, write, and execute (7 = 4+2+1), while the group and others have read and execute only (5 = 4+0+1)."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Linux file permissions are the foundation of the operating system's security model. Understanding how to read and modify permissions is essential for both system administration and cybersecurity.\n\n"
        "When you run ls -l in a directory, the output looks like this:\n"
        "-rwxr-xr--  1  mulenga  staff  1024  Mar 10 10:30  script.sh\n\n"
        "Let's break down every character:\n"
        "- The first character indicates the file type. A dash (-) means a regular file. The letter d means a directory. The letter l means a symbolic link (similar to a shortcut in Windows).\n"
        "- Characters 2-4 (rwx) show the owner's permissions. In this example, the owner (mulenga) can read, write, and execute the file.\n"
        "- Characters 5-7 (r-x) show the group's permissions. The group (staff) can read and execute but cannot write (the dash means no permission).\n"
        "- Characters 8-10 (r--) show permissions for everyone else. Other users can only read the file; they cannot modify it or run it.\n\n"
        "Numeric representation: Each permission is assigned a number. Read (r) = 4. Write (w) = 2. Execute (x) = 1. No permission (-) = 0. To get the numeric value for a category, add the numbers together.\n"
        "- rwx = 4 + 2 + 1 = 7\n"
        "- r-x = 4 + 0 + 1 = 5\n"
        "- r-- = 4 + 0 + 0 = 4\n"
        "So -rwxr-xr-- is equivalent to 754 in numeric notation.\n\n"
        "Security-critical permission examples:\n"
        "- chmod 600 ~/.ssh/id_rsa: SSH private keys should NEVER be readable by anyone except the owner. 600 means only the owner has read and write access.\n"
        "- chmod 700 /home/username: A user's home directory should be private. 700 gives the owner full access and denies everyone else.\n"
        "- chmod 644 /etc/passwd: This file is designed to be world-readable (644 = owner read/write, group read, others read).\n"
        "- chmod 640 /etc/shadow: This file should be readable ONLY by root and the shadow group. 640 = owner read/write, group read, others no access.\n\n"
        "Privilege escalation through misconfigured permissions: If a script is owned by root but writable by everyone (chmod 777), any user can modify it to run malicious commands. When root runs that script - perhaps as part of a scheduled task - the attacker's code executes with root privileges. Tools like LinPEAS (Linux Privilege Escalation Awesome Script) automate the search for such misconfigurations.\n\n"
        "Special permissions: Beyond rwx, Linux has three special permission bits:\n"
        "- Set User ID (SUID): When set on an executable file, it runs with the privileges of the file owner rather than the user who executed it. Attackers hunt for misconfigured SUID binaries to escalate privileges.\n"
        "- Set Group ID (SGID): Similar to SUID but applies to the group.\n"
        "- Sticky Bit: When set on a directory, users can only delete their own files, even if the directory is world-writable. This is used on /tmp/.\n\n"
        "Abbreviations defined:\n"
        "- SUID = Set User ID\n"
        "- SGID = Set Group ID\n"
        "- SSH = Secure Shell",
        img=os.path.join(ASSETS_DIR, "linux_permissions.png")
    )
    
    add_bullet_slide(prs,
        "Topic 2 Summary: Key Takeaways",
        [
            "- The Operating System (OS) is the ultimate security gatekeeper. It manages processes, memory, files, devices, and user access.",
            "- Windows is the dominant desktop OS. Critical directories for security include System32, AppData, and Temp.",
            "- Windows Command Prompt commands like ipconfig, netstat, tasklist, systeminfo, and net user are essential for investigation.",
            "- The Principle of Least Privilege means users should have only the minimum access needed. Administrator accounts should be used sparingly.",
            "- Linux dominates servers and cybersecurity tools. Key directories include /etc/ (configuration), /var/log/ (logs), /tmp/ (temporary files), and /home/ (user directories).",
            "- Essential Linux commands include ls -la, cat, grep, ps aux, chmod, whoami, and history.",
            "- Linux file permissions use read (r=4), write (w=2), and execute (x=1) for owner, group, and others. Misconfigured permissions are a common attack vector."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "This summary ties together the Windows and Linux content. The key message is that both operating systems are essential for cybersecurity professionals to understand. Windows dominates the endpoint (desktop) world, while Linux dominates the server and security tools world.\n\n"
        "Students should leave this topic with practical skills: they should be able to open Command Prompt on Windows and run basic investigation commands. They should be able to open a Linux terminal and navigate the file system, check permissions, read logs, and understand what different permission strings mean.\n\n"
        "The Principle of Least Privilege applies equally to Windows and Linux. In both systems, running with excessive privileges increases the damage that malware can cause. Organizations that enforce least privilege and regularly audit user accounts and permissions dramatically reduce their attack surface.\n\n"
        "For the practical assessment, students should be comfortable with both Windows CMD and Linux terminal navigation, file permission interpretation, and identifying security-critical directories on both platforms."
    )
    
    out = os.path.join(BASE_DIR, "Topic_2_Operating_Systems.pptx")
    prs.save(out)
    print(f"Saved: {out}")


# ---------------------------------------------------------------------------
# TOPIC 3: PROGRAMMING LOGIC
# ---------------------------------------------------------------------------
def build_topic_3():
    prs = new_prs()
    
    add_title_slide(prs,
        "Topic 3: Introduction to Programming Logic",
        "Thinking Like a Defender and an Attacker",
        "Welcome to Topic 3: Introduction to Programming Logic. Programming is not about becoming a software developer - it is about understanding how software works so you can analyze malware, automate security tasks, spot vulnerabilities in code, and communicate effectively with developers. In this topic, we cover variables, data types, control structures (if/else and loops), functions, pseudocode, and basic Bash scripting. Every concept is explained from a cybersecurity perspective, with no prior programming experience assumed."
    )
    
    add_bullet_slide(prs,
        "Why Programming Logic Matters in Cybersecurity",
        [
            "- Automation: Security professionals write scripts to automate repetitive tasks such as parsing log files, scanning networks, and generating reports.",
            "- Malware Analysis: Malware is simply software written with malicious intent. Understanding programming logic allows you to read malware reports, understand indicators of compromise, and even reverse-engineer simple samples.",
            "- Vulnerability Analysis: Many security flaws are logic errors. Missing input validation, incorrect access checks, and flawed business logic can all lead to exploitable vulnerabilities.",
            "- Tool Development: Python is the most popular language in cybersecurity. Over 90 percent of the tools in Kali Linux are written in Python or Bash.",
            "- Career Advantage: Job surveys consistently list Python and scripting skills as the most requested technical abilities for roles like SOC Analyst, Penetration Tester, and Security Analyst."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Programming logic is the ability to think step-by-step about how a computer solves a problem. This skill is essential in cybersecurity because everything attackers and defenders do ultimately involves code.\n\n"
        "Automation in security operations: A Security Operations Center (SOC) analyst might receive tens of thousands of alerts per day. Manually reviewing each one is impossible. Instead, analysts write Python scripts to automatically parse log files, extract suspicious Internet Protocol (IP) addresses, cross-reference them with threat intelligence databases, and generate a summary report. For example, a simple script might read a firewall log, count how many times each external IP attempted to connect, and flag any IP with more than 100 attempts as a potential brute-force attacker.\n\n"
        "Malware analysis: At its core, ransomware is just a program that uses loops to iterate through files, conditionals to decide which files to encrypt, and file system Application Programming Interfaces (APIs) to overwrite data with encrypted versions. If you understand loops and conditionals, you can read technical analyses of malware and understand exactly how it works. For example, the WannaCry ransomware had a built-in 'kill switch' - it checked whether a specific domain name was registered, and if so, it stopped spreading. This was a simple if/else statement that security researchers discovered and exploited to slow the outbreak.\n\n"
        "Vulnerability analysis: Many of the most common vulnerabilities are caused by missing or incorrect logic. SQL injection occurs when a program takes user input and puts it directly into a database query without checking whether it is safe. This is a missing validation step - a logic error. Similarly, if a web application checks if (user.isAdmin) but fails to handle the case where isAdmin is undefined or null, an attacker might be able to bypass the check.\n\n"
        "Tool development: Python dominates the cybersecurity tool ecosystem because it is easy to read, has a massive library of pre-built modules, and works on both Windows and Linux. Tools like Nmap (network scanner), Metasploit (penetration testing framework), and Scapy (packet manipulation) all have Python interfaces or are written in Python. Bash scripting is equally important for automating Linux system administration tasks.\n\n"
        "Career impact: According to the 2023 Cybersecurity Workforce Study by (ISC) squared, over 60 percent of cybersecurity job postings mention Python or scripting skills. Learning even basic programming significantly improves employability.\n\n"
        "Abbreviations defined:\n"
        "- SOC = Security Operations Center\n"
        "- IP = Internet Protocol\n"
        "- API = Application Programming Interface\n"
        "- SQL = Structured Query Language\n"
        "- (ISC) squared = International Information System Security Certification Consortium"
    )
    
    add_bullet_slide(prs,
        "Variables and Data Types",
        [
            "- A variable is a named container that stores a value. Think of it as a labeled box where you can store information.",
            "- String: A sequence of characters used for text. Examples in cybersecurity include usernames, passwords, domain names, and Internet Protocol (IP) addresses.",
            "- Integer: A whole number without decimals. Examples include port numbers (such as 22 for Secure Shell, 80 for HyperText Transfer Protocol, and 443 for HyperText Transfer Protocol Secure) and failed login attempt counters.",
            "- Float: A number with a decimal point. Examples include success rates, response times, and risk scores.",
            "- Boolean: A value that is either True or False. Examples include access_granted, account_locked, and is_admin.",
            "- List (or Array): An ordered collection of values. Examples include a list of suspicious IP addresses or a list of blocked domains."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Variables are the foundation of every program. When you write a script to analyze security data, you use variables to store and manipulate that data.\n\n"
        "String: In Python, strings are created by enclosing text in single or double quotes. For example: username = 'admin' or password = 'SecurePass123!'. In cybersecurity, strings are everywhere - in log files, network packets, and database queries. A common security task is parsing strings to extract useful information. For example, a log entry might look like: '192.168.1.50 - - [10/Mar/2026:14:32:01] GET /login HTTP/1.1'. A script would use string manipulation to extract the IP address (192.168.1.50), the timestamp, and the requested page (/login).\n\n"
        "Integer: Integers are used for counting and indexing. In networking, port numbers are integers in the range 0 to 65535. The well-known ports (0 to 1023) are reserved for standard services and require administrator privileges to use on Linux and Unix systems. For example, port 22 is Secure Shell (SSH), port 80 is HyperText Transfer Protocol (HTTP), and port 443 is HyperText Transfer Protocol Secure (HTTPS). Integer overflow is a serious vulnerability: if a program uses an integer to track the size of something, and that size grows beyond the maximum value the integer can hold, it wraps around to a negative number or zero. This can cause the program to allocate too little memory, leading to a buffer overflow.\n\n"
        "Float: Floats represent real numbers with decimal points. In security analytics, you might calculate the percentage of failed logins (failed / total * 100.0) or the average response time of a server.\n\n"
        "Boolean: Booleans have only two possible values: True or False. Authentication systems rely heavily on booleans. For example, after a user enters a password, the system might set is_authenticated = True. If the password is wrong, it remains False. A logic flaw occurs when the program does not properly check this value. For instance, in some weakly typed languages, a null (empty) value might be treated differently than False, potentially allowing an attacker to bypass checks.\n\n"
        "List: A list is an ordered collection of items. In Python, lists are created with square brackets. For example: suspicious_ips = ['192.168.1.100', '10.0.0.50', '172.16.0.25']. Lists are essential for automation - a penetration tester might have a list of target IP addresses to scan, and a script would loop through the list, scanning each one. However, in secure coding, iterating over a list of user inputs without proper validation can lead to injection attacks if any item in the list contains malicious data.\n\n"
        "Abbreviations defined:\n"
        "- IP = Internet Protocol\n"
        "- SSH = Secure Shell\n"
        "- HTTP = HyperText Transfer Protocol\n"
        "- HTTPS = HyperText Transfer Protocol Secure"
    )
    
    add_two_col_slide(prs,
        "Variables in Action: Python Examples",
        [
            "# String variables",
            "student_name = 'Chanda Mwale'",
            "password = 'Zambia@2026'",
            "target_ip = '192.168.1.50'",
            "",
            "# Integer variables",
            "port_number = 22",
            "failed_attempts = 0",
            "max_attempts = 3",
            "",
            "# Boolean variable",
            "access_granted = False"
        ],
        [
            "# Float variable",
            "risk_score = 7.5",
            "",
            "# List variable",
            "blocked_ips = [",
            "    '10.0.0.99',",
            "    '192.168.1.200',",
            "    '172.16.0.15'",
            "]",
            "",
            "# Print a variable",
            "print('Target:', target_ip)",
            "print('Port:', port_number)"
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "This slide shows concrete Python syntax for declaring variables. In Python, you do not need to declare the data type explicitly - the interpreter figures it out based on the value you assign.\n\n"
        "String variables are enclosed in quotes. The student_name variable stores the text 'Chanda Mwale'. The password variable stores a string. In real applications, passwords should NEVER be stored in plain text like this - they should be hashed using algorithms like the Secure Hash Algorithm 256 (SHA-256) or the bcrypt algorithm. However, for learning basic syntax, this example is fine. The target_ip variable stores an IP address as a string.\n\n"
        "Integer variables store whole numbers. port_number = 22 assigns the integer 22 to the variable port_number. failed_attempts = 0 initializes a counter to zero. This is a common pattern in security scripts - you start a counter at zero and increment it each time an event occurs (like a failed login).\n\n"
        "Boolean variables store True or False. access_granted = False means the user has not yet been authenticated. Later in the program, after the password is checked, this variable might be changed to True.\n\n"
        "Float variables store decimal numbers. risk_score = 7.5 might represent the calculated risk level of a particular user or event on a scale of 1 to 10.\n\n"
        "List variables store multiple values. blocked_ips is a list of three IP addresses. You can access individual items using an index, starting from zero. blocked_ips[0] would give '10.0.0.99', blocked_ips[1] would give '192.168.1.200', and so on. You can also add items to a list using the append method: blocked_ips.append('192.168.1.250').\n\n"
        "The print function displays output to the screen. In Python 3, print is a function that takes the text to display as an argument. The output of print('Target:', target_ip) would be: Target: 192.168.1.50\n\n"
        "Abbreviations defined:\n"
        "- IP = Internet Protocol\n"
        "- SHA-256 = Secure Hash Algorithm 256"
    )
    
    add_bullet_slide(prs,
        "Control Structures: If/Else Statements",
        [
            "- Programs need to make decisions based on conditions. The if/else statement is the most basic decision-making structure in programming.",
            "- The program checks a condition. If the condition is true, it runs one block of code. If the condition is false, it runs a different block (or skips the action entirely).",
            "- In cybersecurity, if/else logic is used everywhere: checking passwords, validating user input, enforcing access controls, and triggering alerts.",
            "",
            "Example - Password Check:",
            "  password = input('Enter password: ')",
            "  if password == 'Secure@2026':",
            "      print('Access granted')",
            "  else:",
            "      print('Access denied')"
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Control structures determine the flow of a program. Without them, a program would simply execute every line from top to bottom in order. If/else statements allow programs to adapt their behavior based on circumstances.\n\n"
        "How if/else works: The keyword if is followed by a condition - an expression that evaluates to either True or False. If the condition is True, the indented block of code under the if statement runs. If the condition is False, the indented block under the else statement runs instead. Python uses indentation (spaces at the beginning of lines) to define blocks of code. Standard practice is to use four spaces per indentation level.\n\n"
        "The example on this slide asks the user to enter a password. The input() function pauses the program and waits for the user to type something and press Enter. The text entered by the user is stored in the variable password. The condition password == 'Secure@2026' uses the double equals sign (==), which is the comparison operator. It checks whether the value of password is exactly equal to the string 'Secure@2026'. A single equals sign (=) is the assignment operator, used to give a value to a variable. This is a common beginner mistake - using = instead of == in a condition.\n\n"
        "Real-world cybersecurity applications of if/else:\n"
        "- Authentication: if (password_hash == stored_hash) { grant_access() } else { deny_access() }. This is the core logic of every login system.\n"
        "- Authorization: if (user.role == 'admin') { show_admin_panel() } else { show_standard_dashboard() }. This determines what a user is allowed to see and do.\n"
        "- Input validation: if (len(user_input) > 1000) { reject_input() } else { process_input() }. This prevents buffer overflow by rejecting oversized input.\n"
        "- Alerting: if (failed_logins > 5) { send_alert_to_admin() } else { allow_retry() }. This is how intrusion detection systems identify brute-force attacks.\n\n"
        "Famous logic bug example: In 2019, Apple discovered a critical bug in FaceTime (their video calling app). Due to an incorrectly structured if statement, a caller could connect to another person's microphone and camera BEFORE the recipient answered the call. This was purely a logic error - no memory corruption or sophisticated hacking was involved. It demonstrates how dangerous logic bugs can be.\n\n"
        "Abbreviations defined:\n"
        "- API = Application Programming Interface (referenced in context of input function)"
    )
    
    add_bullet_slide(prs,
        "Control Structures: Loops",
        [
            "- A loop allows a program to repeat a block of code multiple times. This is essential for tasks that involve repetition, such as checking multiple login attempts or scanning a list of network ports.",
            "- A while loop continues running as long as a specified condition remains true.",
            "",
            "Example - Account Lockout Policy:",
            "  attempts = 0",
            "  while attempts < 3:",
            "      password = input('Enter password: ')",
            "      if password == 'Secure@2026':",
            "          print('Access granted')",
            "          break",
            "      else:",
            "          attempts = attempts + 1",
            "          print('Wrong password. Attempt', attempts, 'of 3')",
            "  if attempts == 3:",
            "      print('Account locked. Too many failed attempts.')"
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Loops are one of the most powerful concepts in programming. They allow a small amount of code to handle large amounts of work. In cybersecurity, loops are used constantly.\n\n"
        "How a while loop works: The keyword while is followed by a condition. Before each iteration (each repetition), the program checks whether the condition is True. If it is, the indented block of code runs. After the block finishes, the condition is checked again. This repeats until the condition becomes False. It is critically important that something inside the loop eventually changes the condition to False. Otherwise, the loop will run forever - this is called an infinite loop.\n\n"
        "Let's trace through the account lockout example step by step:\n"
        "1. attempts = 0: We create a variable called attempts and set it to 0.\n"
        "2. while attempts < 3: The condition is 0 < 3, which is True. So we enter the loop.\n"
        "3. password = input('Enter password: '): The program waits for the user to type a password.\n"
        "4. if password == 'Secure@2026': The program checks if the password is correct.\n"
        "   - If correct: It prints 'Access granted' and then encounters the break statement. The break statement immediately exits the loop, skipping everything else.\n"
        "   - If incorrect: It runs the else block. attempts = attempts + 1 increases the counter by 1. If this was the first failure, attempts is now 1. The program prints a warning message.\n"
        "5. The loop goes back to step 2 and checks attempts < 3 again.\n"
        "6. After three failed attempts, attempts becomes 3. The condition 3 < 3 is now False, so the loop ends.\n"
        "7. if attempts == 3: Since attempts is 3, the program prints 'Account locked. Too many failed attempts.'\n\n"
        "This example mirrors exactly how real-world account lockout policies work. After a configurable number of failed login attempts (often three to five), the account is temporarily or permanently locked to prevent brute-force password guessing.\n\n"
        "Other security uses of loops:\n"
        "- Port scanning: A tool like Nmap uses loops to test thousands of ports. Pseudocode: for port in range(1, 65536): scan_port(port).\n"
        "- Password cracking: A dictionary attack loops through a list of common passwords, trying each one. Pseudocode: for word in wordlist: try_login(word).\n"
        "- Log analysis: A script might loop through millions of log lines, checking each one for suspicious patterns. Pseudocode: for line in logfile: analyze(line).\n\n"
        "Infinite loops as a security threat: If an attacker can trigger an infinite loop in a web application or network service, they can cause a Denial of Service (DoS). The service consumes 100 percent of the CPU and stops responding to legitimate users. Input validation and proper loop termination are essential defenses.\n\n"
        "Abbreviations defined:\n"
        "- DoS = Denial of Service\n"
        "- CPU = Central Processing Unit"
    )
    
    add_bullet_slide(prs,
        "Functions: Reusable Blocks of Code",
        [
            "- A function is a named block of code that performs a specific task. Once defined, it can be called (used) multiple times without rewriting the code.",
            "- Functions make programs easier to read, maintain, and debug.",
            "- In cybersecurity, functions are used to create reusable security checks, such as validating passwords, formatting log entries, or calculating risk scores.",
            "",
            "Example - Password Strength Checker:",
            "  def check_password_strength(password):",
            "      if len(password) < 8:",
            "          return 'Weak - too short'",
            "      elif password.isalpha():",
            "          return 'Weak - letters only'",
            "      else:",
            "          return 'Strong'",
            "",
            "  result = check_password_strength('Zambia@2026')",
            "  print(result)  # Output: Strong"
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Functions are essential for writing clean, maintainable, and secure code. They allow you to define a security check once and reuse it throughout your program.\n\n"
        "How functions work in Python: A function is defined using the def keyword, followed by the function name, parentheses containing any inputs (called parameters or arguments), and a colon. The indented block below is the function body - the code that runs when the function is called. The return statement sends a result back to the code that called the function.\n\n"
        "In the password strength example:\n"
        "- def check_password_strength(password): defines a function named check_password_strength that takes one input called password.\n"
        "- len(password) is a built-in function that returns the number of characters in the string. If the password has fewer than 8 characters, the function immediately returns 'Weak - too short'.\n"
        "- password.isalpha() is a string method that returns True if all characters in the password are letters (A-Z, a-z). If so, the function returns 'Weak - letters only' because the password lacks numbers and symbols.\n"
        "- If neither of the above conditions is true, the password is at least 8 characters long and contains at least one non-letter character, so the function returns 'Strong'.\n"
        "- result = check_password_strength('Zambia@2026') calls the function with the string 'Zambia@2026'. The returned value ('Strong') is stored in the variable result.\n"
        "- print(result) displays the output.\n\n"
        "Password policy best practices: The United States National Institute of Standards and Technology (NIST) publishes guidelines for password security in Special Publication 800-63B. Modern recommendations emphasize password length over complexity. A long passphrase like 'Correct-Horse-Battery-Staple!' is more secure and easier to remember than a short complex password like 'P@ssw0rd1'. NIST recommends a minimum length of 8 characters but encourages even longer passwords. They also advise against forcing users to change passwords frequently unless there is evidence of compromise, because frequent changes lead to predictable patterns (like Password1, Password2, Password3).\n\n"
        "Other security uses of functions:\n"
        "- Hashing a password before storing it in a database.\n"
        "- Validating that an IP address is in the correct format.\n"
        "- Encrypting a message before sending it over a network.\n"
        "- Parsing a log line and extracting the timestamp, source IP, and event type.\n\n"
        "Abbreviations defined:\n"
        "- NIST = National Institute of Standards and Technology"
    )
    
    add_bullet_slide(prs,
        "Pseudocode: Planning Before Coding",
        [
            "- Pseudocode is a plain-language description of the steps a program will take, written before any actual code.",
            "- It helps programmers think through the logic without worrying about the exact syntax of a programming language.",
            "- Pseudocode is widely used by security architects to design authentication flows, access control policies, and incident response playbooks.",
            "",
            "Example - Login System with Account Lockout:",
            "  START",
            "    INPUT username",
            "    INPUT password",
            "    IF username exists in database THEN",
            "        IF password matches stored hash THEN",
            "            GRANT access",
            "            LOG successful login",
            "        ELSE",
            "            INCREMENT failed_attempts",
            "            IF failed_attempts >= 3 THEN",
            "                LOCK account",
            "                SEND alert to administrator",
            "            END IF",
            "        END IF",
            "    ELSE",
            "        DISPLAY 'Invalid credentials'",
            "    END IF",
            "  END"
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Pseudocode bridges the gap between human thinking and computer programming. It uses plain English (or any natural language) combined with programming-like structures such as IF, THEN, ELSE, WHILE, and FOR. Because pseudocode is not tied to any specific programming language, it can be understood by programmers, security analysts, managers, and auditors alike.\n\n"
        "Why pseudocode matters in cybersecurity:\n"
        "- Security architects use pseudocode to design authentication systems before developers write a single line of real code. This helps catch logic flaws early.\n"
        "- Incident response playbooks are often written in pseudocode or flowchart form to ensure everyone follows the same steps during a crisis.\n"
        "- Compliance auditors may review pseudocode descriptions of access control logic to verify that an application meets regulatory requirements.\n\n"
        "Analyzing the login system pseudocode:\n"
        "1. START and END mark the beginning and end of the algorithm.\n"
        "2. INPUT username and INPUT password represent the user providing their credentials.\n"
        "3. The outer IF checks whether the username exists in the database. This prevents the system from revealing whether the username or password was wrong - a security best practice called 'ambiguous error messages.' If an attacker knows a username exists, they can focus on cracking just the password.\n"
        "4. The inner IF checks whether the password matches the stored hash. In real systems, passwords are NEVER stored in plain text. Instead, the system stores a cryptographic hash of the password. When the user logs in, the system hashes the entered password and compares the two hashes.\n"
        "5. If the password is wrong, failed_attempts is incremented (increased by one).\n"
        "6. If failed_attempts reaches 3, the account is locked and an alert is sent to the administrator. This is the account lockout policy.\n\n"
        "Real-world consideration: Account lockout policies must balance security and usability. If the threshold is too low (like 2 attempts), users may be locked out frequently due to typos. If the threshold is too high (like 50 attempts), attackers have more opportunities to guess passwords. A common setting is 3 to 5 failed attempts, with a temporary lockout period of 15 to 30 minutes. Some systems also implement CAPTCHA (Completely Automated Public Turing test to tell Computers and Humans Apart) after a few failed attempts to slow down automated guessing.\n\n"
        "Abbreviations defined:\n"
        "- CAPTCHA = Completely Automated Public Turing test to tell Computers and Humans Apart"
    )
    
    add_bullet_slide(prs,
        "Bash Scripting: Automating Security Tasks",
        [
            "- Bash (Bourne Again Shell) is the default command interpreter on most Linux systems.",
            "- A Bash script is a text file containing multiple commands that can be executed together.",
            "- Bash scripting is essential for automating repetitive tasks such as system health checks, log analysis, and backup routines.",
            "",
            "Example - Simple Network Health Check Script:",
            "  #!/bin/bash",
            "  echo 'Checking network connectivity...'",
            "  ping -c 4 8.8.8.8",
            "  echo ''",
            "  echo 'Your IP address:'",
            "  ip addr | grep 'inet ' | grep -v '127.0.0.1'",
            "  echo ''",
            "  echo 'Open network connections:'",
            "  netstat -tuln"
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Bash is both a command language and a scripting language. When you type commands into a Linux terminal, you are talking to Bash. A Bash script is simply a file that collects multiple Bash commands so they can be run together.\n\n"
        "The first line of a Bash script is always the shebang: #!/bin/bash. This tells the operating system which interpreter (in this case, Bash) should execute the script. Without this line, the system might try to run the script with the wrong program.\n\n"
        "The echo command prints text to the screen. It is the Bash equivalent of Python's print function.\n\n"
        "ping -c 4 8.8.8.8 sends four network test packets to the IP address 8.8.8.8, which is one of Google's public DNS servers. The -c 4 flag means 'send exactly 4 packets.' If the responses come back successfully, it means the computer has internet connectivity. If they time out, there may be a network problem.\n\n"
        "ip addr | grep 'inet ' | grep -v '127.0.0.1' displays the computer's IP addresses, excluding the loopback address (127.0.0.1). The pipe symbol (|) takes the output of one command and feeds it as input to the next command. This is called piping. First, ip addr shows all network interfaces. Then, grep 'inet ' filters for lines containing IP addresses. Finally, grep -v '127.0.0.1' removes the loopback line.\n\n"
        "netstat -tuln shows all listening network connections. The flags mean: -t = TCP, -u = UDP, -l = listening ports only, -n = show numeric addresses instead of resolving names. This command is useful for auditing which services are exposed on the machine. Unexpected open ports may indicate unauthorized software or backdoors.\n\n"
        "To run this script, you would save it as check_network.sh, make it executable with chmod +x check_network.sh, and then run it with ./check_network.sh.\n\n"
        "Security warning: Bash scripts that accept user input without validation are vulnerable to command injection. If a script takes user input and passes it directly to a command, an attacker can inject additional commands. For example, if a script runs ping $user_input and the attacker provides 8.8.8.8; rm -rf /, the script might execute the destructive rm command. Always validate and sanitize inputs in scripts.\n\n"
        "Abbreviations defined:\n"
        "- DNS = Domain Name System\n"
        "- TCP = Transmission Control Protocol\n"
        "- UDP = User Datagram Protocol\n"
        "- IP = Internet Protocol"
    )
    
    add_bullet_slide(prs,
        "Topic 3 Summary: Key Takeaways",
        [
            "- Programming logic enables cybersecurity professionals to automate tasks, analyze malware, spot vulnerabilities, and develop tools.",
            "- Variables store data. Common types include strings (text), integers (whole numbers), floats (decimals), booleans (True/False), and lists (collections).",
            "- If/else statements allow programs to make decisions based on conditions. They are the foundation of access control and input validation.",
            "- Loops allow programs to repeat actions. They are essential for tasks like scanning ports, trying password lists, and processing logs.",
            "- Functions create reusable blocks of code, making programs easier to maintain and security checks easier to apply consistently.",
            "- Pseudocode is a planning tool that describes program logic in plain language before writing actual code.",
            "- Bash scripting automates Linux tasks. Always validate inputs to prevent command injection attacks."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "This summary reinforces that programming is a practical skill for cybersecurity, not just an academic exercise. Students do not need to become software engineers, but they do need to understand how code works so they can 'think like both an attacker and a defender.'\n\n"
        "Key connections to cybersecurity careers:\n"
        "- SOC Analysts use Python to parse logs and automate alert triage.\n"
        "- Penetration testers write scripts to automate scanning and exploitation.\n"
        "- Malware analysts read code to understand how threats behave.\n"
        "- Security engineers write scripts to enforce policies and check compliance.\n\n"
        "For the practical assessment, students should be able to write simple pseudocode for a login system, explain what variables and data types are, and understand how if/else and loops control program flow. If time permits, having them write and run a simple Bash script on Kali Linux or an Ubuntu virtual machine is an excellent hands-on exercise."
    )
    
    out = os.path.join(BASE_DIR, "Topic_3_Programming_Logic.pptx")
    prs.save(out)
    print(f"Saved: {out}")


# ---------------------------------------------------------------------------
# TOPIC 4: MATHEMATICS FOR CYBERSECURITY
# ---------------------------------------------------------------------------
def build_topic_4():
    prs = new_prs()
    
    add_title_slide(prs,
        "Topic 4: Mathematics for Cybersecurity",
        "Binary, Hexadecimal, Boolean Logic, and XOR",
        "Welcome to Topic 4: Mathematics for Cybersecurity. Mathematics is the hidden engine behind almost every security technology you use: encryption, hashing, digital signatures, network addressing, and access control. This topic focuses on the specific math concepts that cybersecurity professionals need to understand: binary and hexadecimal number systems, Boolean logic, and the Exclusive OR (XOR) operation. We will explain every concept clearly, show how it is used in real security systems, and provide plenty of practical examples."
    )
    
    add_bullet_slide(prs,
        "The Binary Number System",
        [
            "- Binary is a base-2 number system that uses only two digits: 0 and 1.",
            "- Each binary digit is called a bit. Eight bits make one byte.",
            "- Computers store all data - text, images, passwords, and network packets - as binary.",
            "- Understanding binary is essential for interpreting memory dumps, network packet captures, and cryptographic hashes.",
            "",
            "Decimal to Binary Examples:",
            "- Decimal 0 = Binary 0000",
            "- Decimal 1 = Binary 0001",
            "- Decimal 2 = Binary 0010",
            "- Decimal 5 = Binary 0101",
            "- Decimal 10 = Binary 1010",
            "- Decimal 255 = Binary 11111111"
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Binary is the native language of computers. At the hardware level, a computer's Central Processing Unit (CPU) and memory chips represent information using electrical signals: a high voltage typically means 1, and a low voltage (or no voltage) means 0. Because there are only two states, the base-2 binary system is the perfect mathematical model for how computers store and process information.\n\n"
        "A bit (short for binary digit) is the smallest unit of data in computing. It can hold exactly one of two values: 0 or 1. A byte is a group of 8 bits. One byte can represent 256 different values (2 to the power of 8), ranging from 0 to 255. This is why an IPv4 address uses four bytes (four numbers from 0 to 255), and why the American Standard Code for Information Interchange (ASCII) character set uses one byte per character.\n\n"
        "How to convert decimal to binary: The easiest method for small numbers is to find the largest power of 2 that fits into the number, subtract it, and repeat.\n"
        "Example: Convert decimal 13 to binary.\n"
        "- The powers of 2 are: 8, 4, 2, 1.\n"
        "- Does 8 fit into 13? Yes. Write 1. Remainder = 13 - 8 = 5.\n"
        "- Does 4 fit into 5? Yes. Write 1. Remainder = 5 - 4 = 1.\n"
        "- Does 2 fit into 1? No. Write 0. Remainder stays 1.\n"
        "- Does 1 fit into 1? Yes. Write 1. Remainder = 0.\n"
        "- Reading the results: 1101. So decimal 13 = binary 1101.\n\n"
        "Real-world cybersecurity context:\n"
        "- Subnet masks are expressed in binary. A /24 subnet mask means the first 24 bits are 1s and the last 8 bits are 0s: 11111111.11111111.11111111.00000000, which is 255.255.255.0 in decimal.\n"
        "- File permissions in Linux are often represented in octal (base-8), which is derived from binary groupings of three bits.\n"
        "- Cryptographic keys are measured in bits. A 256-bit key has 2^256 possible combinations, making it computationally infeasible to guess.\n"
        "- Memory forensics involves analyzing raw binary data extracted from RAM chips to find artifacts like passwords, encryption keys, and malware.\n\n"
        "Abbreviations defined:\n"
        "- CPU = Central Processing Unit\n"
        "- ASCII = American Standard Code for Information Interchange\n"
        "- IPv4 = Internet Protocol version 4\n"
        "- RAM = Random Access Memory"
    )
    
    add_bullet_slide(prs,
        "Binary Addition and Overflow",
        [
            "- Binary addition follows the same rules as decimal addition, but with only two digits.",
            "- Rules: 0 + 0 = 0, 0 + 1 = 1, 1 + 0 = 1, 1 + 1 = 10 (write 0, carry 1).",
            "- Example: Add 0101 (5) and 0011 (3). Result = 1000 (8).",
            "- Overflow occurs when the result requires more bits than the container can hold.",
            "- Integer overflow is a real-world vulnerability that has caused crashes, privilege escalation, and remote code execution in software."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Binary addition is a fundamental operation in computing. Even though humans work in decimal, computers perform all arithmetic in binary.\n\n"
        "Binary addition rules:\n"
        "- 0 + 0 = 0\n"
        "- 0 + 1 = 1\n"
        "- 1 + 0 = 1\n"
        "- 1 + 1 = 10 (which is 2 in decimal). Write down 0 and carry 1 to the next column.\n"
        "- 1 + 1 + 1 (with a carried 1) = 11 (which is 3 in decimal). Write down 1 and carry 1.\n\n"
        "Example: Add 0101 (5) and 0011 (3).\n"
        "  0101\n"
        "+ 0011\n"
        "------\n"
        "  1000\n\n"
        "Let's do it column by column from right to left:\n"
        "- Rightmost column: 1 + 1 = 10. Write 0, carry 1.\n"
        "- Next column: 0 + 1 + carry 1 = 10. Write 0, carry 1.\n"
        "- Next column: 1 + 0 + carry 1 = 10. Write 0, carry 1.\n"
        "- Leftmost column: 0 + 0 + carry 1 = 1. Write 1.\n"
        "- Final result: 1000, which is 8 in decimal.\n\n"
        "Integer overflow: Computers allocate a fixed number of bits to store numbers. For example, an unsigned 8-bit integer can hold values from 0 to 255. If you try to add 255 + 1 in an 8-bit container, the result is 256, but 256 in binary is 1 0000 0000 (9 bits). Since only 8 bits are available, the leading 1 is discarded, and the result wraps around to 0. This is called overflow.\n\n"
        "Security impact of integer overflow: Attackers have exploited integer overflow to cause devastating security problems. In 2014, the 'Heartbleed' vulnerability in OpenSSL allowed attackers to read sensitive memory from servers - though Heartbleed itself was a buffer over-read, it was closely related to length miscalculation. More directly, integer overflows in image parsers, file format handlers, and network protocols have allowed attackers to bypass size checks, allocate too-small memory buffers, and then overflow those buffers with malicious data. The Stagefright vulnerability in Android (2015) was partly caused by integer underflow (the opposite: subtracting too much and wrapping to a huge positive number).\n\n"
        "Mitigation: Modern programming languages and compilers include protections against integer overflow. However, in languages like C and C++, it is still the programmer's responsibility to check for overflow before performing arithmetic. Secure coding standards require that all arithmetic on untrusted input be bounds-checked.\n\n"
        "Abbreviations defined:\n"
        "- OpenSSL = Open Secure Sockets Layer toolkit\n"
        "- Android = Mobile operating system developed by Google"
    )
    
    add_bullet_slide(prs,
        "The Hexadecimal Number System",
        [
            "- Hexadecimal (hex) is a base-16 number system. It uses 16 symbols: 0-9 and A-F.",
            "- Hexadecimal is used as a compact way to represent binary data.",
            "- One hex digit represents exactly 4 bits (a nibble). Two hex digits represent one byte (8 bits).",
            "",
            "Decimal to Hexadecimal Examples:",
            "- Decimal 0 = Hex 00",
            "- Decimal 10 = Hex 0A",
            "- Decimal 15 = Hex 0F",
            "- Decimal 16 = Hex 10",
            "- Decimal 255 = Hex FF",
            "- Decimal 256 = Hex 100"
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Hexadecimal is indispensable in cybersecurity because it provides a compact, human-readable way to represent long sequences of binary data. Without hex, a 256-bit cryptographic key would be a 77-digit decimal number or a 256-digit binary number - both impractical to write down or compare. In hex, the same key is a 64-character string, which is much more manageable.\n\n"
        "Hexadecimal uses 16 digits: 0, 1, 2, 3, 4, 5, 6, 7, 8, 9, A, B, C, D, E, F. The letters A through F represent decimal values 10 through 15.\n\n"
        "The relationship between hex and binary is elegant because 16 is 2 to the power of 4. This means every hex digit maps directly to exactly 4 binary digits (a nibble). Every pair of hex digits maps to exactly 1 byte (8 bits). This makes conversion between hex and binary very easy:\n"
        "- Hex F = Binary 1111\n"
        "- Hex A = Binary 1010\n"
        "- Hex 0 = Binary 0000\n"
        "- Hex FF = Binary 11111111 (one byte)\n"
        "- Hex 1A = Binary 00011010 (one byte)\n\n"
        "Real-world cybersecurity uses of hex:\n"
        "- MAC (Media Access Control) addresses are written in hex, like 00:1A:2B:3C:4D:5E.\n"
        "- IPv6 addresses use hex, like 2001:0DB8:85A3::8A2E:0370:7334.\n"
        "- Memory dumps and packet captures display raw bytes in hex. Tools like Wireshark show packet contents in both hex and ASCII.\n"
        "- Cryptographic hashes are represented in hex. For example, an SHA-256 hash is a 64-character hex string.\n"
        "- Color codes in web development use hex, like #FF5733 for a shade of red.\n\n"
        "How to convert decimal to hex: Divide the decimal number by 16 repeatedly and record the remainders.\n"
        "Example: Convert decimal 255 to hex.\n"
        "- 255 divided by 16 = 15 with a remainder of 15.\n"
        "- 15 in hex is F.\n"
        "- So 255 in decimal = FF in hex.\n\n"
        "Example: Convert decimal 173 to hex.\n"
        "- 173 divided by 16 = 10 with a remainder of 13.\n"
        "- 10 in hex is A. 13 in hex is D.\n"
        "- So 173 in decimal = AD in hex.\n\n"
        "Abbreviations defined:\n"
        "- Hex = Hexadecimal\n"
        "- MAC = Media Access Control\n"
        "- IPv6 = Internet Protocol version 6\n"
        "- ASCII = American Standard Code for Information Interchange\n"
        "- SHA-256 = Secure Hash Algorithm 256",
        img=os.path.join(ASSETS_DIR, "binary_hex_table.png")
    )
    
    add_bullet_slide(prs,
        "Boolean Logic in Computing",
        [
            "- Boolean logic deals with values that are either True (1) or False (0).",
            "- The three fundamental Boolean operations are AND, OR, and NOT.",
            "- AND: The output is True only if BOTH inputs are True.",
            "- OR: The output is True if AT LEAST ONE input is True.",
            "- NOT: The output is the opposite of the input.",
            "- Boolean logic is the foundation of all digital circuits, programming conditions, and access control rules."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Boolean logic, named after the mathematician George Boole (1815-1864), is the mathematical system that underlies all digital computing. Every Central Processing Unit (CPU), every software condition, and every database query filter relies on Boolean operations.\n\n"
        "AND operation (symbol: && or AND):\n"
        "- True AND True = True\n"
        "- True AND False = False\n"
        "- False AND True = False\n"
        "- False AND False = False\n"
        "In cybersecurity, AND is used in access control. For example: if (user.isEmployee == True) AND (user.hasSecurityClearance == True) then grantAccess(). Both conditions must be satisfied.\n\n"
        "OR operation (symbol: || or OR):\n"
        "- True OR True = True\n"
        "- True OR False = True\n"
        "- False OR True = True\n"
        "- False OR False = False\n"
        "In cybersecurity, OR is used in firewall rules. For example: allow traffic if (sourcePort == 80) OR (sourcePort == 443). If either condition is true, the traffic is allowed.\n\n"
        "NOT operation (symbol: ! or NOT):\n"
        "- NOT True = False\n"
        "- NOT False = True\n"
        "In cybersecurity, NOT is used to deny access. For example: if NOT (user.isBanned) then allowLogin().\n\n"
        "Truth tables: A truth table is a chart that lists all possible input combinations and their corresponding outputs for a Boolean operation. Engineers use truth tables to design logic gates - the physical building blocks of computer chips. A logic gate is a tiny electronic circuit that performs a Boolean operation. Modern CPUs contain billions of logic gates working together at incredible speeds.\n\n"
        "Boolean logic in SQL injection: Attackers sometimes exploit Boolean-based SQL injection to extract data one bit at a time. They craft queries that return True or False based on a condition. By observing whether the application behaves differently for True versus False, they can slowly reconstruct secret information like database contents.\n\n"
        "Abbreviations defined:\n"
        "- CPU = Central Processing Unit\n"
        "- SQL = Structured Query Language"
    )
    
    add_bullet_slide(prs,
        "Boolean Truth Tables",
        [
            "",
            "AND Table          OR Table           NOT Table",
            "A B | Output       A B | Output       A | Output",
            "----------------    ----------------   --------------",
            "0 0 | 0            0 0 | 0            0 | 1",
            "0 1 | 0            0 1 | 1            1 | 0",
            "1 0 | 0            1 0 | 1",
            "1 1 | 1            1 1 | 1",
            "",
            "Example Security Application:",
            "- Access requires: (hasBadge = True) AND (isInsideHours = True)",
            "- Only when BOTH conditions are True is access granted.",
            "- This is how multi-factor authentication and dual-authorization systems work."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "This slide presents the truth tables for AND, OR, and NOT in a clean tabular format. Instructors should walk through each row slowly, especially for students who have never seen truth tables before.\n\n"
        "AND truth table walkthrough:\n"
        "- Row 1: Input A=0, Input B=0. Since neither is 1, the output is 0 (False).\n"
        "- Row 2: Input A=0, Input B=1. A is not True, so the output is 0.\n"
        "- Row 3: Input A=1, Input B=0. B is not True, so the output is 0.\n"
        "- Row 4: Input A=1, Input B=1. Both are True, so the output is 1 (True).\n\n"
        "OR truth table walkthrough:\n"
        "- Row 1: Input A=0, Input B=0. Neither is True, so output is 0.\n"
        "- Row 2: Input A=0, Input B=1. B is True, so output is 1.\n"
        "- Row 3: Input A=1, Input B=0. A is True, so output is 1.\n"
        "- Row 4: Input A=1, Input B=1. Both are True, so output is 1.\n\n"
        "NOT truth table walkthrough:\n"
        "- Row 1: Input A=0. The opposite of False is True, so output is 1.\n"
        "- Row 2: Input A=1. The opposite of True is False, so output is 0.\n\n"
        "Security example - Multi-factor authentication (MFA): MFA requires users to provide at least two different types of authentication factors. For example, 'something you know' (password) AND 'something you have' (mobile phone with a one-time code). In Boolean terms: accessGranted = (correctPassword) AND (validOTP). Even if an attacker steals the password (A=1, B=0), the AND operation returns False, and access is denied.\n\n"
        "Security example - Firewall rule: A firewall might allow traffic if it comes from the internal network OR if it uses the HTTPS protocol. Internal network = True, other network = False. HTTPS = True, other protocol = False. If either is True, the OR operation allows the traffic.\n\n"
        "Abbreviations defined:\n"
        "- MFA = Multi-Factor Authentication\n"
        "- OTP = One-Time Password\n"
        "- HTTPS = HyperText Transfer Protocol Secure"
    )
    
    add_bullet_slide(prs,
        "Exclusive OR (XOR): The Heart of Cryptography",
        [
            "- Exclusive OR (XOR) is a Boolean operation that outputs True (1) only when the inputs are DIFFERENT.",
            "- XOR Rules: 0 XOR 0 = 0, 0 XOR 1 = 1, 1 XOR 0 = 1, 1 XOR 1 = 0.",
            "- The XOR operation is reversible. If A XOR B = C, then C XOR B = A.",
            "- This reversibility makes XOR the fundamental building block of many encryption algorithms, including the Data Encryption Standard (DES), Triple DES (3DES), Advanced Encryption Standard (AES), and stream ciphers.",
            "- XOR is also used in hash functions and error detection codes."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "XOR is arguably the single most important Boolean operation in all of cryptography. While AND and OR are useful for general logic, XOR has a unique mathematical property that makes it perfect for encryption: it is its own inverse.\n\n"
        "What does 'its own inverse' mean? If you XOR a plaintext message (P) with a secret key (K), you get ciphertext (C): P XOR K = C. Then, if you XOR the ciphertext (C) with the same key (K), you get the original plaintext back: C XOR K = P. This property is called self-inverse or involution.\n\n"
        "Let's verify this with a simple example:\n"
        "- Plaintext bit: 1\n"
        "- Key bit: 0\n"
        "- Ciphertext = 1 XOR 0 = 1\n"
        "- Decryption = 1 XOR 0 = 1 (original plaintext restored)\n\n"
        "Another example:\n"
        "- Plaintext bit: 1\n"
        "- Key bit: 1\n"
        "- Ciphertext = 1 XOR 1 = 0\n"
        "- Decryption = 0 XOR 1 = 1 (original plaintext restored)\n\n"
        "This works because for any bits A and B: (A XOR B) XOR B = A. The key cancels itself out.\n\n"
        "Real-world cryptographic uses of XOR:\n"
        "- Stream ciphers: In a stream cipher, a keystream generator produces a long sequence of pseudo-random bits. The plaintext is XORed with this keystream to produce ciphertext. The One-Time Pad (OTP) is the theoretically unbreakable stream cipher where the key is truly random, as long as the message, and never reused. Real-world stream ciphers like RC4 (now broken and deprecated) and ChaCha20 (used in modern Transport Layer Security, or TLS) use this principle.\n"
        "- Block ciphers: The Data Encryption Standard (DES) and its successor, the Advanced Encryption Standard (AES), use XOR operations in multiple rounds to mix the plaintext with the key. Each round transforms the data using substitution, permutation, and XOR with round keys derived from the main key.\n"
        "- Hash functions: Some hash algorithms use XOR to combine blocks of data.\n"
        "- Cyclic Redundancy Check (CRC): XOR is used to calculate checksums that detect accidental changes to data during transmission.\n\n"
        "Security warning about XOR-based 'encryption': Because XOR is simple, inexperienced developers sometimes create their own 'encryption' by XORing data with a short, repeating password. This is extremely insecure. If an attacker can guess even a small portion of the plaintext (for example, knowing that an email starts with 'Dear '), they can derive the key and decrypt the entire message. The security of XOR-based encryption depends entirely on the quality and secrecy of the key.\n\n"
        "Abbreviations defined:\n"
        "- XOR = Exclusive OR\n"
        "- DES = Data Encryption Standard\n"
        "- 3DES = Triple Data Encryption Standard\n"
        "- AES = Advanced Encryption Standard\n"
        "- OTP = One-Time Pad\n"
        "- RC4 = Rivest Cipher 4\n"
        "- TLS = Transport Layer Security\n"
        "- CRC = Cyclic Redundancy Check"
    )
    
    add_bullet_slide(prs,
        "XOR in Action: Encryption Example",
        [
            "",
            "Plaintext letter 'A' in binary = 01000001",
            "Key byte                     = 00110101",
            "-------------------------------------------",
            "Ciphertext (XOR result)      = 01110100",
            "",
            "Decryption:",
            "Ciphertext                   = 01110100",
            "Key byte                     = 00110101",
            "-------------------------------------------",
            "Plaintext restored           = 01000001 = 'A'",
            "",
            "This demonstrates why XOR is called 'self-inverting.' Applying the same key twice returns the original data."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "This slide walks through a concrete, byte-level XOR encryption and decryption example. Instructors should present this line by line to ensure every student understands the bitwise operation.\n\n"
        "Step 1: Convert the plaintext to binary. The letter 'A' in the American Standard Code for Information Interchange (ASCII) is decimal 65, which is binary 01000001.\n\n"
        "Step 2: Choose a key byte. For this example, the key is 00110101 (which is decimal 53 in decimal, or hex 35).\n\n"
        "Step 3: XOR each corresponding bit:\n"
        "- Bit 1: 0 XOR 0 = 0\n"
        "- Bit 2: 1 XOR 0 = 1\n"
        "- Bit 3: 0 XOR 1 = 1\n"
        "- Bit 4: 0 XOR 1 = 1\n"
        "- Bit 5: 0 XOR 0 = 0\n"
        "- Bit 6: 0 XOR 1 = 1\n"
        "- Bit 7: 0 XOR 0 = 0\n"
        "- Bit 8: 1 XOR 1 = 0\n\n"
        "Wait - let me recalculate correctly, bit by bit:\n"
        "  Plaintext:  0 1 0 0 0 0 0 1\n"
        "  Key:        0 0 1 1 0 1 0 1\n"
        "  ----------------------------\n"
        "  XOR result: 0 1 1 1 0 1 0 0\n\n"
        "So the ciphertext is 01110100, which is decimal 116. In ASCII, decimal 116 is the lowercase letter 't'. So encrypting 'A' with this key produces the ciphertext byte 't'.\n\n"
        "Step 4: Decryption. To recover the original plaintext, XOR the ciphertext with the same key:\n"
        "  Ciphertext: 0 1 1 1 0 1 0 0\n"
        "  Key:        0 0 1 1 0 1 0 1\n"
        "  ----------------------------\n"
        "  Result:     0 1 0 0 0 0 0 1\n\n"
        "The result is 01000001, which is the ASCII code for 'A'. The original data is perfectly restored.\n\n"
        "Why this matters: This simple example is the exact same principle used by billion-dollar encryption systems. The only difference is that modern algorithms like Advanced Encryption Standard (AES) use much longer keys (128, 192, or 256 bits), perform many rounds of transformations, and include additional operations like substitution and permutation to ensure security. But at the core, XOR is doing the heavy lifting of combining the plaintext with the key.\n\n"
        "Abbreviations defined:\n"
        "- ASCII = American Standard Code for Information Interchange\n"
        "- AES = Advanced Encryption Standard"
    )
    
    add_bullet_slide(prs,
        "Topic 4 Summary: Key Takeaways",
        [
            "- Binary (base-2) is the language of computers. One byte is 8 bits, representing values 0 to 255.",
            "- Binary addition follows simple rules, but overflow can cause serious security vulnerabilities.",
            "- Hexadecimal (base-16) is a compact way to represent binary. It is used for MAC addresses, IPv6 addresses, memory dumps, and cryptographic hashes.",
            "- Boolean logic (AND, OR, NOT) is the foundation of digital circuits, programming conditions, and access control.",
            "- Exclusive OR (XOR) outputs 1 only when inputs differ. Its self-inverse property makes it the core operation in most encryption algorithms.",
            "- Mathematics is not just theory - it is the practical foundation of cryptography, networking, and secure computing."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "This topic provides the mathematical literacy that underpins all technical cybersecurity work. Students should not feel intimidated by the math - every concept introduced here is used directly in real security tools and protocols that they will encounter throughout the course and in their careers.\n\n"
        "Key connections to future topics and careers:\n"
        "- Binary and hex are essential for reading packet captures, memory dumps, and forensic artifacts.\n"
        "- Boolean logic is the basis of firewall rules, intrusion detection signatures, and SQL queries.\n"
        "- XOR is the mathematical heart of encryption. When students later learn about Virtual Private Networks (VPNs), Transport Layer Security (TLS), and disk encryption, they will be building on the XOR concepts learned here.\n\n"
        "For the practical assessment, students should be able to convert small decimal numbers to binary and hex, construct simple truth tables, explain why XOR is used in cryptography, and perform a manual XOR encryption/decryption on a single byte. A practical exercise might involve using an online XOR calculator or a simple Python script to encrypt a short message."
    )
    
    out = os.path.join(BASE_DIR, "Topic_4_Mathematics_for_Cybersecurity.pptx")
    prs.save(out)
    print(f"Saved: {out}")


# ---------------------------------------------------------------------------
# TOPIC 5: NETWORKING ESSENTIALS
# ---------------------------------------------------------------------------
def build_topic_5():
    prs = new_prs()
    
    add_title_slide(prs,
        "Topic 5: Networking Essentials",
        "How Data Travels and How Attackers Intercept It",
        "Welcome to Topic 5: Networking Essentials. Computer networks are the highways of the digital world. Every email, website visit, bank transaction, and video call depends on networks. For cybersecurity professionals, understanding how data moves across a network is not optional - it is essential. You cannot defend what you do not understand. This topic covers network types, the OSI and TCP/IP models, IP addressing, protocols, ports, common network devices, and the security threats that target each layer."
    )
    
    add_bullet_slide(prs,
        "Types of Networks",
        [
            "- Local Area Network (LAN): A network that covers a small geographic area, such as a home, office, school, or building. LANs are fast and privately controlled.",
            "- Wide Area Network (WAN): A network that spans a large geographic area, such as a city, country, or the entire world. The internet is the largest WAN.",
            "- Personal Area Network (PAN): A very small network centered around a single person, typically using Bluetooth or USB connections. Examples include wireless earphones and fitness trackers.",
            "- Wireless Local Area Network (WLAN): A LAN that uses Wi-Fi instead of physical cables.",
            "- Metropolitan Area Network (MAN): A network that covers a city or large campus, such as a cable television network or municipal fiber network."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Understanding network types helps cybersecurity professionals assess the scope of their responsibility and the nature of threats they face.\n\n"
        "Local Area Network (LAN): A LAN connects devices in a limited area, typically using Ethernet cables or Wi-Fi. In an office environment, all employee computers, printers, and servers are on the same LAN. Because LANs are controlled by the organization, administrators can enforce security policies like firewalls, network segmentation, and intrusion detection systems. However, if an attacker gains access to the LAN - for example, by plugging into an open network jack or compromising an employee's laptop - they can move laterally, scanning and attacking other devices on the same network. This is why internal network security is just as important as perimeter security.\n\n"
        "Wide Area Network (WAN): WANs connect multiple LANs across large distances. A Zambian bank with branches in Lusaka, Kitwe, and Livingstone uses a WAN to connect those branch LANs. WANs are typically leased from telecommunications providers like MTN Zambia or Airtel Zambia. Because data travels over third-party infrastructure, WAN traffic is often encrypted using Virtual Private Network (VPN) tunnels to prevent interception. The internet itself is a massive, global WAN composed of millions of interconnected networks.\n\n"
        "Personal Area Network (PAN): PANs are the smallest type of network. They connect personal devices like smartphones, wireless keyboards, Bluetooth speakers, and smartwatches. Because PANs often use wireless protocols like Bluetooth, they are vulnerable to attacks like Bluejacking (sending unsolicited messages) and Bluesnarfing (stealing data over Bluetooth). Users should disable Bluetooth when not in use and avoid pairing with unknown devices.\n\n"
        "Wireless Local Area Network (WLAN): WLANs use radio waves instead of cables. The most common WLAN technology is Wi-Fi (based on the IEEE 802.11 standards). WLANs are convenient but introduce unique security challenges. Wireless signals pass through walls, meaning an attacker with a laptop and antenna can sit in a parking lot and attempt to access the network. Securing WLANs requires strong encryption (WPA3 is the current standard, replacing WPA2), hidden Service Set Identifiers (SSIDs) if desired (though this is not a strong security measure), and MAC address filtering.\n\n"
        "Metropolitan Area Network (MAN): MANs cover a city or large campus. Examples include municipal Wi-Fi projects and cable TV distribution networks. From a cybersecurity perspective, MANs are less commonly managed by individual organizations but are relevant for understanding internet service provider (ISP) infrastructure.\n\n"
        "Abbreviations defined:\n"
        "- LAN = Local Area Network\n"
        "- WAN = Wide Area Network\n"
        "- PAN = Personal Area Network\n"
        "- WLAN = Wireless Local Area Network\n"
        "- Wi-Fi = Wireless Fidelity (based on IEEE 802.11 standards)\n"
        "- VPN = Virtual Private Network\n"
        "- WPA3 = Wi-Fi Protected Access 3\n"
        "- WPA2 = Wi-Fi Protected Access 2\n"
        "- SSID = Service Set Identifier\n"
        "- MAC = Media Access Control\n"
        "- ISP = Internet Service Provider\n"
        "- IEEE = Institute of Electrical and Electronics Engineers"
    )
    
    add_bullet_slide(prs,
        "The OSI Model: Seven Layers of Networking",
        [
            "- The Open Systems Interconnection (OSI) model is a conceptual framework that divides network communication into seven distinct layers.",
            "- Each layer has a specific function and communicates with the layers directly above and below it.",
            "- The seven layers are:",
            "   7. Application  (user-facing software like web browsers and email clients)",
            "   6. Presentation  (data formatting, encryption, and compression)",
            "   5. Session       (establishing, managing, and terminating connections)",
            "   4. Transport     (reliable data delivery, error checking)",
            "   3. Network       (logical addressing and routing)",
            "   2. Data Link     (physical addressing and local network access)",
            "   1. Physical      (raw electrical, optical, or radio signals)"
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "The OSI model was developed by the International Organization for Standardization (ISO) in 1984 as a universal reference model for networking. While modern networks do not perfectly map to OSI, it remains the standard teaching framework and the language that network engineers and cybersecurity professionals use to describe where problems and attacks occur.\n\n"
        "The model is often remembered using mnemonics. One popular version from top to bottom is: All People Seem To Need Data Processing (Application, Presentation, Session, Transport, Network, Data Link, Physical). Another from bottom to top is: Please Do Not Throw Sausage Pizza Away (Physical, Data Link, Network, Transport, Session, Presentation, Application). Instructors can encourage students to create their own memorable phrases.\n\n"
        "Detailed layer descriptions:\n"
        "- Layer 7 - Application: This is the layer closest to the end user. It includes web browsers, email clients, file transfer programs, and Application Programming Interfaces (APIs). Protocols at this layer include HyperText Transfer Protocol (HTTP), HyperText Transfer Protocol Secure (HTTPS), File Transfer Protocol (FTP), Simple Mail Transfer Protocol (SMTP), and Post Office Protocol version 3 (POP3). Attacks at this layer include SQL injection, cross-site scripting (XSS), and malware delivered through email attachments.\n"
        "- Layer 6 - Presentation: This layer is responsible for translating data between the application layer and the network. It handles data formatting (for example, converting Extended Binary Coded Decimal Interchange Code (EBCDIC) to ASCII), encryption and decryption, and compression. Transport Layer Security (TLS) and Secure Sockets Layer (SSL) operate here. If encryption fails or certificates are invalid, the presentation layer is where the error manifests.\n"
        "- Layer 5 - Session: This layer manages sessions - sustained interactions between applications. It establishes, maintains, and terminates connections. For example, when you log into a website, the session layer helps maintain your login state. Session hijacking is an attack where an attacker steals a session token (like a cookie) and impersonates the legitimate user.\n"
        "- Layer 4 - Transport: This layer ensures reliable end-to-end data delivery. The two most important protocols here are Transmission Control Protocol (TCP) and User Datagram Protocol (UDP). TCP provides reliable, ordered delivery with error checking. UDP provides fast, connectionless delivery without guarantees. Attacks at this layer include SYN flood (a type of Denial of Service, or DoS) and TCP session hijacking.\n"
        "- Layer 3 - Network: This layer handles logical addressing and routing. The Internet Protocol (IP) operates here. It is responsible for getting packets from source to destination across multiple networks. Devices at this layer are routers. Attacks include IP spoofing (sending packets with a fake source address) and routing table poisoning.\n"
        "- Layer 2 - Data Link: This layer handles physical addressing using Media Access Control (MAC) addresses, and it manages access to the physical network medium. Devices at this layer are switches and network interface cards (NICs). Attacks include ARP (Address Resolution Protocol) spoofing and MAC flooding.\n"
        "- Layer 1 - Physical: This is the hardware layer. It includes cables, fiber optics, radio waves, hubs, and repeaters. Attacks include wiretapping, cutting cables, and jamming wireless signals.\n\n"
        "Abbreviations defined:\n"
        "- OSI = Open Systems Interconnection\n"
        "- ISO = International Organization for Standardization\n"
        "- API = Application Programming Interface\n"
        "- HTTP = HyperText Transfer Protocol\n"
        "- HTTPS = HyperText Transfer Protocol Secure\n"
        "- FTP = File Transfer Protocol\n"
        "- SMTP = Simple Mail Transfer Protocol\n"
        "- POP3 = Post Office Protocol version 3\n"
        "- SQL = Structured Query Language\n"
        "- XSS = Cross-Site Scripting\n"
        "- EBCDIC = Extended Binary Coded Decimal Interchange Code\n"
        "- ASCII = American Standard Code for Information Interchange\n"
        "- TLS = Transport Layer Security\n"
        "- SSL = Secure Sockets Layer\n"
        "- TCP = Transmission Control Protocol\n"
        "- UDP = User Datagram Protocol\n"
        "- DoS = Denial of Service\n"
        "- IP = Internet Protocol\n"
        "- MAC = Media Access Control\n"
        "- NIC = Network Interface Card\n"
        "- ARP = Address Resolution Protocol",
        img=os.path.join(ASSETS_DIR, "osi_model.png")
    )
    
    add_bullet_slide(prs,
        "The TCP/IP Model",
        [
            "- The TCP/IP model is the practical framework used by the internet and most modern networks.",
            "- It combines some OSI layers into four layers:",
            "   4. Application    (combines OSI Application, Presentation, and Session)",
            "   3. Transport      (same as OSI Transport - TCP and UDP)",
            "   2. Internet       (same as OSI Network - IP addressing and routing)",
            "   1. Network Access (combines OSI Data Link and Physical)",
            "- TCP/IP was developed by the United States Department of Defense and is the foundation of the internet.",
            "- While OSI is a teaching model, TCP/IP is the implementation model."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "The TCP/IP model predates the OSI model. It was developed in the 1970s by the United States Department of Defense (DoD) Advanced Research Projects Agency (DARPA) as part of the ARPANET project - the precursor to the modern internet. Because TCP/IP was built for real-world use rather than theoretical perfection, it is simpler and more practical than OSI.\n\n"
        "Layer 4 - Application: In the TCP/IP model, this single layer covers everything that the OSI model splits into Application, Presentation, and Session. Protocols like HTTP, HTTPS, FTP, SMTP, and Domain Name System (DNS) all live here. When a web browser requests a page, it uses HTTP at the application layer. When an email client sends a message, it uses SMTP. When you type a domain name like google.com into your browser, DNS resolves it to an IP address at this layer.\n\n"
        "Layer 3 - Transport: This layer is nearly identical to OSI Layer 4. TCP is the workhorse of the internet. It breaks data into segments, numbers them, sends them across the network, and reassembles them in the correct order at the destination. It also handles error checking and retransmission of lost segments. UDP is a simpler, faster alternative used by applications that prioritize speed over reliability, such as online gaming, live video streaming, and DNS queries.\n\n"
        "Layer 2 - Internet: This layer corresponds to OSI Layer 3 (Network). The Internet Protocol (IP) is responsible for logical addressing and routing. Every device on a TCP/IP network has an IP address. IPv4 addresses are 32 bits long and are written in dotted-decimal notation (for example, 192.168.1.50). IPv6 addresses are 128 bits long and are written in hexadecimal (for example, 2001:0DB8:85A3::8A2E:0370:7334). Routing protocols like Open Shortest Path First (OSPF) and Border Gateway Protocol (BGP) operate at this layer, determining the best path for packets to travel across the internet.\n\n"
        "Layer 1 - Network Access: This layer combines OSI Layers 1 and 2. It handles everything related to putting data onto the physical medium - whether that is an Ethernet cable, a fiber optic line, or a Wi-Fi radio signal. Ethernet frames, MAC addresses, and physical signaling all belong here.\n\n"
        "OSI vs TCP/IP in practice: Network engineers and cybersecurity professionals use both models. OSI provides a more detailed conceptual framework, which is excellent for learning and troubleshooting. TCP/IP reflects how the internet actually works. When analyzing a network problem, you might ask: 'Is this a Layer 2 issue (Data Link), a Layer 3 issue (Network), or a Layer 7 issue (Application)?' This layered language is universal in the industry.\n\n"
        "Abbreviations defined:\n"
        "- DoD = Department of Defense\n"
        "- DARPA = Defense Advanced Research Projects Agency\n"
        "- ARPANET = Advanced Research Projects Agency Network\n"
        "- HTTP = HyperText Transfer Protocol\n"
        "- HTTPS = HyperText Transfer Protocol Secure\n"
        "- FTP = File Transfer Protocol\n"
        "- SMTP = Simple Mail Transfer Protocol\n"
        "- DNS = Domain Name System\n"
        "- TCP = Transmission Control Protocol\n"
        "- UDP = User Datagram Protocol\n"
        "- IP = Internet Protocol\n"
        "- IPv4 = Internet Protocol version 4\n"
        "- IPv6 = Internet Protocol version 6\n"
        "- OSPF = Open Shortest Path First\n"
        "- BGP = Border Gateway Protocol\n"
        "- MAC = Media Access Control",
        img=os.path.join(ASSETS_DIR, "osi_tcpip_comparison.png")
    )
    
    add_bullet_slide(prs,
        "How TCP Works: The Three-Way Handshake",
        [
            "- TCP is a connection-oriented protocol. Before data is sent, the sender and receiver must establish a connection using a three-way handshake.",
            "- Step 1 - SYN (Synchronize): The client sends a SYN packet to the server, requesting to open a connection.",
            "- Step 2 - SYN-ACK (Synchronize-Acknowledge): The server responds with a SYN-ACK packet, acknowledging the request and sending its own synchronization message.",
            "- Step 3 - ACK (Acknowledge): The client sends an ACK packet back to the server. The connection is now established, and data transfer can begin.",
            "- This handshake ensures both sides are ready to communicate and agree on initial sequence numbers."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "The TCP three-way handshake is one of the most important concepts in networking. It is the foundation of reliable internet communication, and it is also the target of one of the most common network attacks: the SYN flood.\n\n"
        "Why a handshake is needed: Unlike UDP, which simply sends data without checking whether the recipient is ready, TCP ensures that both the client and server are alive, reachable, and ready to exchange data. The handshake also synchronizes sequence numbers - unique numbers assigned to each byte of data so the receiver can reassemble packets in the correct order and detect missing segments.\n\n"
        "Detailed handshake walkthrough:\n"
        "1. SYN: The client chooses a random initial sequence number (let's call it 1000) and sends a TCP packet with the SYN flag set to the server's IP address and port. This says, 'I want to start a conversation, and my starting number is 1000.'\n"
        "2. SYN-ACK: The server receives the SYN, allocates resources for the connection (creating a data structure in memory called a Transmission Control Block), chooses its own random initial sequence number (let's say 5000), and sends back a packet with both the SYN and ACK flags set. The acknowledgment number is 1001, which means 'I received your SYN, and I expect your next packet to start at 1001.'\n"
        "3. ACK: The client receives the SYN-ACK, sends an ACK packet with acknowledgment number 5001 (acknowledging the server's SYN), and the connection is officially established. Now both sides can send data.\n\n"
        "TCP flags: TCP packets contain several flags that control the connection state. The key flags are:\n"
        "- SYN (Synchronize): Initiates a connection.\n"
        "- ACK (Acknowledge): Confirms receipt of data.\n"
        "- FIN (Finish): Gracefully terminates a connection.\n"
        "- RST (Reset): Abruptly terminates a connection, often used when a port is closed or something goes wrong.\n"
        "- PSH (Push): Asks the receiver to pass data to the application immediately.\n"
        "- URG (Urgent): Indicates that the data should be processed urgently.\n\n"
        "The SYN flood attack: An attacker sends a flood of SYN packets to a server, often using spoofed (fake) source IP addresses. The server responds with SYN-ACK packets and allocates resources for each half-open connection. However, because the source addresses are fake, the final ACK packets never arrive. The server's connection table fills up with half-open connections, leaving no room for legitimate users. This is a classic Denial of Service (DoS) attack. Defenses include SYN cookies (a technique where the server avoids allocating resources until the handshake is complete) and rate limiting.\n\n"
        "Abbreviations defined:\n"
        "- TCP = Transmission Control Protocol\n"
        "- UDP = User Datagram Protocol\n"
        "- SYN = Synchronize\n"
        "- ACK = Acknowledge\n"
        "- FIN = Finish\n"
        "- RST = Reset\n"
        "- PSH = Push\n"
        "- URG = Urgent\n"
        "- IP = Internet Protocol\n"
        "- DoS = Denial of Service",
        img=os.path.join(ASSETS_DIR, "tcp_handshake.png")
    )
    
    add_bullet_slide(prs,
        "IP Addressing and Subnetting",
        [
            "- An IP address is a unique identifier assigned to every device on a network.",
            "- IPv4 addresses are 32 bits long, divided into four octets (bytes) separated by dots. Example: 192.168.1.50.",
            "- IPv6 addresses are 128 bits long, written in hexadecimal. Example: 2001:0DB8:85A3::8A2E:0370:7334.",
            "- Private IP ranges (used inside organizations) include:",
            "   + 10.0.0.0 to 10.255.255.255",
            "   + 172.16.0.0 to 172.31.255.255",
            "   + 192.168.0.0 to 192.168.255.255",
            "- A subnet mask defines which part of an IP address is the network portion and which part is the host portion.",
            "- The Classless Inter-Domain Routing (CIDR) notation, such as /24, indicates how many bits are used for the network portion."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "IP addressing is the foundation of routing. Without unique addresses, devices would not know where to send data.\n\n"
        "IPv4 structure: An IPv4 address is a 32-bit number. Because humans struggle to read long binary strings, it is written in dotted-decimal notation - four numbers between 0 and 255 separated by dots. Each number represents one byte (8 bits). For example, 192.168.1.50 in binary is 11000000.10101000.00000001.00110010.\n\n"
        "IPv6: Because the world was running out of IPv4 addresses (there are only about 4.3 billion possible IPv4 addresses), IPv6 was developed with 128-bit addresses, providing approximately 340 undecillion (3.4 x 10^38) unique addresses. IPv6 addresses are written in hexadecimal, grouped into eight blocks of four hex digits separated by colons. Leading zeros in each block can be omitted, and consecutive blocks of zeros can be replaced with a double colon (::) once per address. For example, 2001:0DB8:0000:0000:0000:0000:1428:57AB can be shortened to 2001:DB8::1428:57AB.\n\n"
        "Private IP addresses: Not every device on the internet needs a globally unique IP address. Private IP ranges are reserved for use inside organizations, and Network Address Translation (NAT) allows many private devices to share a single public IP address. The three private IPv4 ranges are:\n"
        "- 10.0.0.0 /8 (10.0.0.0 to 10.255.255.255)\n"
        "- 172.16.0.0 /12 (172.16.0.0 to 172.31.255.255)\n"
        "- 192.168.0.0 /16 (192.168.0.0 to 192.168.255.255)\n"
        "If you check the IP address of your home router or office computer, it will almost certainly be in one of these ranges.\n\n"
        "Subnetting: A subnet mask separates the network portion of an IP address from the host portion. For example, a subnet mask of 255.255.255.0 means the first three octets identify the network, and the last octet identifies the specific host. In CIDR notation, this is written as /24 because the first 24 bits are the network portion. A /24 network can have up to 256 IP addresses (0 to 255), but two are reserved: the network address (x.x.x.0) and the broadcast address (x.x.x.255). So a /24 provides 254 usable host addresses.\n\n"
        "Security implications of subnetting: Network segmentation is a critical security practice. By dividing a large network into smaller subnets, organizations can isolate sensitive systems. For example, the finance department might be on subnet 192.168.10.0/24, while guest Wi-Fi users are on 192.168.50.0/24. Firewalls can enforce rules that prevent guest users from accessing the finance subnet. This limits the 'blast radius' if one part of the network is compromised.\n\n"
        "Abbreviations defined:\n"
        "- IP = Internet Protocol\n"
        "- IPv4 = Internet Protocol version 4\n"
        "- IPv6 = Internet Protocol version 6\n"
        "- CIDR = Classless Inter-Domain Routing\n"
        "- NAT = Network Address Translation"
    )
    
    add_bullet_slide(prs,
        "Network Devices and Their Security Roles",
        [
            "- Network Interface Card (NIC): A hardware component that connects a computer to a network. Every NIC has a unique Media Access Control (MAC) address burned into it at the factory.",
            "- Switch: A Layer 2 device that forwards data frames to specific devices within a LAN based on MAC addresses. A managed switch supports security features like Virtual LANs (VLANs) and port security.",
            "- Router: A Layer 3 device that forwards packets between different networks based on IP addresses. Routers use routing tables and protocols to determine the best path.",
            "- Firewall: A security device that monitors and controls incoming and outgoing network traffic based on predefined rules. It can operate at multiple OSI layers.",
            "- Access Point (AP): A device that allows wireless devices to connect to a wired network."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Understanding network devices is essential for both designing secure networks and investigating incidents. Each device has specific capabilities and vulnerabilities.\n\n"
        "Network Interface Card (NIC): Every networked computer, smartphone, printer, and IoT device has at least one NIC. Each NIC has a globally unique MAC address assigned by the manufacturer. The MAC address is 48 bits long and is typically written as six pairs of hex digits separated by colons or hyphens, like 00:1A:2B:3C:4D:5E. The first three pairs (the Organizationally Unique Identifier, or OUI) identify the manufacturer. For example, addresses starting with 00:50:56 belong to VMware. MAC addresses are used for communication within a local network (Layer 2). MAC spoofing is an attack where a device changes its MAC address to impersonate another device, often to bypass MAC filtering or to hide the attacker's identity.\n\n"
        "Switch: A switch is smarter than an older network device called a hub. A hub broadcasts all incoming data to every connected device, which creates a major security risk - any computer on the hub can see all traffic. A switch, by contrast, learns which MAC addresses are connected to which ports and forwards frames only to the intended recipient. This is much more secure, but it is not perfect. Attackers can use ARP spoofing to trick the switch into sending traffic to the wrong port. Managed switches offer advanced security features like VLANs (which create separate logical networks on the same physical switch) and port security (which limits which MAC addresses can connect to each port).\n\n"
        "Router: Routers are the traffic directors of the internet. They examine the destination IP address of each packet and forward it toward its destination using routing tables. Routers can connect different types of networks (for example, a LAN to a WAN, or Ethernet to Wi-Fi). Home routers combine multiple functions: routing, Network Address Translation (NAT), Dynamic Host Configuration Protocol (DHCP), and sometimes firewalling. A misconfigured router - for example, one with default passwords still in place - is a common entry point for attackers. The 2016 Mirai botnet infected hundreds of thousands of routers and IoT devices by trying a short list of default passwords.\n\n"
        "Firewall: Firewalls are the primary perimeter defense for networks. They inspect traffic and allow or block it based on rules. Traditional firewalls operate at Layers 3 and 4, filtering by IP address and port number. Next-Generation Firewalls (NGFWs) can inspect traffic at the application layer (Layer 7), detecting malware hidden inside allowed protocols. Firewalls can be hardware appliances or software programs (like the built-in Windows Firewall or Linux iptables). A poorly configured firewall - one with overly permissive rules - provides a false sense of security.\n\n"
        "Access Point (AP): An AP bridges wireless clients to a wired LAN. Modern enterprise APs support strong encryption (WPA2/WPA3), multiple SSIDs (network names), and client isolation (preventing wireless clients from communicating directly with each other). Rogue access points - unauthorized APs placed on a network by attackers or careless employees - are a serious security threat because they can provide an unmonitored backdoor into the LAN.\n\n"
        "Abbreviations defined:\n"
        "- NIC = Network Interface Card\n"
        "- MAC = Media Access Control\n"
        "- IoT = Internet of Things\n"
        "- OUI = Organizationally Unique Identifier\n"
        "- ARP = Address Resolution Protocol\n"
        "- VLAN = Virtual Local Area Network\n"
        "- NAT = Network Address Translation\n"
        "- DHCP = Dynamic Host Configuration Protocol\n"
        "- NGFW = Next-Generation Firewall\n"
        "- WPA2 = Wi-Fi Protected Access 2\n"
        "- WPA3 = Wi-Fi Protected Access 3\n"
        "- SSID = Service Set Identifier\n"
        "- LAN = Local Area Network\n"
        "- IoT = Internet of Things"
    )
    
    add_bullet_slide(prs,
        "Network Protocols and Port Numbers",
        [
            "- A protocol is a set of rules that define how devices communicate.",
            "- A port number is a 16-bit identifier that tells the operating system which application should receive incoming data.",
            "- Well-known ports (0 to 1023) are reserved for standard services and require administrator privileges on most systems.",
            "- Common well-known ports include:",
            "   + Port 20/21: File Transfer Protocol (FTP)",
            "   + Port 22: Secure Shell (SSH)",
            "   + Port 25: Simple Mail Transfer Protocol (SMTP)",
            "   + Port 53: Domain Name System (DNS)",
            "   + Port 80: HyperText Transfer Protocol (HTTP)",
            "   + Port 110: Post Office Protocol version 3 (POP3)",
            "   + Port 143: Internet Message Access Protocol (IMAP)",
            "   + Port 443: HyperText Transfer Protocol Secure (HTTPS)",
            "   + Port 3389: Remote Desktop Protocol (RDP)"
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Protocols and ports are the vocabulary of network communication. Every service running on a server listens on a specific port, waiting for incoming connections. Understanding which ports are used by which services is essential for configuring firewalls, diagnosing connectivity issues, and identifying malicious activity.\n\n"
        "How ports work: When a client wants to connect to a server, it sends a packet to the server's IP address and port number. The operating system on the server checks the port number and delivers the data to the correct application. For example, when you visit a website, your browser sends a request to the web server's IP address on port 80 (for HTTP) or port 443 (for HTTPS). The web server software (like Apache, Nginx, or Microsoft Internet Information Services) is listening on that port and responds with the requested web page.\n\n"
        "Port ranges:\n"
        "- Well-known ports: 0 to 1023. Reserved by the Internet Assigned Numbers Authority (IANA) for standard services. On Linux and Unix systems, only the root user can open a well-known port.\n"
        "- Registered ports: 1024 to 49151. Used by vendors for proprietary applications. Examples include port 3306 (MySQL database), port 5432 (PostgreSQL), and port 8080 (alternate HTTP).\n"
        "- Dynamic/Private ports: 49152 to 65535. Used temporarily by client applications for outbound connections. When your browser connects to a web server, it uses a high-numbered dynamic port on your computer as the source port.\n\n"
        "Detailed protocol and port descriptions:\n"
        "- FTP (Ports 20 and 21): File Transfer Protocol is used for transferring files between computers. Port 21 is for control commands, and port 20 is for data transfer. FTP sends passwords in plain text, making it insecure. Secure alternatives include SFTP (SSH File Transfer Protocol) and FTPS (FTP over SSL/TLS).\n"
        "- SSH (Port 22): Secure Shell provides encrypted remote access to servers and network devices. It is the standard remote management protocol for Linux and Unix systems. Attackers constantly scan the internet for open port 22 and attempt brute-force password attacks. Best practices include disabling password authentication in favor of SSH key pairs, changing the default port, and using fail2ban to block repeated failed attempts.\n"
        "- SMTP (Port 25): Simple Mail Transfer Protocol is used for sending email between mail servers. It is also used by email clients to send outgoing mail. Like FTP, traditional SMTP does not encrypt traffic. Modern implementations use STARTTLS to upgrade to encrypted connections.\n"
        "- DNS (Port 53): The Domain Name System translates human-readable domain names (like google.com) into IP addresses. It uses UDP for most queries and TCP for zone transfers. DNS is a frequent attack target. DNS spoofing (also called DNS cache poisoning) redirects users to malicious websites. DNS tunneling encodes data in DNS queries to bypass firewalls.\n"
        "- HTTP (Port 80): HyperText Transfer Protocol is the foundation of the World Wide Web. It is unencrypted, meaning anyone who intercepts the traffic can read the contents. For this reason, modern websites redirect HTTP traffic to HTTPS.\n"
        "- POP3 (Port 110): Post Office Protocol version 3 is used by email clients to download messages from a mail server. It typically downloads emails to the local device and deletes them from the server. POP3S (POP3 Secure) uses port 995 and encrypts the connection with SSL/TLS.\n"
        "- IMAP (Port 143): Internet Message Access Protocol is another email retrieval protocol. Unlike POP3, IMAP leaves messages on the server and synchronizes them across multiple devices. IMAPS (IMAP Secure) uses port 993.\n"
        "- HTTPS (Port 443): HyperText Transfer Protocol Secure is the encrypted version of HTTP. It uses Transport Layer Security (TLS) to encrypt data between the browser and the web server. HTTPS is essential for protecting passwords, credit card numbers, and personal information. Modern browsers warn users if a website does not use HTTPS.\n"
        "- RDP (Port 3389): Remote Desktop Protocol is Microsoft's protocol for remote graphical access to Windows computers. Like SSH, it is a common target for brute-force attacks. The 2017 WannaCry ransomware spread partly by exploiting a vulnerability in the Server Message Block (SMB) protocol, but RDP is also frequently attacked directly. Best practices include placing RDP behind a VPN, using Network Level Authentication (NLA), and restricting access by IP address.\n\n"
        "Abbreviations defined:\n"
        "- IANA = Internet Assigned Numbers Authority\n"
        "- HTTP = HyperText Transfer Protocol\n"
        "- HTTPS = HyperText Transfer Protocol Secure\n"
        "- FTP = File Transfer Protocol\n"
        "- SSH = Secure Shell\n"
        "- SMTP = Simple Mail Transfer Protocol\n"
        "- DNS = Domain Name System\n"
        "- POP3 = Post Office Protocol version 3\n"
        "- IMAP = Internet Message Access Protocol\n"
        "- TLS = Transport Layer Security\n"
        "- SSL = Secure Sockets Layer\n"
        "- RDP = Remote Desktop Protocol\n"
        "- SMB = Server Message Block\n"
        "- NLA = Network Level Authentication\n"
        "- VPN = Virtual Private Network\n"
        "- SQL = Structured Query Language\n"
        "- UDP = User Datagram Protocol\n"
        "- TCP = Transmission Control Protocol"
    )
    
    add_bullet_slide(prs,
        "Network Security Threats by OSI Layer",
        [
            "- Network attacks can target any layer of the OSI model. Defenders must understand each layer's vulnerabilities.",
            "",
            "Layer 1 (Physical): Wiretapping, cable cutting, electromagnetic interference, device theft.",
            "Layer 2 (Data Link): ARP spoofing, MAC flooding, unauthorized switch access.",
            "Layer 3 (Network): IP spoofing, ICMP (Internet Control Message Protocol) redirection, routing attacks.",
            "Layer 4 (Transport): SYN flood, TCP reset attacks, session hijacking.",
            "Layer 5 (Session): Session hijacking and fixation.",
            "Layer 6 (Presentation): Weak or misconfigured encryption (SSL stripping).",
            "Layer 7 (Application): SQL injection, cross-site scripting (XSS), malware, phishing."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Understanding attacks by OSI layer is a core skill for cybersecurity professionals. It helps defenders choose the right countermeasures and speak the common language of the industry.\n\n"
        "Layer 1 - Physical attacks: These are often overlooked because they seem low-tech, but they can be devastating. Wiretapping involves physically connecting to a network cable to intercept data. Although modern networks use switches rather than hubs, specialized tools can still extract useful information from cables. Cable cutting is a Denial of Service (DoS) attack that physically disrupts connectivity. Electromagnetic interference can disrupt wireless signals. Finally, device theft - stealing a laptop, server, or backup tape - gives an attacker direct access to the data stored on it. This is why full disk encryption and physical security controls (locks, guards, surveillance cameras) are essential.\n\n"
        "Layer 2 - Data Link attacks:\n"
        "- ARP Spoofing (ARP Poisoning): The attacker sends fake Address Resolution Protocol (ARP) messages onto the LAN, associating their own MAC address with the IP address of another device (like the default gateway). Other devices on the network then send traffic intended for the gateway to the attacker instead. The attacker can intercept, inspect, or modify the traffic before forwarding it. This is a man-in-the-middle (MitM) attack.\n"
        "- MAC Flooding: The attacker sends a massive number of frames with fake source MAC addresses to a network switch. The switch's MAC address table overflows, causing it to fail open and broadcast all traffic to every port - effectively turning the switch into a hub. The attacker can then capture traffic from all devices on the switch. Defenses include enabling port security on managed switches, which limits the number of MAC addresses per port.\n\n"
        "Layer 3 - Network attacks:\n"
        "- IP Spoofing: The attacker crafts packets with a false source IP address. This is commonly used in Denial of Service attacks to hide the attacker's identity and to bypass IP-based access controls.\n"
        "- ICMP Redirection: An attacker sends fake Internet Control Message Protocol (ICMP) redirect messages to a host, telling it to send traffic through a router controlled by the attacker.\n"
        "- Routing Attacks: Attackers compromise routers or manipulate routing protocols to redirect traffic through their systems. Border Gateway Protocol (BGP) hijacking is a sophisticated attack where an attacker advertises false IP address ranges, causing internet traffic to be routed through their network.\n\n"
        "Layer 4 - Transport attacks:\n"
        "- SYN Flood: As described earlier, the attacker overwhelms a server with half-open TCP connections, exhausting resources.\n"
        "- TCP Reset Attack: An attacker injects forged TCP RST (reset) packets into a connection, causing it to terminate abruptly. This was used by some internet service providers to block peer-to-peer file sharing.\n"
        "- Session Hijacking: After a TCP connection is established, an attacker guesses or steals sequence numbers and injects malicious packets into the session, impersonating one of the parties.\n\n"
        "Layer 5 - Session attacks: Session hijacking and session fixation are attacks against the session management mechanisms of web applications. In session hijacking, the attacker steals a session cookie or token and uses it to impersonate the user. In session fixation, the attacker tricks the user into using a session ID that the attacker already knows.\n\n"
        "Layer 6 - Presentation attacks: Weak encryption or misconfigured Transport Layer Security (TLS) settings allow attackers to intercept and decrypt traffic. SSL stripping is an attack where an attacker downgrades an HTTPS connection to HTTP, removing encryption. Tools like SSLstrip automate this process on open Wi-Fi networks.\n\n"
        "Layer 7 - Application attacks: This is where the majority of cyberattacks occur because applications are complex and often poorly coded. Common attacks include:\n"
        "- SQL Injection: Inserting malicious SQL commands into input fields to steal or manipulate database contents.\n"
        "- Cross-Site Scripting (XSS): Injecting malicious scripts into web pages viewed by other users.\n"
        "- Malware Delivery: Trick users into downloading and executing malicious files.\n"
        "- Phishing: Deceptive emails or websites designed to steal credentials.\n\n"
        "Abbreviations defined:\n"
        "- OSI = Open Systems Interconnection\n"
        "- ARP = Address Resolution Protocol\n"
        "- MAC = Media Access Control\n"
        "- IP = Internet Protocol\n"
        "- ICMP = Internet Control Message Protocol\n"
        "- BGP = Border Gateway Protocol\n"
        "- TCP = Transmission Control Protocol\n"
        "- RST = Reset\n"
        "- TLS = Transport Layer Security\n"
        "- SSL = Secure Sockets Layer\n"
        "- HTTPS = HyperText Transfer Protocol Secure\n"
        "- HTTP = HyperText Transfer Protocol\n"
        "- SQL = Structured Query Language\n"
        "- XSS = Cross-Site Scripting\n"
        "- MitM = Man-in-the-Middle\n"
        "- DoS = Denial of Service\n"
        "- LAN = Local Area Network"
    )
    
    add_bullet_slide(prs,
        "Topic 5 Summary: Key Takeaways",
        [
            "- Networks connect devices over distances ranging from personal space (PAN) to the entire world (WAN / internet).",
            "- The OSI model divides networking into seven conceptual layers. The TCP/IP model compresses these into four practical layers.",
            "- TCP uses a three-way handshake (SYN, SYN-ACK, ACK) to establish reliable connections. This handshake is the target of SYN flood attacks.",
            "- IP addresses uniquely identify devices. IPv4 uses 32-bit dotted-decimal notation. IPv6 uses 128-bit hexadecimal notation.",
            "- Network devices include NICs, switches, routers, firewalls, and access points. Each has specific security capabilities and vulnerabilities.",
            "- Protocols use port numbers to deliver data to the correct application. Well-known ports include 22 (SSH), 80 (HTTP), and 443 (HTTPS).",
            "- Attackers target every OSI layer, from physical wiretapping to application-layer SQL injection. Defense requires layered security."
        ],
        "INSTRUCTOR NOTES - Detailed Explanation:\n"
        "Topic 5 provides the networking foundation that every cybersecurity professional needs. Without understanding how packets move, how addresses work, and where attacks occur, it is impossible to design effective defenses or investigate incidents.\n\n"
        "The concept of 'defense in depth' is directly tied to the OSI model. No single security control can protect against all threats. A secure network uses multiple layers of defense:\n"
        "- Physical security (locks, cameras) protects Layer 1.\n"
        "- Port security and VLANs protect Layer 2.\n"
        "- Firewalls and intrusion detection systems protect Layers 3 and 4.\n"
        "- Encryption (TLS, VPNs) protects Layers 5 and 6.\n"
        "- Web application firewalls, secure coding, and user awareness protect Layer 7.\n\n"
        "For the practical assessment, students should be able to identify common network devices, explain the TCP three-way handshake, convert simple binary and decimal numbers, list well-known ports and their services, and describe at least one attack for each of the lower OSI layers (1-4). A hands-on lab using Wireshark to capture and inspect network packets is an excellent complement to this topic."
    )
    
    out = os.path.join(BASE_DIR, "Topic_5_Networking_Essentials.pptx")
    prs.save(out)
    print(f"Saved: {out}")

# ---------------------------------------------------------------------------
# MAIN
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    build_topic_1()
    build_topic_2()
    build_topic_3()
    build_topic_4()
    build_topic_5()
    print("All Module 1 presentations generated successfully.")
    print("Generated files:")
    for i in range(1, 6):
        fname = os.path.join(BASE_DIR, f"Topic_{i}_*.pptx")
        for f in sorted(glob.glob(fname)):
            print("  -", os.path.basename(f))

