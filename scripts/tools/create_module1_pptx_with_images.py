#!/usr/bin/env python3
"""
Module 1: Foundations of Computing and Mathematics - PPTX Generator WITH IMAGES
Cybersecurity Certificate Program | EduTrack LMS
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

ASSETS = "/Users/michael/Documents/Programming/edutrack-lms/module1_pptx_assets"

# Color palette
DARK_BLUE = RGBColor(0x1A, 0x23, 0x7E)
ACCENT_ORANGE = RGBColor(0xF4, 0x4C, 0x00)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1F, 0x29, 0x33)

def add_title_slide(prs, title, subtitle, notes=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.333), Inches(1.0))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0xFF, 0xB7, 0x4D)
    p.alignment = PP_ALIGN.CENTER
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes
    return slide

def add_section_slide(prs, section_title, notes=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = ACCENT_ORANGE
    bg.line.fill.background()
    
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes
    return slide

def add_content_slide(prs, title, bullets, notes="", img_path=None, img_left=Inches(7.2), img_top=Inches(1.4), img_w=Inches(5.8), img_h=Inches(5.6)):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Text width depends on whether image is present
    text_w = Inches(6.6) if img_path else Inches(12.333)
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), text_w, Inches(5.8))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20 if img_path else 22)
        p.font.color.rgb = DARK_TEXT
        p.level = 0
        p.space_after = Pt(10)
    
    if img_path and os.path.exists(img_path):
        try:
            slide.shapes.add_picture(img_path, img_left, img_top, width=img_w)
        except Exception as e:
            print(f"Could not add image {img_path}: {e}")
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes
    return slide

# ======================= SLIDES =======================

add_title_slide(prs,
    "Module 1: Foundations of Computing and Mathematics",
    "Cybersecurity Certificate Program | EduTrack LMS | Zambia",
    "Welcome to Module 1: Foundations of Computing and Mathematics. This module is the bedrock of the entire Cybersecurity Certificate Program. Every advanced concept in cybersecurity depends on understanding how computers process data, how operating systems manage resources, how logical thinking enables automation, how number systems underpin cryptography, and how networks enable both communication and attack vectors. As instructors, emphasize that this is NOT a basic computing class. Every topic here has direct security implications."
)

add_content_slide(prs,
    "Why Foundations Matter in Cybersecurity",
    [
        "• You cannot secure what you do not understand.",
        "• Every cyberattack exploits how computers, OSs, networks, and code work.",
        "• Module 1 builds the mental models needed for threat detection, incident response, and defense.",
        "• From Zambia's banks to mobile money platforms — all run on these same fundamentals.",
        "• Direct career path: SOC Analyst, Network Security Technician, Junior Security Analyst."
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Open with a compelling example: In 2021, the Colonial Pipeline ransomware attack shut down fuel distribution across the US East Coast. The attackers gained access through a compromised VPN account — exploiting basic networking and OS weaknesses. In Zambia, MTN and Airtel mobile money systems, Zanaco banking infrastructure, and ZESCO systems all depend on the exact fundamentals we cover here. Tell students: 'Attackers know these basics better than many defenders. Our job is to out-learn them.'"
)

add_section_slide(prs, "1. Computer Fundamentals",
    "Section transition: 1. Computer Fundamentals. Use this moment to reset student attention and preview what security relevance is coming."
)

add_content_slide(prs,
    "What is a Computer? The IPOS Model",
    [
        "• Input:    Keyboard, mouse, biometric scanner, card reader",
        "• Processing: CPU executes instructions (the 'brain')",
        "• Output:   Monitor, printer, speakers",
        "• Storage:  HDD, SSD, USB — permanent data retention",
        "",
        "Cybersecurity angle: Every stage is an attack surface.",
        "• Keyloggers target Input",
        "• Side-channel attacks exploit Processing",
        "• Shoulder surfing attacks Output",
        "• Unencrypted USB drives risk Storage"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "The IPOS model maps directly to attack surfaces. Input stage: Keyloggers capture keystrokes. Processing stage: Spectre and Meltdown exploit CPU speculative execution, leaking sensitive data. Output stage: Shoulder surfing in Lusaka internet cafes targets the output phase. Storage stage: Unencrypted USB drives are a massive data loss vector. In 2022, the UK government fined organizations millions for lost USB drives containing personal data. Zambian banks and ZESCO handle customer data on local servers. Understanding where data lives at each IPOS stage is essential for proper data classification and protection.",
    img_path=os.path.join(ASSETS, "ipos_cycle.png"),
    img_left=Inches(7.0), img_top=Inches(1.4), img_w=Inches(6.0), img_h=Inches(5.6)
)

add_content_slide(prs,
    "Hardware Components & Security Relevance",
    [
        "• CPU (Central Processing Unit)",
        "   - Executes instructions; speed in GHz",
        "   - Security: Cryptographic ops, side-channel attacks",
        "",
        "• RAM (Random Access Memory)",
        "   - Temporary, volatile storage",
        "   - Security: Fileless malware lives ONLY in RAM",
        "",
        "• Storage (HDD / SSD / USB)",
        "   - Permanent retention of data",
        "   - Security: Encryption at rest protects stolen devices",
        "",
        "• NIC (Network Interface Card)",
        "   - Every NIC has a unique MAC address",
        "   - Security: MAC spoofing, device identification"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "CPU: Modern processors include security extensions like Intel SGX and AMD SEV. However, Spectre (2018) and Meltdown (2018) showed speculative execution could leak passwords and encryption keys. Over 90% of processors worldwide were affected. RAM: Fileless malware is one of the fastest-growing threat categories. By never writing to disk, it evades traditional antivirus. Tools like PowerShell Empire and Cobalt Strike operate primarily in memory. Storage: Full Disk Encryption using BitLocker (Windows) or LUKS (Linux) ensures stolen laptops remain unreadable. In 2023, the average cost of a data breach involving lost devices was $4.45 million globally. NIC: MAC addresses are theoretically unique, but tools like macchanger allow attackers to spoof them to bypass network access controls.",
    img_path=os.path.join(ASSETS, "hardware_components.png"),
    img_left=Inches(7.0), img_top=Inches(1.4), img_w=Inches(6.0), img_h=Inches(5.6)
)

add_content_slide(prs,
    "System Architecture: 32-bit vs 64-bit & Von Neumann",
    [
        "• 32-bit Systems:",
        "   - Can address up to 4GB of RAM",
        "   - Limited security features, being phased out",
        "",
        "• 64-bit Systems:",
        "   - Can address up to 16 exabytes theoretically",
        "   - Required for modern security tools (ASLR, DEP)",
        "",
        "• Von Neumann Architecture:",
        "   - Input → [CPU ↔ RAM] → Output",
        "            ↕",
        "         [Storage]",
        "",
        "Security implication: Data and instructions share memory → buffer overflow attacks."
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "The transition from 32-bit to 64-bit is fundamentally about security. 64-bit operating systems implement stronger Address Space Layout Randomization (ASLR), randomizing where programs load in memory, making exploits harder. Data Execution Prevention (DEP) also works more effectively on 64-bit systems. The Von Neumann architecture stores both data and instructions in the same memory space. This design decision is the root cause of buffer overflow attacks. When a program receives more input than expected, the excess can overwrite adjacent memory containing executable instructions. Modern defenses like NX bit and stack canaries were invented to mitigate this. Example: The 2017 WannaCry ransomware exploited a buffer overflow in Windows SMB protocol (EternalBlue).",
    img_path=os.path.join(ASSETS, "von_neumann.png"),
    img_left=Inches(7.0), img_top=Inches(1.4), img_w=Inches(6.0), img_h=Inches(5.6)
)

add_content_slide(prs,
    "Types of Computers & Their Security Posture",
    [
        "• Personal Computer (PC)",
        "   - Desktops, laptops — endpoint security focus",
        "",
        "• Server",
        "   - High-value targets; powers websites, databases",
        "   - Often runs Linux; needs hardening",
        "",
        "• Embedded Systems",
        "   - ATMs, traffic lights, medical devices",
        "   - Often outdated and unpatchable",
        "",
        "• Mobile & IoT Devices",
        "   - Smartphones, smart meters, cameras",
        "   - Massive attack surface; default passwords common"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Different computer types require different security strategies. PCs need antivirus, EDR, patch management, and user awareness. Servers are the crown jewels — the 2017 Equifax breach (147 million records) happened because of an unpatched Apache Struts server. In Zambia, servers at mobile network operators and financial institutions run predominantly Linux. Embedded systems and IoT represent the fastest-growing attack surface. The 2016 Mirai botnet infected 600,000+ IoT devices with default passwords and launched record-breaking DDoS attacks. In Zambia, smart meters and cheap CCTV cameras often ship with default credentials like admin/admin."
)

add_section_slide(prs, "2. Operating Systems",
    "Section transition: 2. Operating Systems. Use this moment to reset student attention and preview what security relevance is coming."
)

add_content_slide(prs,
    "What is an OS? The Security Layer",
    [
        "User",
        "  ↓",
        "Applications (Browser, Wireshark)",
        "  ↓",
        "Operating System (Windows, Linux, macOS)",
        "  ↓",
        "Hardware (CPU, RAM, Disk, Network)",
        "",
        "Key OS functions with security impact:",
        "• Process management → Malware hides as legit processes",
        "• Memory management → Buffer overflow protection",
        "• File system → Permissions and access control",
        "• Security & access → Authentication",
        "• Device management → Driver vulnerabilities"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "The OS is the ultimate security gatekeeper. Every application runs WITHIN the OS. If the OS is compromised, nothing running on it can be trusted. This is why rooting (Android) or jailbreaking (iOS) disables critical security protections. Process management: Attackers name malware to mimic legitimate processes (e.g., svch0st.exe). Memory management: Modern OSs implement DEP/NX bits, ASLR, and stack canaries. File system security: The 2020 Twitter hack happened because of poorly configured access controls. Device management: The 2021 PrintNightmare vulnerability was a Windows driver flaw allowing remote code execution."
)

add_content_slide(prs,
    "Windows File System: Where Threats Hide",
    [
        "Security-critical directories:",
        "• C:\\Windows\\System32",
        "   → Core OS files; common malware target",
        "",
        "• C:\\Users\\[name]\\AppData",
        "   → Hidden app data; persistence location",
        "",
        "• C:\\Temp",
        "   → Temporary files; check during investigations",
        "",
        "Malware frequently drops payloads in AppData and Temp because these locations are writable by standard users and often excluded from antivirus scanning."
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Windows dominates the desktop market globally and in Zambia. Incident responders spend significant portions of their careers investigating Windows systems. C:\\Windows\\System32 contains critical DLLs and executables. Malware frequently injects malicious DLLs here (DLL hijacking). The 2020 SolarWinds supply-chain attack compromised Orion software components in system directories. AppData is the #1 location for malware persistence. Ransomware like Ryuk and TrickBot store configuration files here. Temp directories: Malware often extracts from %TEMP%. Command: dir /o-d %TEMP% shows most recent files first. Practical tip: Teach students to enable 'Show hidden items' and navigate these directories.",
    img_path=os.path.join(ASSETS, "windows_fs.png"),
    img_left=Inches(7.0), img_top=Inches(1.4), img_w=Inches(6.0), img_h=Inches(5.6)
)

add_content_slide(prs,
    "Windows Command Prompt: Essential Security Commands",
    [
        "Command        Purpose",
        "─────────      ─────────────────────────────",
        "ipconfig       Show IP config",
        "netstat -an    Show active connections",
        "tasklist       List running processes",
        "systeminfo     System details & patches",
        "net user       Manage user accounts",
        "dir /a         List hidden files",
        "",
        "Principle of Least Privilege:",
        "• Administrator = Full access (use sparingly!)",
        "• Standard User = Limited access (safer)",
        "• Guest = Very limited"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "The Windows Command Prompt and PowerShell are primary tools for incident response. ipconfig /all reveals MAC addresses, DNS servers, and DHCP lease info. Unexpected DNS servers might indicate DNS hijacking. netstat -an lists all active connections. Seeing a connection to an unusual foreign IP on port 4444 (Metasploit) may indicate compromise. tasklist shows running processes — malware often injects into svchost.exe. systeminfo outputs OS version and hotfixes. The 'Hotfix(s)' line is critical — unpatched systems are vulnerable. net user lists local accounts; attackers often create hidden backdoor accounts."
)

add_content_slide(prs,
    "Linux File System: The Security Professional's OS",
    [
        "Why Linux for cybersecurity?",
        "• Free, open-source, highly configurable",
        "• 95%+ of cybersecurity tools are Linux-first",
        "• Powers most web servers and cloud platforms",
        "• Kali Linux is the standard for penetration testing",
        "",
        "Security-critical directories:",
        "• /etc/      → Config files & passwords",
        "• /var/log/  → System logs (security goldmine)",
        "• /tmp/      → Temporary files (malware favorite)",
        "• /root/     → Root admin home"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Linux is the dominant OS in cybersecurity and server infrastructure. Over 80% of web servers run Linux. In Zambia, MTN and Airtel core network infrastructure, banking servers, and government systems heavily rely on Linux. /etc/ contains system-wide configs: /etc/passwd (user accounts), /etc/shadow (hashed passwords), /etc/hosts (can be poisoned by attackers), /etc/crontab (persistence). /var/log/ is the security goldmine: auth.log shows failed logins, apache2/access.log shows web requests. /tmp/ is world-writable and a favorite malware drop location. Zambian context: ZICTA and major banks run Linux-based infrastructure. Linux skills are directly employable.",
    img_path=os.path.join(ASSETS, "linux_fs.png"),
    img_left=Inches(7.0), img_top=Inches(1.4), img_w=Inches(6.0), img_h=Inches(5.6)
)

add_content_slide(prs,
    "Linux Commands & File Permissions",
    [
        "Essential Linux Commands:",
        "• ls -la /etc       → List files with details",
        "• cat /etc/passwd   → View file contents",
        "• grep 'Failed' /var/log/auth.log",
        "• ps aux            → Show all running processes",
        "• sudo chmod 600 file.txt",
        "",
        "Permission Breakdown:  -rwxr-xr--",
        "• - = file type",
        "• rwx = Owner (read/write/execute)",
        "• r-x = Group        r-- = Others",
        "",
        "Numeric: r=4, w=2, x=1"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Linux command-line proficiency is non-negotiable. ls -la reveals hidden files (those starting with a dot). Attackers hide config files and SSH keys as hidden files. .bashrc modifications can establish persistence. grep is one of the most powerful tools for log analysis: grep 'Failed password' /var/log/auth.log | wc -l counts failed login attempts. File permissions: chmod 600 ~/.ssh/id_rsa is critical — SSH private keys should NEVER be readable by others. If it's 644, anyone can steal it. Privilege escalation often starts with misconfigured permissions. find / -perm -4000 2>/dev/null discovers all SUID files — a common privilege escalation step.",
    img_path=os.path.join(ASSETS, "linux_permissions.png"),
    img_left=Inches(6.8), img_top=Inches(1.4), img_w=Inches(6.2), img_h=Inches(5.6)
)

add_section_slide(prs, "3. Introduction to Programming Logic",
    "Section transition: 3. Introduction to Programming Logic. Use this moment to reset student attention and preview what security relevance is coming."
)

add_content_slide(prs,
    "Why Programming Logic Matters in Cybersecurity",
    [
        "Logic skills enable you to:",
        "",
        "• Read and write basic scripts for automation",
        "   → Automate log parsing, network scans",
        "",
        "• Understand how malware and exploits work",
        "   → Malware is code written with malicious intent",
        "",
        "• Analyze code for vulnerabilities",
        "   → Spot SQL injection, buffer overflow",
        "",
        "• Write simple security tools",
        "   → Python is the #1 language in cybersecurity",
        "",
        "• Over 90% of tools in Kali Linux are Python or Bash"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Programming is the closest thing to a superpower in cybersecurity. It allows automation of repetitive tasks and deep understanding of software. SOC analysts write Python scripts to parse 100,000 firewall logs. Penetration testers write Bash scripts to run Nmap against target lists. Malware analysts use Python to unpack suspicious files. Understanding malware: At its core, malware is just software. Ransomware uses loops to iterate through files, conditionals to check extensions, and file system APIs to encrypt data. Vulnerability analysis: Many vulnerabilities are logic errors. A missing input validation check leads to SQL injection. Industry data: Python is the most requested programming skill in cybersecurity job surveys."
)

add_content_slide(prs,
    "Variables, Data Types & Security Examples",
    [
        "A variable is a named container that stores a value.",
        "",
        "  student_name = 'Chanda Mwale'  # String",
        "  age = 22                        # Integer",
        "  gpa = 3.5                       # Float",
        "  is_enrolled = True              # Boolean",
        "",
        "Security-relevant data types:",
        "• String  → usernames, passwords, IPs, domains",
        "• Integer → port numbers (22, 80, 443)",
        "• Boolean → access_granted, is_admin",
        "• List    → ['192.168.1.1', '10.0.0.1']"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Variables are the building blocks of every program, including every piece of malware and security tool. String security: Passwords stored as plain strings in memory can be recovered using memory dumps. IP addresses in strings are parsed in log analysis using regex. Integer security: Port numbers are integers 0-65535. Well-known ports (0-1023) require admin privileges. Integer overflow is serious — when an integer exceeds its max value, it wraps around. The 2014 Heartbleed vulnerability involved integer-related flaws. Boolean security: Authentication systems use booleans like is_authenticated, is_admin. Logic flaws where booleans are improperly checked lead to unauthorized access. List security: Lists of IPs are used in firewall rules. Iterating over user inputs without validation can lead to injection attacks."
)

add_content_slide(prs,
    "Control Structures: If/Else & Loops in Security",
    [
        "If/Else — Decision Making:",
        "  if password == 'Secure@2026':",
        "      print('Access granted')",
        "  else:",
        "      print('Access denied')",
        "",
        "While Loop — Account Lockout:",
        "  attempts = 0",
        "  while attempts < 3:",
        "      password = input('Enter password: ')",
        "      if password == 'Secure@2026':",
        "          print('Access granted'); break",
        "      else:",
        "          attempts += 1",
        "  if attempts == 3:",
        "      print('Account locked.')"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Control structures are where security logic lives. Every access control decision, rate limit, and account lockout policy uses conditionals and loops. If/Else: Authentication if (password_hash == stored_hash) is the core of every login system. Authorization if (user.role == 'admin') determines access. Input validation if (input_length > 1000) prevents buffer overflow. A real bug: In 2019, Apple FaceTime had a logic flaw where calls connected before the callee answered. Loops: Port scanners like Nmap use loops: for port in range(1, 65536): scan(port). Password crackers use loops over wordlists. Loop security bugs: Infinite loops can cause Denial of Service if attackers trigger them."
)

add_content_slide(prs,
    "Functions, Pseudocode & Bash Scripting",
    [
        "Function — Reusable security check:",
        "  def check_password_strength(pw):",
        "      if len(pw) < 8: return 'Weak'",
        "      elif pw.isalpha(): return 'Weak'",
        "      else: return 'Strong'",
        "",
        "Pseudocode — Planning before coding:",
        "  START → INPUT username → INPUT password",
        "  IF username EXISTS AND password MATCHES:",
        "      GRANT access → LOG success",
        "  ELSE:",
        "      INCREMENT failed_attempts",
        "      IF failed_attempts >= 3: LOCK account",
        "",
        "Bash — Network health check:",
        "  #!/bin/bash",
        "  ping -c 4 8.8.8.8",
        "  ip a | grep 'inet '"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Functions allow security checks to be defined once and reused everywhere. The password strength function mirrors real password policy enforcement. NIST recommends 8+ characters. Modern guidance emphasizes LENGTH over complexity because users reuse complex passwords. A passphrase like 'Correct-Horse-Battery-Staple!' is more secure than 'P@ssw0rd1'. Pseudocode is invaluable for security planning. Before writing code, architects pseudocode authentication flows and incident response playbooks. Bash scripting: Automated backups ensure logs are preserved. Network monitoring scripts ping critical infrastructure. Important Bash security note: Always validate inputs. A script that passes user input directly to a command without sanitization is vulnerable to command injection — an OWASP Top 10 vulnerability."
)

add_section_slide(prs, "4. Mathematics for Cybersecurity",
    "Section transition: 4. Mathematics for Cybersecurity. Use this moment to reset student attention and preview what security relevance is coming."
)

add_content_slide(prs,
    "Number Systems: Decimal, Binary, Hexadecimal",
    [
        "Decimal (Base 10) — Human system",
        "   Digits: 0-9     Example: 255",
        "",
        "Binary (Base 2) — Computer system",
        "   Digits: 0, 1    Example: 1111 1111",
        "   1 bit = one digit    8 bits = 1 byte",
        "",
        "Hexadecimal (Base 16) — Compact notation",
        "   Digits: 0-9, A=10..F=15",
        "   Example: FF",
        "",
        "Cybersecurity uses:",
        "• Binary  → IP addressing, subnet masks, Boolean logic",
        "• Hex     → MAC addresses, memory addresses, hashes",
        "• Decimal → Human-readable port numbers, IP octets"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Number systems are the language of computers and cybersecurity professionals. Every IP address, MAC address, cryptographic hash, and memory pointer uses these systems. Binary is essential for IP subnetting. Subnet masks like 255.255.255.0 are /24 in CIDR notation (24 binary 1s). Boolean logic operates on individual bits. Hexadecimal is the 'shorthand' for binary — each hex digit represents exactly 4 bits. Cybersecurity applications: MAC addresses (00:1A:2B:3C:4D:5E) are 6 bytes in hex. Cryptographic hashes like SHA-256 are displayed in hex. Memory addresses in debuggers are shown in hex. A SHA-256 hash looks like e3b0c44298fc1c149afbf4c8996fb924... — 64 hex characters representing 256 bits."
)

add_content_slide(prs,
    "Binary Conversion with Examples",
    [
        "Decimal to Binary: Divide by 2, collect remainders",
        "   Convert 45 to binary:",
        "   45 ÷ 2 = 22 r 1",
        "   22 ÷ 2 = 11 r 0",
        "   11 ÷ 2 =  5 r 1",
        "    5 ÷ 2 =  2 r 1",
        "    2 ÷ 2 =  1 r 0",
        "    1 ÷ 2 =  0 r 1",
        "   Read bottom-up: 101101",
        "",
        "Binary to Decimal: Multiply by position values",
        "   128  64  32  16   8   4   2   1",
        "     1   0   1   0   1   1   0   0",
        "   = 128 + 32 + 8 + 4 = 172",
        "",
        "This octet (172) appears in IP addresses like 172.16.0.1"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Binary conversion is a practical skill used constantly in network engineering and forensics. Each octet in an IPv4 address is an 8-bit binary number. Example: 192.168.1.1 = 11000000.10101000.00000001.00000001. Subnetting application: To determine if two IPs are on the same network, perform a bitwise AND between each IP and the subnet mask. Forensics application: File signatures (magic numbers) are examined in hex/binary to identify files even when extensions are changed. Example: A PNG file starts with hex 89 50 4E 47. Teaching tip: Use a whiteboard to manually work through conversions. Have students practice with numbers relevant to them."
)

add_content_slide(prs,
    "Hexadecimal Conversion: The Binary Shortcut",
    [
        "Binary to Hex: Group into 4 bits",
        "   Binary:  1010  1100",
        "   Hex:        A     C",
        "   So 10101100 binary = AC hex = 172 decimal",
        "",
        "Decimal to Hex: Divide by 16",
        "   Convert 255 to hex:",
        "   255 ÷ 16 = 15 remainder 15 → F",
        "    15 ÷ 16 =  0 remainder 15 → F",
        "   Result: FF",
        "",
        "Quick Reference:",
        "   Dec  0   5   9   10   15   16   255",
        "   Hex  0    5    9    A    F    10       FF"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "The binary-to-hex 'group into 4 bits' method is the most important conversion technique for cybersecurity professionals. Why hex is everywhere: MAC addresses are 48-bit addresses displayed as 6 pairs of hex digits. The first 3 bytes identify the manufacturer. Attackers use MAC vendor prefixes to identify device types on networks. Cryptographic hashes: SHA-256 hashes are long hex strings. If even one bit changes in a file, the hash changes completely — this is file integrity checking. Memory forensics: Memory dumps display addresses in hex. Malware analysis might show MOV EAX, 0x00401234. Practical exercise: Convert web colors from hex to binary. Royal blue #4169E1 breaks down as 41=01000001, 69=01101001, E1=11100001. This demonstrates hex is just compact binary.",
    img_path=os.path.join(ASSETS, "binary_hex_table.png"),
    img_left=Inches(7.0), img_top=Inches(1.4), img_w=Inches(6.0), img_h=Inches(5.6)
)

add_content_slide(prs,
    "Boolean Algebra & Logic Gates",
    [
        "Boolean algebra uses only TRUE (1) and FALSE (0).",
        "",
        "Logic Operations:",
        "• AND (∧)  → Both must be true      1 AND 1 = 1",
        "• OR (∨)   → At least one true      0 OR 1 = 1",
        "• NOT (¬)  → Inverts the value      NOT 1 = 0",
        "• XOR (⊕)  → One or other, not both 1 XOR 1 = 0",
        "",
        "AND Gate Truth Table:",
        "   A | B | A AND B",
        "   0 | 0 |    0",
        "   0 | 1 |    0",
        "   1 | 0 |    0",
        "   1 | 1 |    1"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Boolean logic is the mathematical foundation of digital computing. Every CPU instruction, firewall rule, and access control decision uses Boolean operations. AND in security: Multi-factor authentication is essentially AND: access_granted = password_correct AND otp_valid. Both must be true. OR in security: Role-based access control uses OR: can_view_report = (role == 'admin') OR (role == 'analyst'). NOT in security: Deny lists use NOT: allow_access = NOT(ip_in_blocklist). Logic gates in hardware: CPUs are built from billions of transistors arranged as logic gates. Set theory connection: Finding the intersection between two threat feeds is an AND operation. Combining feeds is a union (OR)."
)

add_content_slide(prs,
    "XOR: The Foundation of Encryption",
    [
        "XOR Truth Table:",
        "   A | B | A XOR B",
        "   0 | 0 |    0",
        "   0 | 1 |    1",
        "   1 | 0 |    1",
        "   1 | 1 |    0",
        "",
        "The Magic Property:",
        "   If you XOR a value with a key, then XOR the result",
        "   with the SAME key, you get the original value back!",
        "",
        "Example:",
        "   Plaintext:  1010",
        "   Key:        1100",
        "   Ciphertext: 0110   (1010 XOR 1100)",
        "   Decrypt:    1010   (0110 XOR 1100) ← Original!"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "XOR (Exclusive OR) is arguably the single most important Boolean operation in cryptography. The reversible property: A XOR B = C, and C XOR B = A. Real-world applications: 1. One-Time Pad (OTP): The only encryption proven unconditionally secure. Cold War spies used OTPs. 2. Stream Ciphers: ChaCha20 in TLS/HTTPS uses XOR to combine a keystream with plaintext. 3. Block Ciphers: AES uses XOR extensively in its rounds to mix data with round keys. AES protects everything from WhatsApp to government classified documents. 4. Hash Functions: SHA-256 uses XOR-like operations to create avalanche effects. Malware obfuscation: Attackers also use XOR! Simple malware XORs payloads with one-byte keys to evade antivirus. Security analysts brute-force XOR keys to decode hidden strings."
)

add_section_slide(prs, "5. Networking Essentials",
    "Section transition: 5. Networking Essentials. Use this moment to reset student attention and preview what security relevance is coming."
)

add_content_slide(prs,
    "What is a Network? Types & Security Relevance",
    [
        "A network is two or more computers connected to share resources.",
        "Virtually all cybersecurity work involves networks.",
        "",
        "Network Types:",
        "• LAN  → School lab, office building",
        "• WAN  → National backbone (ZAMTEL)",
        "• MAN  → City-wide network (Lusaka)",
        "• WLAN → Home WiFi, coffee shop hotspots",
        "• VPN  → Encrypted tunnel for remote work",
        "",
        "Security perspective:",
        "• LANs need segmentation to limit breach spread",
        "• WANs need encryption",
        "• WLANs are vulnerable to eavesdropping",
        "• VPN misconfigurations are common breach causes"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Networking is the backbone of cybersecurity. Every attack, defense mechanism, and investigation relies on how devices communicate. LAN Security: A flat LAN allows malware to spread rapidly. WannaCry spread laterally across LANs using SMB vulnerabilities. Network segmentation contains outbreaks. WAN Security: Because traffic traverses third-party infrastructure, encryption is mandatory. In Zambia, the ZAMTEL national backbone connects cities. WLAN Security: Wireless networks broadcast over radio waves, making them easy to intercept. Public WiFi in Lusaka cafes is risky — attackers set up rogue access points. VPN Security: The 2021 Colonial Pipeline attack began through a compromised VPN account with no MFA.",
    img_path=os.path.join(ASSETS, "network_types.png"),
    img_left=Inches(7.0), img_top=Inches(1.4), img_w=Inches(6.0), img_h=Inches(5.6)
)

add_content_slide(prs,
    "The OSI Model: 7 Layers of Communication",
    [
        "Layer | Name          | Function            | Protocol",
        "───── | ───────────── | ─────────────────── | ──────────────",
        "  7   | Application   | User interface      | HTTP, DNS, FTP",
        "  6   | Presentation  | Encryption/format   | SSL/TLS, JPEG",
        "  5   | Session       | Manages connections | NetBIOS, RPC",
        "  4   | Transport     | End-to-end delivery | TCP, UDP",
        "  3   | Network       | Routing, IP address | IP, ICMP",
        "  2   | Data Link     | MAC addressing      | Ethernet, WiFi",
        "  1   | Physical      | Cables, signals     | Cables, radio",
        "",
        "Memory Aid: All People Seem To Need Data Processing",
        "",
        "Cybersecurity relevance: Attacks happen at specific layers."
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "The OSI model is the conceptual framework every cybersecurity professional must internalize. Layer 7 — Application: SQL injection, XSS, phishing, malware downloads. Layer 6 — Presentation: TLS/SSL encryption, certificate spoofing. Layer 5 — Session: Session hijacking steals valid session tokens. Layer 4 — Transport: SYN flood attacks exhaust server resources. Port scanning targets this layer. Layer 3 — Network: IP spoofing, DDoS amplification, route hijacking. Layer 2 — Data Link: ARP spoofing and MAC flooding. Layer 1 — Physical: Physical security and signal interception. Teaching analogy: Use the Zampost letter analogy. Application = what you write. Presentation = language. Session = ongoing correspondence. Transport = registered mail vs flyer. Network = routing Lusaka to Kitwe. Data Link = truck route. Physical = roads.",
    img_path=os.path.join(ASSETS, "osi_model.png"),
    img_left=Inches(7.0), img_top=Inches(1.3), img_w=Inches(5.8), img_h=Inches(5.8)
)

add_content_slide(prs,
    "Mapping Cyber Attacks to OSI Layers",
    [
        "Layer    Attack Type                    Defense",
        "─────    ───────────────────────────    ──────────────────────────────",
        "L7       SQL Injection, Phishing        WAF, Secure Coding",
        "L6       Weak TLS/SSL                   Patch Management",
        "L5       Session Hijacking              HTTPS, Secure Cookies",
        "L4       SYN Flood, Port Scanning       Firewalls, IDS/IPS",
        "L3       IP Spoofing, DDoS              ACLs, DDoS Mitigation",
        "L2       ARP Spoofing                   Dynamic ARP Inspection",
        "L1       Wiretapping, Physical Theft    Encryption, Locks",
        "",
        "Key insight: Defense in Depth means protecting EVERY layer."
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "This slide connects abstract theory to real attacks. Effective defense requires protection at multiple layers — Defense in Depth. Real-world examples: L7 — SQL Injection: The 2017 Equifax breach (147 million records) was partly due to unpatched web app vulnerabilities. L4 — SYN Flood: The 2000 attack on Yahoo! and eBay used SYN floods. L3 — DDoS Amplification: The 2016 Dyn attack used IoT botnets and DNS amplification to take down Twitter, Netflix, and Reddit. L2 — ARP Spoofing: On local networks, attackers send fake ARP messages to become the man-in-the-middle. Defense: Dynamic ARP Inspection. L1 — Physical attacks: The 2013 Target breach involved physical access to POS terminals. In Zambia, server rooms need biometric access and CCTV."
)

add_content_slide(prs,
    "The TCP/IP Model: Real-World Implementation",
    [
        "TCP/IP Layer       | Corresponds to OSI    | Protocols",
        "─────────────────  | ──────────────────    | ───────────────────────────",
        "Application        | Layers 5, 6, 7        | HTTP, HTTPS, DNS, FTP, SSH",
        "Transport          | Layer 4               | TCP, UDP",
        "Internet           | Layer 3               | IP, ICMP, ARP",
        "Network Access     | Layers 1, 2           | Ethernet, WiFi",
        "",
        "OSI is theory → used for learning and explaining.",
        "TCP/IP is implementation → used by actual devices.",
        "",
        "In practice, cybersecurity professionals use BOTH models."
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "While the OSI model is excellent for learning, the TCP/IP model is what routers, operating systems, firewalls, and Wireshark actually implement. Application Layer (TCP/IP): HTTP (port 80) — never use for sensitive data. HTTPS (port 443) — encrypted with TLS, required for banking. DNS (port 53) — vulnerable to poisoning, hijacking, tunneling. SSH (port 22) — encrypted remote administration. FTP (port 21) — no encryption, use SFTP instead. Transport Layer: Determines reliable (TCP) vs fast/unreliable (UDP). Firewalls make most decisions here. Internet Layer: IP addresses, routing, ICMP (ping). VPNs like IPsec operate here. Network Access Layer: Ethernet frames, MAC addresses, WiFi signals. Critical for local network attacks. Tool context: Wireshark shows the TCP/IP model in action — Ethernet header, IP header, TCP/UDP header, and application payload.",
    img_path=os.path.join(ASSETS, "osi_tcpip_comparison.png"),
    img_left=Inches(7.0), img_top=Inches(1.4), img_w=Inches(6.0), img_h=Inches(5.6)
)

add_content_slide(prs,
    "TCP vs UDP: The Security Trade-off",
    [
        "Feature        TCP                         UDP",
        "─────────      ─────────────────────────   ─────────────────────────",
        "Connection     Connection-oriented         Connectionless",
        "Reliability    Guaranteed delivery         No guarantee",
        "Speed          Slower (handshake)          Faster",
        "Use Cases      Web, email, file transfer   Streaming, DNS, VoIP",
        "Cyber Risk     SYN flood attacks           DDoS amplification",
        "",
        "TCP Three-Way Handshake:",
        "   Client → SYN → Server",
        "   Client ← SYN-ACK ← Server",
        "   Client → ACK → Server",
        "",
        "Attackers exploit the handshake with SYN floods."
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Understanding TCP vs UDP is fundamental for network security. TCP: Connection-oriented with a three-way handshake (SYN → SYN-ACK → ACK). Reliable with sequence numbers and retransmissions. Used for web browsing, email, file transfer, remote access. Security risk: SYN flood attacks exploit the handshake by sending thousands of SYN packets without completing them. Mitigation: SYN cookies, rate limiting. UDP: Connectionless — no handshake. No guarantee of delivery. Used for DNS, video streaming, VoIP, online gaming. Security risk: UDP amplification attacks. Attackers send small spoofed requests to servers (DNS/NTP), which respond with larger packets to the victim. The 2016 Dyn attack used this. Analogy: TCP = registered mail with return receipt. UDP = dropping a flyer through a door.",
    img_path=os.path.join(ASSETS, "tcp_handshake.png"),
    img_left=Inches(6.8), img_top=Inches(1.4), img_w=Inches(6.2), img_h=Inches(5.6)
)

add_content_slide(prs,
    "IP Addressing & Subnetting Basics",
    [
        "IPv4 Address:",
        "   32-bit address written as 4 octets",
        "   Example: 192.168.1.100",
        "   Each octet = 8 bits = 0-255",
        "",
        "Private IP Ranges (not routable on internet):",
        "   10.0.0.0      – 10.255.255.255",
        "   172.16.0.0    – 172.31.255.255",
        "   192.168.0.0   – 192.168.255.255",
        "   127.0.0.1     → Loopback / localhost",
        "",
        "Subnetting Example:",
        "   IP:       192.168.1.100",
        "   Mask:     255.255.255.0  (/24)",
        "   Network:  192.168.1.0",
        "   Hosts:    192.168.1.1 – 192.168.1.254"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "IP addressing is the addressing system of the internet. IPv4 is still dominant in most networks. 32-bit address space provides ~4.3 billion unique addresses. Private IP ranges (RFC 1918) are reserved for internal networks: 10.0.0.0/8, 172.16.0.0/12, 192.168.0.0/16. When you connect to MTN or Airtel in Zambia, your phone typically gets a private IP (192.168.x.x) while a public IP faces the internet. Subnetting and security: Subnetting divides networks into smaller segments. From a security perspective, this is network segmentation — one of the most effective security controls. If finance is on 192.168.10.0/24 and engineering on 192.168.20.0/24, a router/firewall between them restricts access. If ransomware infects one subnet, segmentation prevents spread. CIDR notation: /24 means the first 24 bits are the network portion, equivalent to 255.255.255.0."
)

add_content_slide(prs,
    "Key Protocols: DNS, DHCP, HTTP vs HTTPS",
    [
        "DNS (Domain Name System) — Port 53",
        "   Translates names to IP addresses",
        "   ⚠ Threat: DNS poisoning redirects to fake sites",
        "",
        "DHCP — Ports 67/68",
        "   Automatically assigns IP addresses",
        "   ⚠ Threat: Rogue DHCP server hijacks traffic",
        "",
        "HTTP (Port 80)  →  Unencrypted — readable on public WiFi",
        "HTTPS (Port 443) →  Encrypted with TLS — safe for banking",
        "",
        "SSH (Port 22) →  Encrypted remote access (replaces Telnet)"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "DNS is the 'phonebook of the internet.' Security threats: DNS Spoofing/Poisoning corrupts DNS cache to redirect users to malicious sites. DNS Tunneling: Malware uses DNS queries to secretly exfiltrate data because DNS is usually allowed through firewalls. DNS Hijacking: Attackers modify DNS settings to point to malicious servers. Defenses: DNSSEC adds cryptographic verification. DNS filtering blocks known malicious domains. DHCP automatically assigns IP addresses. Rogue DHCP Server: An attacker sets up a malicious DHCP server, giving devices attacker-controlled configurations — typically a rogue gateway intercepting all traffic. DHCP Starvation floods the DHCP server with requests, exhausting the IP pool. HTTP vs HTTPS: HTTP sends data in plaintext. Anyone on the same network can use Wireshark to read passwords. HTTPS encrypts with TLS. The padlock icon indicates HTTPS. Never enter passwords on HTTP sites. SSH provides encrypted remote command-line access, replacing insecure Telnet. Brute force attacks target SSH on port 22. Best practices: use key-based auth, disable password auth, implement fail2ban."
)

add_content_slide(prs,
    "Common Port Numbers & Their Security Risks",
    [
        "Port    Protocol    Use",
        "────    ────────    ─────────────────────────────",
        "22      SSH         Secure remote access",
        "23      Telnet      Insecure — AVOID",
        "25      SMTP        Sending email",
        "53      DNS         Domain resolution",
        "80      HTTP        Web (unencrypted)",
        "443     HTTPS       Web (encrypted)",
        "445     SMB         Windows file sharing",
        "3389    RDP         Windows remote desktop",
        "3306    MySQL       Database",
        "",
        "Attackers scan for open ports to find ways in."
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Port numbers are the 'doors' through which network services communicate. Port 22 (SSH): The most attacked port. Automated bots constantly brute-force SSH. Defense: disable root login, use keys, implement fail2ban. Port 23 (Telnet): Ancient plaintext protocol. Absolutely no reason to use it in 2026. Critical finding if open. Port 25 (SMTP): Misconfigured SMTP servers can be exploited as open relays for spam. Port 53 (DNS): Critical but vulnerable to poisoning and tunneling. Port 80 (HTTP): Unencrypted web traffic. Port 443 (HTTPS): Standard for secure web. When students use Airtel Money or Zanaco Internet Banking, data travels over port 443. Port 445 (SMB): Target of EternalBlue exploit used by WannaCry and NotPetya. Microsoft strongly recommends disabling SMBv1. Port 3389 (RDP): Major ransomware entry point. FBI reported RDP as the most common ransomware entry in 2019. Never expose RDP directly to the internet. Port 3306 (MySQL): Database ports should never be internet-facing."
)

add_section_slide(prs, "Module 1 Review & What's Next",
    "Section transition: Module 1 Review & What's Next. Use this moment to reset student attention and preview what security relevance is coming."
)

add_content_slide(prs,
    "Module 1 Revision: Key Takeaways",
    [
        "Computer Fundamentals:",
        "   ✓ IPOS model — every stage is an attack surface",
        "   ✓ CPU, RAM, Storage, NIC — security implications",
        "",
        "Operating Systems:",
        "   ✓ Windows & Linux file systems",
        "   ✓ Command-line navigation & Linux permissions",
        "",
        "Programming Logic:",
        "   ✓ Variables, if/else, loops, functions",
        "   ✓ Pseudocode and Bash scripting",
        "",
        "Mathematics:",
        "   ✓ Binary, decimal, hexadecimal conversions",
        "   ✓ Boolean algebra and XOR in cryptography",
        "",
        "Networking:",
        "   ✓ OSI & TCP/IP models, IP addressing, subnetting",
        "   ✓ DNS, DHCP, ports, and protocol security risks"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "This revision summarizes foundational knowledge needed before Module 2. Connect the topics: A SOC analyst investigating Windows malware uses OS knowledge to find malware in AppData. They use programming logic to understand persistence mechanisms. They use networking knowledge to trace C2 connections over TCP port 443. They use binary/hex knowledge to analyze file signatures. They use Linux skills to run analysis tools on Kali Linux. All these skills work together. Module 1 provides the vocabulary and mental models. Module 2 applies them to actual attacks and defenses. Assessment reminder: Practical labs (OS installation and configuration) and a quiz covering number systems and networking fundamentals."
)

add_content_slide(prs,
    "What's Next? Module 2 Preview",
    [
        "Now that you understand how computers, OSs, and networks work...",
        "",
        "Module 2: Cybersecurity Principles and Defense",
        "",
        "• We will learn HOW attackers exploit everything we just studied:",
        "   → Port 23 Telnet? That's a vulnerability.",
        "   → XOR operations? That's how malware hides.",
        "   → /tmp and C:\\Temp? Malware's favorite hiding spots.",
        "",
        "Coming up:",
        "   • CIA Triad, Threats & Attack Types",
        "   • Network Security: Firewalls, VPNs, Wireless",
        "   • Cryptography: Encryption, Hashing, Digital Signatures",
        "   • Authentication, Access Control, IDS/IPS"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Final hook for Module 2. Key message: 'You now understand the machinery. In Module 2, we start looking at how that machinery gets attacked — and how to defend it.' Make explicit connections: 'Remember the OSI model? In Module 2, you'll learn where firewalls (Layer 3-4), IDS/IPS (Layer 4-7), and WAFs (Layer 7) fit.' 'Remember binary and XOR? In Module 2, you'll see how AES and SHA-256 use these concepts.' 'Remember Windows user accounts? In Module 2, you'll learn about Active Directory, RBAC, and MFA.' 'Remember TCP ports? In Module 2, you'll configure firewall rules.' Motivational closing: The cybersecurity industry in Zambia and globally is growing. MTN Zambia, Airtel, Zanaco, ZESCO, ZICTA, NGOs, and international organizations all need skilled cybersecurity professionals. The foundation built in Module 1 is the first step toward an in-demand, well-paying career."
)

# Save presentation
output_path = "/Users/michael/Documents/Programming/edutrack-lms/Module_1_Foundations_of_Computing_and_Mathematics.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
