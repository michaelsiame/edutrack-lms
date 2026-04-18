#!/usr/bin/env python3
"""
Module 1: Foundations of Computing and Mathematics - PPTX Generator
Cybersecurity Certificate Program | EduTrack LMS
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
DARK_BLUE = RGBColor(0x1A, 0x23, 0x7E)
ACCENT_ORANGE = RGBColor(0xF4, 0x4C, 0x00)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1F, 0x29, 0x33)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.333), Inches(1.0))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0xFF, 0xB7, 0x4D)
    p.alignment = PP_ALIGN.CENTER
    
    # Add notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = (
        "Welcome to Module 1: Foundations of Computing and Mathematics. "
        "This module is the bedrock of the entire Cybersecurity Certificate Program. "
        "Every advanced concept in cybersecurity — from malware analysis to network defense — "
        "depends on understanding how computers process data, how operating systems manage resources, "
        "how logical thinking enables automation, how number systems underpin cryptography, and how networks enable both communication and attack vectors. "
        "As instructors, emphasize that this is NOT a basic computing class. Every topic here has direct security implications."
    )
    return slide

def add_section_slide(prs, section_title):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = ACCENT_ORANGE
    bg.line.fill.background()
    
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = f"Section transition: {section_title}. Use this moment to reset student attention and preview what security relevance is coming."
    return slide

def add_content_slide(prs, title, bullets, notes=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.8))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(22)
        p.font.color.rgb = DARK_TEXT
        p.level = 0
        p.space_after = Pt(12)
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes
    return slide

# ======================= SLIDES =======================

# Slide 1: Title
add_title_slide(prs,
    "Module 1: Foundations of Computing and Mathematics",
    "Cybersecurity Certificate Program | EduTrack LMS | Zambia"
)

# Slide 2: Why Foundations Matter
add_content_slide(prs,
    "Why Foundations Matter in Cybersecurity",
    [
        "• You cannot secure what you do not understand.",
        "• Every cyberattack exploits how computers, OSs, networks, and code work.",
        "• Module 1 builds the mental models needed for threat detection, incident response, and defense.",
        "• From Zambia's banks to mobile money platforms — all run on these same fundamentals.",
        "• Direct career path: SOC Analyst, Network Security Technician, Junior Security Analyst."
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Open with a compelling example: In 2021, the Colonial Pipeline ransomware attack shut down fuel distribution across the US East Coast. "
    "The attackers gained access through a compromised VPN account — exploiting basic networking and OS weaknesses. "
    "In Zambia, MTN and Airtel mobile money systems, Zanaco banking infrastructure, and ZESCO systems all depend on the exact fundamentals we cover here. "
    "Tell students: 'Attackers know these basics better than many defenders. Our job is to out-learn them.' "
    "The CIA Triad (Confidentiality, Integrity, Availability) depends on hardware working correctly, OS permissions being set properly, scripts automating defenses, cryptographic math protecting data, and networks being segmented securely."
)

# ======================= SECTION: COMPUTER FUNDAMENTALS =======================
add_section_slide(prs, "1. Computer Fundamentals")

# Slide 3: What is a Computer - IPOS Model
slide = add_content_slide(prs,
    "What is a Computer? The IPOS Model",
    [
        "Input  →  Processing  →  Output  →  Storage",
        "",
        "• Input:    Keyboard, mouse, biometric scanner, card reader",
        "• Processing: CPU executes instructions (the 'brain')",
        "• Output:   Monitor, printer, speakers",
        "• Storage:  HDD, SSD, USB — permanent data retention",
        "",
        "Cybersecurity angle: Every stage is an attack surface."
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "The IPOS model is not just academic theory — it maps directly to attack surfaces. "
    "Input stage: Keyloggers capture keystrokes before they reach the CPU. Hardware keyloggers can be physically inserted between keyboard and USB port. "
    "Processing stage: Side-channel attacks like Spectre and Meltdown exploit how CPUs predict and execute instructions, leaking sensitive data across processes. "
    "Output stage: Shoulder surfing in Lusaka internet cafes targets the output phase. Screen privacy filters are a physical security control. "
    "Storage stage: Unencrypted USB drives are a massive data loss vector. In 2022, the UK government fined organizations millions for lost USB drives containing personal data. "
    "Zambian context: Banks like Zanaco and ZESCO offices handle customer data on local servers. Understanding where data lives at each IPOS stage is essential for proper data classification and protection."
)

# Slide 4: Hardware Components
add_content_slide(prs,
    "Hardware Components & Security Relevance",
    [
        "• CPU (Central Processing Unit)",
        "   - Executes instructions; speed in GHz",
        "   - Security relevance: Cryptographic operations, side-channel attacks (Spectre/Meltdown)",
        "",
        "• RAM (Random Access Memory)",
        "   - Temporary, volatile storage",
        "   - Security relevance: Fileless malware lives ONLY in RAM to avoid disk detection",
        "",
        "• Storage (HDD / SSD / USB)",
        "   - Permanent retention of data",
        "   - Security relevance: Encryption at rest protects stolen devices",
        "",
        "• NIC (Network Interface Card)",
        "   - Every NIC has a unique MAC address",
        "   - Security relevance: MAC spoofing, device identification on networks"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "CPU: Modern processors include security extensions like Intel SGX and AMD SEV that create encrypted enclaves for sensitive operations. However, vulnerabilities like Spectre (2018) and Meltdown (2018) showed that speculative execution — a performance optimization — could leak passwords and encryption keys across process boundaries. Over 90% of processors worldwide were affected.\n\n"
    "RAM: Fileless malware (also called 'living-off-the-land' attacks) is one of the fastest-growing threat categories. By never writing to disk, it evades traditional antivirus signature scanning. Tools like PowerShell Empire and Cobalt Strike operate primarily in memory. Memory forensics tools like Volatility are essential for incident responders to detect these threats.\n\n"
    "Storage: Full Disk Encryption (FDE) using BitLocker (Windows) or LUKS (Linux) ensures that if a laptop is stolen from a Zambian government office or NGO, the data remains unreadable without the decryption key. In 2023, the average cost of a data breach involving lost devices was $4.45 million globally.\n\n"
    "NIC: MAC addresses are burned into hardware and theoretically unique. However, tools like macchanger (Linux) and registry edits (Windows) allow attackers to spoof MAC addresses to bypass network access controls (NAC) or frame innocent users."
)

# Slide 5: System Architecture
add_content_slide(prs,
    "System Architecture: 32-bit vs 64-bit & Von Neumann",
    [
        "• 32-bit Systems:",
        "   - Can address up to 4GB of RAM",
        "   - Limited security features, being phased out",
        "",
        "• 64-bit Systems:",
        "   - Can address up to 16 exabytes theoretically",
        "   - Required for modern security tools (ASLR, DEP, modern encryption)",
        "",
        "• Von Neumann Architecture:",
        "   Input → [CPU ↔ RAM] → Output",
        "            ↕",
        "         [Storage]",
        "",
        "Security implication: Data and instructions share memory → buffer overflow attacks."
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "The transition from 32-bit to 64-bit is not just about more RAM — it's fundamentally about security. 64-bit operating systems implement stronger versions of Address Space Layout Randomization (ASLR), which randomizes where programs load in memory, making it harder for attackers to predict memory addresses for exploits. Data Execution Prevention (DEP) also works more effectively on 64-bit systems.\n\n"
    "The Von Neumann architecture stores both data and instructions in the same memory space. This design decision — made in the 1940s — is the root cause of buffer overflow attacks, one of the most common and dangerous vulnerability classes. When a program receives more input than expected, the excess can overwrite adjacent memory that contains executable instructions, allowing an attacker to inject and run malicious code. Modern defenses like NX bit (No-eXecute) and stack canaries were invented specifically to mitigate this architectural weakness.\n\n"
    "Example for students: The 2017 WannaCry ransomware exploited a buffer overflow in Windows SMB protocol (EternalBlue), propagating across networks and encrypting files globally. Understanding architecture helps us understand WHY such vulnerabilities are possible."
)

# Slide 6: Types of Computers
add_content_slide(prs,
    "Types of Computers & Their Security Posture",
    [
        "• Personal Computer (PC)",
        "   - Desktops, laptops — endpoint security focus",
        "",
        "• Server",
        "   - High-value targets; powers websites, databases, email",
        "   - Often runs Linux; needs hardening and monitoring",
        "",
        "• Embedded Systems",
        "   - ATMs, traffic lights, medical devices",
        "   - Often outdated and unpatchable → IoT vulnerability",
        "",
        "• Mobile & IoT Devices",
        "   - Smartphones, smart meters, security cameras",
        "   - Massive attack surface; default passwords common"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Different computer types require different security strategies. PCs are endpoints — they need antivirus, EDR (Endpoint Detection and Response), patch management, and user awareness training. In Zambia, most office workers at banks and government offices use Windows PCs.\n\n"
    "Servers are the crown jewels. A compromised web server can expose thousands of customer records. The 2017 Equifax breach (147 million records) happened because of an unpatched Apache Struts server. In Zambia, servers at mobile network operators and financial institutions run predominantly Linux, making Linux security skills directly employable.\n\n"
    "Embedded systems and IoT represent the fastest-growing attack surface. The 2016 Mirai botnet infected 600,000+ IoT devices (cameras, routers) with default passwords and launched record-breaking DDoS attacks. In Zambia, smart meters and cheap CCTV cameras often ship with default credentials like admin/admin. The first step in securing them is changing default passwords and isolating them on separate network segments.\n\n"
    "Class discussion: Ask students what IoT devices they have at home and whether they changed the default passwords."
)

# ======================= SECTION: OPERATING SYSTEMS =======================
add_section_slide(prs, "2. Operating Systems")

# Slide 7: What is an OS?
add_content_slide(prs,
    "What is an Operating System? The Security Layer",
    [
        "User",
        "  ↓",
        "Applications (Browser, Wireshark, Office)",
        "  ↓",
        "Operating System (Windows, Linux, macOS)",
        "  ↓",
        "Hardware (CPU, RAM, Disk, Network)",
        "",
        "Key OS functions with security impact:",
        "• Process management    → Malware hides as legitimate processes",
        "• Memory management     → Buffer overflow protection",
        "• File system           → Permissions and access control",
        "• Security & access     → Authentication and authorization",
        "• Device management     → Driver vulnerabilities"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "The OS is the ultimate security gatekeeper. Every application runs WITHIN the OS, and the OS decides what resources each application can access. If the OS is compromised, nothing running on it can be trusted. This is why 'rooting' (Android) or 'jailbreaking' (iOS) disables critical security protections.\n\n"
    "Process management: Attackers often name their malware processes to mimic legitimate Windows processes (e.g., 'svch0st.exe' instead of 'svchost.exe'). SOC analysts must be able to distinguish legitimate from malicious processes. Tools like Task Manager, Process Explorer, and `ps aux` are essential.\n\n"
    "Memory management: Modern OSs implement DEP/NX bits, ASLR, and stack canaries to prevent buffer overflow exploitation. When you hear that Windows or Linux has a 'security patch,' it often fixes a flaw in how the OS manages memory.\n\n"
    "File system security: The OS enforces who can read, write, or execute files. In the 2020 Twitter hack, attackers gained access to internal admin tools because of poorly configured access controls.\n\n"
    "Device management: Drivers run with kernel-level privileges. A vulnerability in a printer driver or graphics driver can give an attacker complete system control. The 2021 PrintNightmare vulnerability in Windows was a driver flaw that allowed remote code execution."
)

# Slide 8: Windows File System
add_content_slide(prs,
    "Windows File System: Where Threats Hide",
    [
        "C:\\",
        "├── Windows\\          ← OS system files (malware targets System32)",
        "├── Program Files\\    ← Installed applications",
        "├── Users\\",
        "│   └── Chanda\\",
        "│       ├── Desktop\\",
        "│       ├── Documents\\",
        "│       └── Downloads\\  ← Common malware delivery point",
        "└── Temp\\             ← Temporary files (classic malware hideout)",
        "",
        "Security-critical directories:",
        "• C:\\Windows\\System32  → Core OS files; common malware target",
        "• C:\\Users\\[name]\\AppData  → Hidden app data; persistence location",
        "• C:\\Temp  → Check during incident response"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Windows dominates the desktop market globally and in Zambia. Incident responders and SOC analysts spend a significant portion of their careers investigating Windows systems. Knowing the file system layout is essential for finding malware, analyzing logs, and securing systems.\n\n"
    "C:\\Windows\\System32: This directory contains critical DLLs and executables. Malware frequently injects malicious DLLs here or replaces legitimate files (DLL hijacking). The 2020 SolarWinds supply-chain attack compromised the Orion software, which had components in system directories, giving the attackers a foothold in thousands of organizations including US government agencies.\n\n"
    "AppData (Roaming/Local/LocalLow): This hidden directory is the #1 location for malware persistence. Ransomware like Ryuk and TrickBot store configuration files and payloads in AppData because it's writable by standard users and often excluded from antivirus scanning. To view it, enable 'Show hidden files' or use CMD: `dir /a C:\\Users\\[username]\\AppData`.\n\n"
    "Temp directories: Malware often extracts and executes from %TEMP% because any user can write here. During an investigation, listing recently modified files in Temp can reveal the initial infection vector. Command: `dir /o-d %TEMP%` shows most recent files first.\n\n"
    "Practical tip: Teach students to open File Explorer, enable 'View hidden items,' and navigate to these directories on their own computers."
)

# Slide 9: Windows CMD for Security
add_content_slide(prs,
    "Windows Command Prompt: Essential Security Commands",
    [
        "Command        Purpose                          Security Use",
        "─────────      ─────────────────────────────    ──────────────────────────────",
        "ipconfig       Show IP configuration            Find your IP, default gateway",
        "netstat -an    Show active connections          Spot suspicious connections",
        "tasklist       List running processes           Identify rogue processes",
        "systeminfo     System details                   Check patches and OS version",
        "net user       Manage user accounts             Audit accounts for backdoors",
        "dir /a         List hidden files                Find concealed malware",
        "",
        "Principle of Least Privilege:",
        "• Administrator = Full access (use sparingly!)",
        "• Standard User = Limited access (safer for daily use)",
        "• Guest = Very limited"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "The Windows Command Prompt (CMD) and PowerShell are the primary tools for Windows incident response and system administration. PowerShell has largely replaced CMD for advanced tasks, but these basic CMD commands remain essential and are guaranteed to work on every Windows system.\n\n"
    "`ipconfig /all` reveals not just IP addresses but also MAC addresses, DNS servers, and DHCP lease information. During an investigation, unexpected DNS servers might indicate DNS hijacking. Command: `ipconfig /displaydns` shows the DNS resolver cache, which can reveal connections to malicious domains.\n\n"
    "`netstat -an` lists all active network connections with their local and remote IP:port combinations. If you see a connection to an unusual foreign IP on port 4444 (default Metasploit port) or port 3389 (RDP) from an unknown process, it may indicate compromise. For more detail, use `netstat -anb` (requires admin) to see which executable owns each connection.\n\n"
    "`tasklist` and `tasklist /svc` show running processes and which Windows services they host. Malware often injects into legitimate processes like svchost.exe. Comparing `tasklist` output against a known-good baseline helps spot anomalies.\n\n"
    "`systeminfo` outputs the OS version, install date, hotfixes, and hardware info. The 'Hotfix(s)' line is critical — it shows which security patches are installed. Unpatched systems are vulnerable. For example, if KB5014699 (a critical patch) is missing, the system may be exploitable.\n\n"
    "`net user` lists local user accounts. Attackers often create hidden backdoor accounts with names like 'helpdesk' or 'backup.' The command `net localgroup administrators` shows who has admin rights. Any unexpected account here is a red flag."
)

# Slide 10: Linux File System
add_content_slide(prs,
    "Linux File System: The Security Professional's OS",
    [
        "/                       ← Root directory",
        "├── /home/              ← User home directories",
        "├── /etc/               ← Configuration files (passwords, services)",
        "├── /var/log/           ← System logs (security goldmine)",
        "├── /tmp/               ← Temporary files (malware favorite)",
        "├── /usr/               ← User programs and tools",
        "├── /bin/               ← Essential system commands",
    "└── /root/              ← Root (admin) home",
        "",
        "Why Linux for cybersecurity?",
        "• Free, open-source, highly configurable",
        "• 95%+ of cybersecurity tools are Linux-first",
        "• Powers most web servers, cloud platforms, and security appliances",
        "• Kali Linux is the standard OS for penetration testing"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Linux is the dominant operating system in cybersecurity, cloud computing, and server infrastructure. According to W3Techs, over 80% of web servers run Linux. In Zambia, MTN and Airtel's core network infrastructure, banking servers, and government systems heavily rely on Linux. A cybersecurity professional who cannot navigate Linux is severely limited.\n\n"
    "`/etc/` (pronounced 'et-see'): This directory contains system-wide configuration files. Key security files include:\n"
    "- `/etc/passwd` — user account information (readable by all, but passwords are stored in `/etc/shadow`)\n"
    "- `/etc/shadow` — hashed passwords (should be readable ONLY by root)\n"
    "- `/etc/hosts` — static hostname-to-IP mappings (attackers sometimes add malicious entries here to redirect traffic)\n"
    "- `/etc/crontab` — scheduled tasks (persistence mechanism for attackers)\n\n"
    "`/var/log/` — The security goldmine. Critical log files include:\n"
    "- `/var/log/auth.log` (Ubuntu/Debian) or `/var/log/secure` (RHEL/CentOS) — authentication attempts, including failed logins and sudo usage\n"
    "- `/var/log/syslog` — general system messages\n"
    "- `/var/log/apache2/access.log` — web server access logs\n"
    "- `/var/log/audit/audit.log` — SELinux audit logs\n\n"
    "`/tmp/` — Like Windows Temp, this is world-writable and a favorite malware drop location. Many Linux distributions now mount `/tmp` with `noexec` flag to prevent execution of dropped binaries.\n\n"
    "Zambian context: ZICTA (Zambia Information and Communications Technology Authority) and major banks run Linux-based infrastructure. Linux skills are directly employable in these organizations."
)

# Slide 11: Linux Commands & Permissions
add_content_slide(prs,
    "Linux Commands & File Permissions",
    [
        "Essential Linux Commands:",
        "• ls -la /etc       → List files with details",
        "• cat /etc/passwd   → View file contents",
        "• grep 'Failed' /var/log/auth.log  → Search logs for failed logins",
        "• ps aux            → Show all running processes",
        "• sudo chmod 600 file.txt  → Restrict file to owner only",
        "• whoami            → Show current user",
        "",
        "Permission Breakdown:  -rwxr-xr--",
        "• - = file type        rwx = Owner (read/write/execute)",
        "• r-x = Group          r-- = Others",
        "",
        "Numeric values: r=4, w=2, x=1",
        "Example: chmod 755 script.sh"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Linux command-line proficiency is non-negotiable for cybersecurity careers. The commands listed here are used daily by SOC analysts, penetration testers, and system administrators.\n\n"
    "`ls -la` reveals hidden files (those starting with a dot). Attackers often hide configuration files, SSH keys, and backdoors as hidden files. For example, `.bashrc` modifications can establish persistence — every time a user opens a terminal, malicious code in `.bashrc` executes.\n\n"
    "`grep` is one of the most powerful tools for log analysis. Example: `grep 'Failed password' /var/log/auth.log | wc -l` counts failed login attempts — a brute-force indicator. More advanced: `grep -i 'error\\|warning' /var/log/syslog` finds both errors and warnings. SOC analysts spend hours grepping through logs.\n\n"
    "File permissions are the foundation of Linux security. The permission string `-rwxr-xr--` breaks down as:\n"
    "- First character: file type (- = regular file, d = directory, l = symbolic link)\n"
    "- Characters 2-4: owner permissions\n"
    "- Characters 5-7: group permissions\n"
    "- Characters 8-10: others permissions\n\n"
    "Security-critical permission settings:\n"
    "- `chmod 600 ~/.ssh/id_rsa` — SSH private key should NEVER be readable by others. If it's 644, anyone on the system can steal it and impersonate you.\n"
    "- `chmod 700 /home/username` — Home directories should be private.\n"
    "- `chmod 644 /etc/passwd` — Should be world-readable (by design).\n"
    "- `chmod 640 /etc/shadow` — Should be readable ONLY by root and shadow group.\n\n"
    "Privilege escalation often starts with misconfigured permissions. If a world-writable script runs as root (via SUID), any user can modify it to execute arbitrary commands as root. Tools like `find / -perm -4000 2>/dev/null` discover all SUID files — a common privilege escalation reconnaissance step."
)

# ======================= SECTION: PROGRAMMING LOGIC =======================
add_section_slide(prs, "3. Introduction to Programming Logic")

# Slide 12: Why Programming Logic
add_content_slide(prs,
    "Why Programming Logic Matters in Cybersecurity",
    [
        "You don't need to be a software developer, but logic skills enable you to:",
        "",
        "• Read and write basic scripts for automation",
        "   → Automate log parsing, network scans, report generation",
        "",
        "• Understand how malware and exploits work",
        "   → Malware is just code written with malicious intent",
        "",
        "• Analyze code for vulnerabilities",
        "   → Spot SQL injection, buffer overflow, logic flaws",
        "",
        "• Write simple security tools",
        "   → Python is the #1 language in cybersecurity",
        "",
        "• Over 90% of tools in Kali Linux are written in Python or Bash"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Programming is the closest thing to a superpower in cybersecurity. It allows you to automate repetitive tasks, customize tools, and deeply understand how software works — which is essential because all cyberattacks target software in some form.\n\n"
    "Automation examples from real security work:\n"
    "- A SOC analyst might write a Python script to parse 100,000 firewall logs and flag connections to known malicious IPs.\n"
    "- A penetration tester might write a Bash script to run Nmap against a list of targets, save results, and generate a report.\n"
    "- A malware analyst uses Python to unpack and analyze suspicious files automatically.\n\n"
    "Understanding malware: At its core, malware is just software. Ransomware uses loops to iterate through files, conditionals to check file extensions, and file system APIs to encrypt data. If you understand programming, you can read malware reports, understand indicators of compromise (IoCs), and even reverse-engineer simple samples.\n\n"
    "Vulnerability analysis: Many vulnerabilities are logic errors. A missing input validation check (an `if` statement) leads to SQL injection. An unchecked buffer size leads to buffer overflow. Understanding control structures helps you 'think like an attacker' and spot these flaws in code reviews.\n\n"
    "Industry data: According to multiple cybersecurity job surveys, Python is the most requested programming skill, followed by Bash/PowerShell scripting. Learning these basics directly improves employability for roles like SOC Analyst, Junior Penetration Tester, and Security Analyst."
)

# Slide 13: Variables and Data Types
add_content_slide(prs,
    "Variables, Data Types & Security Examples",
    [
        "A variable is a named container that stores a value.",
        "",
        "Python Example:",
        "  student_name = 'Chanda Mwale'    # String (text)",
        "  age = 22                          # Integer (whole number)",
        "  gpa = 3.5                         # Float (decimal)",
        "  is_enrolled = True                # Boolean (True/False)",
        "",
        "Security-relevant data types:",
        "• String  → usernames, passwords, IP addresses, domain names",
        "• Integer → port numbers (22, 80, 443), attempt counters",
        "• Boolean → access_granted, is_admin, account_locked",
        "• List    → ['192.168.1.1', '10.0.0.1'] — list of target IPs"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Variables are the building blocks of every program, including every piece of malware and every security tool. Understanding how data is stored and typed helps security professionals work with logs, scripts, and forensic artifacts.\n\n"
    "String security implications:\n"
    "- Passwords stored as plain strings in memory can sometimes be recovered using memory dumping tools.\n"
    "- IP addresses and domain names in strings are constantly parsed in log analysis. Regex patterns are used to extract them.\n"
    "- In web application security, user input that arrives as strings must be validated before use in database queries to prevent SQL injection.\n\n"
    "Integer security implications:\n"
    "- Port numbers are integers in the range 0-65535. Well-known ports (0-1023) require administrative privileges to bind on Linux/Unix.\n"
    "- Integer overflow is a serious vulnerability class. When an integer exceeds its maximum value, it can wrap around to negative numbers, causing unexpected behavior. The 2014 OpenSSL 'Heartbleed' vulnerability and many buffer overflows involve integer-related flaws.\n\n"
    "Boolean security implications:\n"
    "- Authentication systems use booleans like `is_authenticated`, `is_admin`, and `account_active`. Logic flaws where these booleans are improperly checked can lead to unauthorized access. For example, if a web app checks `if (is_admin)` but the variable is null/undefined, JavaScript might treat it as false — but some frameworks behave differently.\n\n"
    "List/Array security implications:\n"
    "- Lists of IP addresses are used in firewall rules, blocklists, and scanning targets.\n"
    "- In secure coding, iterating over a list of user inputs without proper validation can lead to injection attacks."
)

# Slide 14: Control Structures
add_content_slide(prs,
    "Control Structures: If/Else & Loops in Security",
    [
        "If/Else — Decision Making:",
        "  if password == 'Secure@2026':",
        "      print('Access granted')",
        "  else:",
        "      print('Access denied')",
        "",
        "While Loop — Account Lockout Policy:",
        "  attempts = 0",
        "  while attempts < 3:",
        "      password = input('Enter password: ')",
        "      if password == 'Secure@2026':",
        "          print('Access granted'); break",
        "      else:",
        "          attempts += 1",
        "  if attempts == 3:",
        "      print('Account locked. Too many failed attempts.')"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Control structures are where security logic lives. Every access control decision, every rate limit, every account lockout policy is implemented using conditionals and loops.\n\n"
    "If/Else in security:\n"
    "- Authentication: `if (password_hash == stored_hash)` — this is the core of every login system.\n"
    "- Authorization: `if (user.role == 'admin')` — determines whether a user can access admin functions.\n"
    "- Input validation: `if (input_length > 1000)` — prevents buffer overflow by rejecting oversized input.\n"
    "- A real-world bug: In 2019, Apple FaceTime had a logic flaw where the call connected before the callee answered, because an `if` statement checked the wrong condition. Logic bugs can be just as dangerous as memory corruption.\n\n"
    "Loops in security:\n"
    "- Account lockout: The example on this slide shows exactly how real systems implement brute-force protection. After N failed attempts, the account is locked.\n"
    "- Port scanners like Nmap use loops to test thousands of ports: `for port in range(1, 65536): scan(port)`\n"
    "- Password crackers use loops to try dictionary words: `for word in wordlist: try_login(word)`\n"
    "- Log parsers use loops to process millions of lines: `for line in logfile: analyze(line)`\n\n"
    "Loop security bugs: Infinite loops can cause Denial of Service (DoS). If an attacker can trigger an infinite loop in a web application (e.g., by providing specially crafted input), they can exhaust CPU resources and crash the service."
)

# Slide 15: Functions, Pseudocode & Bash
add_content_slide(prs,
    "Functions, Pseudocode & Bash Scripting",
    [
        "Function — Reusable security check:",
        "  def check_password_strength(pw):",
        "      if len(pw) < 8: return 'Weak - too short'",
        "      elif pw.isalpha(): return 'Weak - letters only'",
        "      else: return 'Strong'",
        "",
        "Pseudocode — Planning before coding:",
        "  START → INPUT username → INPUT password",
        "  IF username EXISTS AND password MATCHES hash:",
        "      GRANT access → LOG success",
        "  ELSE:",
        "      INCREMENT failed_attempts",
        "      IF failed_attempts >= 3: LOCK account",
        "",
        "Bash Script — Network health check:",
        "  #!/bin/bash",
        "  ping -c 4 8.8.8.8",
        "  ip a | grep 'inet '"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Functions are essential for writing maintainable, secure code. They allow security checks to be defined once and reused everywhere.\n\n"
    "The password strength function example mirrors real password policy enforcement. Organizations typically require:\n"
    "- Minimum length (NIST recommends 8+ characters)\n"
    "- Complexity (uppercase, lowercase, numbers, symbols)\n"
    "- No dictionary words\n"
    "- No personal information\n\n"
    "However, modern guidance (NIST SP 800-63B) emphasizes LENGTH over complexity because users tend to reuse complex passwords or write them down. A passphrase like 'Correct-Horse-Battery-Staple!' is more secure and memorable than 'P@ssw0rd1'.\n\n"
    "Pseudocode is invaluable for security planning. Before writing a single line of code, security architects pseudocode authentication flows, incident response playbooks, and access control policies. It removes syntax distractions and focuses on logic.\n"
    "Example: When designing a multi-factor authentication system, you might pseudocode:\n"
    "`IF password_correct AND (otp_valid OR biometric_match) THEN grant_access ELSE deny_access`\n\n"
    "Bash scripting in security operations:\n"
    "- Automated backup scripts ensure critical logs are preserved.\n"
    "- Network monitoring scripts ping critical infrastructure and alert when services are down.\n"
    "- The example script combines `ping` (connectivity) and `ip a` (local network info). SOC analysts in Zambia could adapt this to monitor connections to local MTN/Airtel DNS servers or internal bank gateways.\n\n"
    "Important Bash security note: Always validate inputs in scripts. A script that takes user input and passes it directly to a command without sanitization is vulnerable to command injection — one of the OWASP Top 10 vulnerabilities."
)

# ======================= SECTION: MATHEMATICS =======================
add_section_slide(prs, "4. Mathematics for Cybersecurity")

# Slide 16: Number Systems
add_content_slide(prs,
    "Number Systems: Decimal, Binary, Hexadecimal",
    [
        "Decimal (Base 10) — Human system",
        "   Digits: 0-9     Example: 255",
        "",
        "Binary (Base 2) — Computer system",
        "   Digits: 0, 1    Example: 1111 1111",
        "   1 bit = one digit    8 bits = 1 byte",
        "",
        "Hexadecimal (Base 16) — Compact notation",
        "   Digits: 0-9, A=10, B=11, C=12, D=13, E=14, F=15",
        "   Example: FF",
        "",
        "Cybersecurity uses:",
        "• Binary  → IP addressing, subnet masks, Boolean logic",
        "• Hex     → MAC addresses, memory addresses, color codes, hashes",
        "• Decimal → Human-readable port numbers, IP octets"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Number systems are the language of computers and therefore the language of cybersecurity professionals. Every IP address, MAC address, cryptographic hash, and memory pointer is expressed in these systems.\n\n"
    "Decimal (Base 10): This is what humans naturally understand. In cybersecurity, we use decimal for port numbers (22, 80, 443), IP octets (192, 168, 1, 1), and version numbers. However, computers do not process decimal directly — everything is converted to binary internally.\n\n"
    "Binary (Base 2): All digital data is ultimately binary. Understanding binary is essential for:\n"
    "- IP subnetting: Subnet masks like 255.255.255.0 are /24 in CIDR notation, meaning 24 binary 1s followed by 8 binary 0s.\n"
    "- Boolean logic: AND, OR, NOT, XOR operate on individual bits.\n"
    "- Forensics: File headers (magic numbers) are often examined in binary/hex to determine file types regardless of extensions.\n\n"
    "Hexadecimal (Base 16): Hex is the 'shorthand' for binary. Since each hex digit represents exactly 4 bits, it is much more compact and readable than long binary strings.\n"
    "Cybersecurity applications of hex:\n"
    "- MAC addresses: `00:1A:2B:3C:4D:5E` — 6 bytes in hex.\n"
    "- Cryptographic hashes: MD5 and SHA-256 outputs are displayed in hex.\n"
    "- Memory addresses: When debugging exploits, memory locations are shown in hex.\n"
    "- Web colors: `#FF5733` is hex RGB notation.\n\n"
    "Example connection: A SHA-256 hash of a file looks like `e3b0c44298fc1c149afbf4c8996fb92427ae41e4649b934ca495991b7852b855`. This 64-character hex string represents 256 bits of data — the foundation of file integrity checking in incident response."
)

# Slide 17: Binary Conversion
add_content_slide(prs,
    "Binary Conversion with Examples",
    [
        "Decimal to Binary: Divide by 2, collect remainders",
        "   Convert 45 to binary:",
        "   45 ÷ 2 = 22 r 1",
        "   22 ÷ 2 = 11 r 0",
        "   11 ÷ 2 =  5 r 1",
        "    5 ÷ 2 =  2 r 1",
        "    2 ÷ 2 =  1 r 0",
        "    1 ÷ 2 =  0 r 1",
        "   Read bottom-up: 101101",
        "",
        "Binary to Decimal: Multiply by position values",
        "   128  64  32  16   8   4   2   1",
        "     1   0   1   0   1   1   0   0",
        "   = 128 + 32 + 8 + 4 = 172",
        "",
        "This octet (172) appears in IP addresses like 172.16.0.1"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Binary conversion is not just a math exercise — it is a practical skill used constantly in network engineering, forensics, and security analysis.\n\n"
    "Decimal to binary: The division-by-2 method works for any positive integer. In network security, this skill is needed when converting subnet masks or analyzing binary representations of permissions.\n\n"
    "Binary to decimal: The position-value method (powers of 2) is the fastest way to convert IP address octets. Each octet in an IPv4 address is an 8-bit binary number.\n"
    "Example: The IP address `192.168.1.1` breaks down as:\n"
    "- 192 = 11000000\n"
    "- 168 = 10101000\n"
    "- 1   = 00000001\n"
    "- 1   = 00000001\n\n"
    "Subnetting application: To determine if two IPs are on the same network, you perform a bitwise AND between each IP and the subnet mask. This requires understanding binary.\n"
    "Example: Is 192.168.1.100 on the same network as 192.168.1.200 with mask 255.255.255.0?\n"
    "- 192.168.1.100 AND 255.255.255.0 = 192.168.1.0\n"
    "- 192.168.1.200 AND 255.255.255.0 = 192.168.1.0\n"
    "Yes, same network.\n\n"
    "Forensics application: File signatures (magic numbers) are often expressed in hex, which directly maps to binary. For example, a PNG file starts with the hex signature `89 50 4E 47` (`89504E47` in continuous hex). In binary, `89` = `10001001`. Forensic tools examine these binary patterns to identify files even when extensions have been changed to hide them.\n\n"
    "Teaching tip: Use a whiteboard to manually work through conversions. Have students practice with numbers relevant to them — their birth year, their phone number digits, or common IP octets like 192, 168, and 255."
)

# Slide 18: Hexadecimal Conversion
add_content_slide(prs,
    "Hexadecimal Conversion: The Binary Shortcut",
    [
        "Binary to Hex: Group into 4 bits",
        "   Binary:  1010  1100",
        "   Hex:        A     C",
        "   So 10101100 binary = AC hex = 172 decimal",
        "",
        "Decimal to Hex: Divide by 16",
        "   Convert 255 to hex:",
        "   255 ÷ 16 = 15 remainder 15 → F",
        "    15 ÷ 16 =  0 remainder 15 → F",
        "   Result: FF",
        "",
        "Quick Reference:",
        "   Dec  0   5   9   10   15   16   255",
        "   Bin  0000 0101 1001 1010 1111 00010000 11111111",
        "   Hex  0    5    9    A    F    10       FF"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "The binary-to-hex 'group into 4 bits' method is the most important conversion technique for cybersecurity professionals. It is fast, reliable, and used constantly in the field.\n\n"
    "Why hex is everywhere in security:\n"
    "1. MAC addresses: Every network card has a unique 48-bit MAC address, displayed as 6 pairs of hex digits. Example: `00:50:56:C0:00:08`. The first 3 bytes (00:50:56) identify the manufacturer (VMware, in this case). Attackers can use MAC vendor prefixes to identify what types of devices are on a network.\n\n"
    "2. Cryptographic hashes: When you download security software, the vendor provides a SHA-256 hash so you can verify file integrity. These hashes are long strings of hex characters. If even one bit changes in the file, the hash changes completely.\n\n"
    "3. Memory forensics: Memory dumps and debugger outputs display addresses in hex. When analyzing malware, you might see instructions like `MOV EAX, 0x00401234` — the `0x` prefix indicates hexadecimal.\n\n"
    "4. Color coding in web security: Web developers use hex colors like `#FF0000` (pure red). While not directly a security topic, understanding hex notation helps security professionals review web application code and CSS during assessments.\n\n"
    "Practical exercise: Have students convert their favorite color from hex to binary. For example, royal blue `#4169E1` breaks down as:\n"
    "- 41 = 01000001\n"
    "- 69 = 01101001\n"
    "- E1 = 11100001\n\n"
    "This demonstrates that hex is just a compact way to write binary — and everything in computing is ultimately binary."
)

# Slide 19: Boolean Algebra
add_content_slide(prs,
    "Boolean Algebra & Logic Gates",
    [
        "Boolean algebra uses only TRUE (1) and FALSE (0).",
        "",
        "Logic Operations:",
        "• AND (∧)  → Both must be true    1 AND 1 = 1",
        "• OR (∨)   → At least one true    0 OR 1 = 1",
        "• NOT (¬)  → Inverts the value    NOT 1 = 0",
        "• XOR (⊕)  → One or other, not both    1 XOR 1 = 0",
        "",
        "Truth Table — AND Gate:",
        "   A | B | A AND B",
        "   0 | 0 |    0",
        "   0 | 1 |    0",
        "   1 | 0 |    0",
        "   1 | 1 |    1",
        "",
        "Truth Table — OR Gate:",
        "   A | B | A OR B",
        "   0 | 0 |   0",
        "   0 | 1 |   1",
        "   1 | 0 |   1",
        "   1 | 1 |   1"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Boolean logic is the mathematical foundation of digital computing. Every CPU instruction, every firewall rule, every access control decision is ultimately evaluated using Boolean operations.\n\n"
    "AND in security:\n"
    "- Multi-factor authentication (MFA) is essentially an AND operation: `access_granted = password_correct AND otp_valid`. Both must be true.\n"
    "- Firewall rules often use AND logic: `allow_traffic = (source_ip_in_whitelist) AND (destination_port == 443)`.\n\n"
    "OR in security:\n"
    "- Role-based access control might use OR: `can_view_report = (role == 'admin') OR (role == 'analyst')`. Either condition satisfies access.\n"
    "- Alert conditions in SIEM systems: `trigger_alert = (failed_logins > 5) OR (malware_detected == true)`.\n\n"
    "NOT in security:\n"
    "- Deny lists use NOT logic: `allow_access = NOT(ip_in_blocklist)`.\n"
    "- Network segmentation: `permitted = NOT(traffic_from_untrusted_zone_to_critical_asset)`.\n\n"
    "Logic gates in hardware: CPUs are built from billions of transistors arranged as logic gates. Understanding AND, OR, NOT, and XOR helps students appreciate how the physical hardware implements the software they use. While SOC analysts don't design CPUs, this knowledge builds the deep technical foundation that distinguishes junior technicians from senior engineers.\n\n"
    "Set theory connection: Set operations (union, intersection, complement) map directly to Boolean operations. In threat intelligence, analysts work with sets of indicators of compromise (IoCs). Finding the intersection between two threat feeds (common IPs) is an AND operation. Combining two feeds is a union (OR operation)."
)

# Slide 20: XOR and Encryption
add_content_slide(prs,
    "XOR: The Foundation of Encryption",
    [
        "XOR Truth Table:",
        "   A | B | A XOR B",
        "   0 | 0 |    0",
        "   0 | 1 |    1",
        "   1 | 0 |    1",
        "   1 | 1 |    0",
        "",
        "The Magic Property:",
        "   If you XOR a value with a key, then XOR the result",
        "   with the SAME key, you get the original value back!",
        "",
        "Example:",
        "   Plaintext:  1010",
        "   Key:        1100",
        "   Ciphertext: 0110   (1010 XOR 1100)",
        "   Decrypt:    1010   (0110 XOR 1100) ← Original!",
        "",
        "This is the basis of stream ciphers and many encryption algorithms."
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "XOR (Exclusive OR) is arguably the single most important Boolean operation in cryptography. Its unique mathematical properties make it ideal for encryption.\n\n"
    "The reversible property of XOR:\n"
    "- `A XOR B = C`\n"
    "- `C XOR B = A`\n"
    "- `C XOR A = B`\n"
    "This means if you XOR plaintext with a secret key to get ciphertext, XORing the ciphertext with the same key perfectly recovers the plaintext. No information is lost.\n\n"
    "Real-world applications of XOR in cryptography:\n"
    "1. One-Time Pad (OTP): The only encryption scheme proven to be unconditionally secure. It XORs plaintext with a truly random key that is as long as the message and never reused. During the Cold War, spies used OTPs to send unbreakable messages.\n\n"
    "2. Stream Ciphers: Modern protocols like RC4 (now deprecated) and ChaCha20 use XOR to combine a keystream with plaintext. In TLS/HTTPS, ChaCha20-Poly1305 is widely used to encrypt web traffic.\n\n"
    "3. Block Ciphers: AES (Advanced Encryption Standard) uses XOR extensively in its rounds to mix data with round keys. AES protects everything from WhatsApp messages to government classified documents.\n\n"
    "4. Hash Functions: SHA-256 and other hash algorithms use XOR-like operations to create avalanche effects — changing one input bit drastically changes the output.\n\n"
    "Malware obfuscation: Attackers also use XOR! Simple malware samples XOR their payload with a one-byte key to evade signature-based antivirus detection. Security analysts use tools that brute-force XOR keys to decode hidden strings.\n\n"
    "Python demonstration for class:\n"
    "```python\n"
    "plaintext = 'A'  # ASCII 65 = 01000001\n"
    "key = 0x2A       # 00101010\n"
    "cipher = ord(plaintext) ^ key  # 01101011 (107)\n"
    "decrypt = cipher ^ key         # 01000001 (65) = 'A'\n"
    "```\n\n"
    "This simple demo proves that XOR encryption and decryption use the exact same operation — a beautiful symmetry that underpins modern cryptography."
)

# ======================= SECTION: NETWORKING =======================
add_section_slide(prs, "5. Networking Essentials")

# Slide 21: What is a Network?
add_content_slide(prs,
    "What is a Network? Types & Security Relevance",
    [
        "A network is two or more computers connected to share resources and communicate.",
        "Virtually all cybersecurity work involves networks.",
        "",
        "Network Types:",
        "• LAN  (Local Area Network)      → School lab, office building",
        "• WAN  (Wide Area Network)       → National backbone (ZAMTEL)",
        "• MAN  (Metropolitan Area Network) → City-wide network (Lusaka)",
        "• WLAN (Wireless LAN)            → Home WiFi, coffee shop hotspots",
        "• VPN  (Virtual Private Network)  → Encrypted tunnel for remote work",
        "",
        "Security perspective:",
        "• LANs need segmentation to limit breach spread",
        "• WANs expose traffic to interception — encryption is essential",
        "• WLANs are especially vulnerable to eavesdropping",
        "• VPNs protect remote workers but misconfigurations are common"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Networking is the backbone of cybersecurity. Every attack, defense mechanism, and investigation relies on how devices communicate. From securing corporate infrastructures in Lusaka to analyzing suspicious traffic in a SOC, strong networking fundamentals are essential.\n\n"
    "LAN Security:\n"
    "- Local Area Networks are where most day-to-day business happens. However, a flat LAN (where all devices can communicate freely) allows malware like ransomware to spread rapidly. The WannaCry ransomware (2017) spread laterally across LANs using SMB protocol vulnerabilities. Network segmentation — dividing a LAN into smaller subnets — contains such outbreaks.\n\n"
    "WAN Security:\n"
    "- Wide Area Networks span large geographic areas and typically use leased lines or MPLS. Because traffic traverses third-party infrastructure, encryption (IPsec VPNs) is mandatory for sensitive data. In Zambia, the ZAMTEL national backbone connects cities and government offices. Compromises at WAN level can have national impact.\n\n"
    "WLAN Security:\n"
    "- Wireless networks broadcast data over radio waves, making them easy to intercept. Tools like Aircrack-ng can capture WiFi handshakes and crack weak passwords. Public WiFi in Lusaka cafes is particularly risky because attackers can set up rogue access points or perform man-in-the-middle attacks. Always use WPA3 (or at least WPA2) with strong passwords. Never use WEP — it can be cracked in minutes.\n\n"
    "VPN Security:\n"
    "- Virtual Private Networks create encrypted tunnels over untrusted networks (like the internet). They are essential for remote work. However, VPN misconfigurations are a leading cause of breaches. The 2021 Colonial Pipeline attack began through a compromised VPN account with no multi-factor authentication.\n\n"
    "Discussion question: Ask students which network types they use daily and what security measures they currently take."
)

# Slide 22: OSI Model
add_content_slide(prs,
    "The OSI Model: 7 Layers of Communication",
    [
        "SENDER                          RECEIVER",
        "  │                                │",
        "7 │ Application  ←────────────→  Application  │ 7  (HTTP, DNS, FTP)",
        "6 │ Presentation ←────────────→  Presentation │ 6  (SSL/TLS, encoding)",
        "5 │ Session      ←────────────→  Session      │ 5  (NetBIOS, RPC)",
        "4 │ Transport    ←────────────→  Transport    │ 4  (TCP, UDP)",
        "3 │ Network      ←────────────→  Network      │ 3  (IP, routing)",
        "2 │ Data Link    ←────────────→  Data Link    │ 2  (Ethernet, MAC)",
        "1 │ Physical     ──────────────  Physical     │ 1  (Cables, WiFi)",
        "",
        "Memory Aid: All People Seem To Need Data Processing",
        "",
        "Cybersecurity relevance: Attacks happen at specific layers."
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "The OSI model is the conceptual framework that every cybersecurity professional must internalize. It describes how data travels from one application to another across a network. While the TCP/IP model is what devices actually implement, OSI provides the language security professionals use to describe where attacks and defenses occur.\n\n"
    "Layer-by-layer security relevance:\n\n"
    "Layer 7 — Application: Where users interact. Attacks include SQL injection, cross-site scripting (XSS), malware downloads, and phishing. Defenses include web application firewalls (WAFs), input validation, and secure coding.\n\n"
    "Layer 6 — Presentation: Handles data formatting, compression, and encryption. TLS/SSL encryption operates here conceptually. Attacks include weak cipher suites and certificate spoofing.\n\n"
    "Layer 5 — Session: Manages connections between applications. Session hijacking attacks steal valid session tokens (like cookies) to impersonate users.\n\n"
    "Layer 4 — Transport: Responsible for end-to-end delivery, port numbers, and connection management. SYN flood attacks exhaust server resources by initiating but never completing TCP handshakes. Port scanning (Nmap) targets this layer.\n\n"
    "Layer 3 — Network: Handles IP addressing and routing between networks. IP spoofing, DDoS amplification, and route hijacking occur here. Firewalls operate at this layer.\n\n"
    "Layer 2 — Data Link: Manages device-to-device communication on the same network using MAC addresses. ARP spoofing and MAC flooding attacks target this layer. Switches operate here.\n\n"
    "Layer 1 — Physical: The actual transmission medium — cables, radio waves, fiber optics. Physical security (preventing unauthorized access to equipment) and signal interception (TEMPEST attacks) are concerns here.\n\n"
    "Teaching analogy: Use the Zampost letter analogy from the lesson plan. Application = what you write. Presentation = the language. Session = the ongoing correspondence. Transport = registered mail vs flyer. Network = routing from Lusaka to Kitwe. Data Link = the truck route. Physical = the roads."
)

# Slide 23: OSI Model & Cyber Attacks
add_content_slide(prs,
    "Mapping Cyber Attacks to OSI Layers",
    [
        "Layer    Attack Type                    Defense",
        "─────    ───────────────────────────    ──────────────────────────────",
        "L7       SQL Injection, Phishing, XSS   WAF, Secure Coding, User Awareness",
        "L6       Weak TLS/SSL, Certificate Fraud  Patch Management, Certificate Pinning",
        "L5       Session Hijacking                HTTPS, Secure Cookies, Short TTLs",
        "L4       SYN Flood, Port Scanning         Firewalls, IDS/IPS, Rate Limiting",
        "L3       IP Spoofing, DDoS, Routing Hijack  ACLs, DDoS Mitigation, BGP Security",
        "L2       ARP Spoofing, MAC Flooding       Dynamic ARP Inspection, Port Security",
        "L1       Wiretapping, Physical Theft      Encryption, Locks, Biometrics",
        "",
        "Key insight: Defense in Depth means protecting EVERY layer."
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "This slide is critical for helping students connect abstract theory to real attacks. Every major cyberattack can be mapped to an OSI layer, and effective defense requires protection at multiple layers simultaneously — this is the 'Defense in Depth' strategy.\n\n"
    "Real-world attack examples mapped to layers:\n\n"
    "L7 — SQL Injection: The 2017 Equifax breach exposed 147 million records partly due to unpatched web application vulnerabilities. SQL injection has been in the OWASP Top 10 for over 15 years because developers fail to validate user input at the application layer.\n\n"
    "L4 — SYN Flood: This classic DDoS attack exploits the TCP three-way handshake. The attacker sends SYN packets but never completes the handshake, exhausting server connection tables. The 2000 attack on Yahoo!, eBay, and CNN used SYN floods. Modern mitigation includes SYN cookies and rate limiting.\n\n"
    "L3 — DDoS Amplification: Attackers send small requests to open DNS or NTP servers with a spoofed victim IP. The servers respond with much larger packets, overwhelming the victim. The 2016 Dyn attack (which took down Twitter, Netflix, and Reddit) used IoT botnets and DNS amplification.\n\n"
    "L2 — ARP Spoofing: On a local network, an attacker sends fake ARP messages to associate their MAC address with the IP address of the default gateway. This enables man-in-the-middle attacks where all traffic passes through the attacker's machine. Tools like Ettercap automate this. Defense: Dynamic ARP Inspection (DAI) on managed switches.\n\n"
    "L1 — Physical attacks: The 2013 Target breach began when attackers stole HVAC vendor credentials, but the physical access to install malware on POS terminals was also a factor. In Zambia, protecting server rooms with biometric access, CCTV, and mantraps is essential for critical infrastructure.\n\n"
    "Career connection: SOC analysts investigate alerts that often specify the layer. A Layer 7 WAF alert is very different from a Layer 3 DDoS alert. Understanding the layer helps analysts respond appropriately."
)

# Slide 24: TCP/IP Model
add_content_slide(prs,
    "The TCP/IP Model: Real-World Implementation",
    [
        "TCP/IP Layer          Corresponds to OSI        Protocols / Examples",
        "─────────────         ──────────────────        ───────────────────────────",
        "Application           Layers 5, 6, 7              HTTP, HTTPS, DNS, FTP, SSH",
        "Transport             Layer 4                     TCP, UDP",
        "Internet              Layer 3                     IP, ICMP, ARP",
        "Network Access        Layers 1, 2                 Ethernet, WiFi (802.11)",
        "",
        "OSI is theory → used for learning and explaining.",
        "TCP/IP is implementation → used by actual devices, routers, firewalls, and tools.",
        "",
        "In practice, cybersecurity professionals use BOTH models depending on context."
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "While the OSI model is excellent for learning and communication, the TCP/IP model is what routers, operating systems, firewalls, and packet analyzers like Wireshark actually implement. Every cybersecurity professional must be fluent in both.\n\n"
    "Application Layer (TCP/IP): Combines OSI layers 5-7. This is where most security tools focus because it's where users interact with systems. Key protocols:\n"
    "- HTTP (port 80): Unencrypted web traffic. Never use for sensitive data.\n"
    "- HTTPS (port 443): Encrypted with TLS. Required for banking, shopping, and any authentication.\n"
    "- DNS (port 53): The 'phonebook of the internet.' Vulnerable to poisoning, hijacking, and tunneling.\n"
    "- SSH (port 22): Encrypted remote administration. Replaces insecure Telnet.\n"
    "- FTP (port 21): File transfer without encryption. Use SFTP (port 22) instead.\n\n"
    "Transport Layer (TCP/IP): Equivalent to OSI Layer 4. This layer determines whether communication is reliable (TCP) or fast/unreliable (UDP). Port numbers live here. Firewalls make most of their decisions at the Transport and Internet layers.\n\n"
    "Internet Layer (TCP/IP): Equivalent to OSI Layer 3. IP addresses, routing, and ICMP (ping) operate here. This is where network-layer firewalls and intrusion detection systems examine packets. VPNs like IPsec operate primarily at this layer.\n\n"
    "Network Access Layer (TCP/IP): Equivalent to OSI Layers 1-2. Ethernet frames, MAC addresses, WiFi signals, and network interface cards operate here. This layer is often overlooked but is critical for local network attacks.\n\n"
    "Tool context: When using Wireshark, you see the TCP/IP model in action. Each packet shows Ethernet header (Network Access), IP header (Internet), TCP/UDP header (Transport), and HTTP/DNS payload (Application)."
)

# Slide 25: TCP vs UDP
add_content_slide(prs,
    "TCP vs UDP: The Security Trade-off",
    [
        "Feature           TCP                              UDP",
        "─────────         ─────────────────────────────    ──────────────────────────────",
        "Connection        Connection-oriented              Connectionless",
        "Reliability       Guaranteed delivery              No guarantee",
        "Speed             Slower (handshake overhead)      Faster",
        "Use Cases         Web, email, file transfer        Streaming, DNS, VoIP",
        "Cyber Risk        SYN flood attacks                DDoS amplification",
        "",
        "TCP Three-Way Handshake:",
        "   Client → SYN → Server",
        "   Client ← SYN-ACK ← Server",
        "   Client → ACK → Server",
        "",
        "Attackers exploit the handshake with SYN floods."
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Understanding the difference between TCP and UDP is fundamental for network security. These protocols determine how data is transmitted, and each has distinct attack vectors and defensive requirements.\n\n"
    "TCP (Transmission Control Protocol):\n"
    "- Connection-oriented: Before any data is sent, TCP performs a three-way handshake (SYN → SYN-ACK → ACK).\n"
    "- Reliable: TCP guarantees delivery through sequence numbers, acknowledgments, and retransmissions.\n"
    "- Used for: Web browsing (HTTP/HTTPS), email (SMTP, IMAP), file transfer (FTP, SFTP), remote access (SSH, RDP).\n"
    "- Security risk: SYN flood attacks exploit the handshake. An attacker sends thousands of SYN packets but never completes the handshake, filling up the server's connection queue and denying service to legitimate users. Mitigation includes SYN cookies, rate limiting, and SYN proxies.\n\n"
    "UDP (User Datagram Protocol):\n"
    "- Connectionless: No handshake. Packets are sent without establishing a connection.\n"
    "- Unreliable: No guarantee of delivery, no retransmissions, no sequencing.\n"
    "- Used for: DNS queries, video streaming (YouTube, Netflix), VoIP (WhatsApp calls), online gaming.\n"
    "- Security risk: UDP amplification attacks. Because UDP has no handshake, attackers can send small spoofed requests to servers (like DNS or NTP), which respond with much larger UDP packets to the victim's IP address. The 2016 Dyn attack used this technique. Mitigation includes rate limiting, source validation, and DDoS protection services.\n\n"
    "Analogy for teaching:\n"
    "- TCP = Registered mail with return receipt. You know it was delivered. Slower but reliable.\n"
    "- UDP = Dropping a letter in a mailbox. Fast, cheap, but no confirmation of delivery.\n\n"
    "Port scanning context: Nmap uses different scan types for TCP (-sS SYN scan, -sT connect scan) versus UDP (-sU). UDP scanning is much slower because there is no response for open ports — the scanner must wait for timeouts."
)

# Slide 26: IP Addressing & Subnetting
add_content_slide(prs,
    "IP Addressing & Subnetting Basics",
    [
        "IPv4 Address:",
        "   32-bit address written as 4 octets",
        "   Example: 192.168.1.100",
        "   Each octet = 8 bits = 0-255",
        "",
        "Private IP Ranges (not routable on internet):",
        "   10.0.0.0      – 10.255.255.255     (Large enterprises)",
        "   172.16.0.0    – 172.31.255.255     (Medium networks)",
        "   192.168.0.0   – 192.168.255.255    (Home/small office)",
        "   127.0.0.1                          (Loopback / localhost)",
        "",
        "Subnetting Example:",
        "   IP:       192.168.1.100",
        "   Mask:     255.255.255.0  (/24)",
        "   Network:  192.168.1.0",
        "   Hosts:    192.168.1.1 – 192.168.1.254"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "IP addressing is the addressing system of the internet. Every device that communicates over a network has an IP address. Cybersecurity professionals must be comfortable with IPv4 addressing because it is still the dominant protocol in most networks, even as IPv6 adoption grows.\n\n"
    "IPv4 basics:\n"
    "- 32-bit address space provides approximately 4.3 billion unique addresses.\n"
    "- Written in dotted-decimal notation: four octets separated by periods.\n"
    "- Each octet is 8 bits, ranging from 0 to 255 in decimal.\n"
    "- Example: `192.168.1.100` = `11000000.10101000.00000001.01100100` in binary.\n\n"
    "Private IP ranges (RFC 1918):\n"
    "These addresses are reserved for use inside private networks and are not routable on the public internet. Network Address Translation (NAT) translates private IPs to public IPs at the network boundary.\n"
    "- `10.0.0.0/8`: Used by very large organizations.\n"
    "- `172.16.0.0/12`: Used by medium-sized networks.\n"
    "- `192.168.0.0/16`: Used by home routers and small businesses.\n"
    "When you connect to MTN or Airtel mobile data in Zambia, your phone typically receives a private IP (like `192.168.x.x`) from the carrier's internal network, while a public IP is used for internet-facing traffic.\n\n"
    "Subnetting and security:\n"
    "Subnetting divides a network into smaller segments. From a security perspective, this is network segmentation — one of the most effective security controls.\n"
    "- If the finance department is on `192.168.10.0/24` and engineering is on `192.168.20.0/24`, a router or firewall between them can restrict access.\n"
    "- If ransomware infects one subnet, segmentation prevents it from spreading to others.\n"
    "- Critical assets (domain controllers, databases) should be on isolated subnets with strict access controls.\n\n"
    "CIDR notation: `/24` means the first 24 bits are the network portion. This is equivalent to subnet mask `255.255.255.0`. CIDR is the standard way network professionals communicate subnet sizes.\n\n"
    "IPv6 note: With 128-bit addresses, IPv6 provides virtually unlimited addresses. Format: `2001:0db8:85a3::8a2e:0370:7334`. While not yet dominant in Zambia, IPv6 skills are increasingly important globally."
)

# Slide 27: Key Protocols
add_content_slide(prs,
    "Key Protocols: DNS, DHCP, HTTP vs HTTPS",
    [
        "DNS (Domain Name System) — Port 53",
        "   Translates names to IP addresses",
        "   www.zamtel.co.zm  →  196.32.X.X",
        "   ⚠ Security threat: DNS poisoning redirects to fake sites",
        "",
        "DHCP (Dynamic Host Configuration Protocol) — Ports 67/68",
        "   Automatically assigns IP addresses",
        "   ⚠ Security threat: Rogue DHCP server hijacks traffic",
        "",
        "HTTP (Port 80)  →  Unencrypted — readable by anyone on the network",
        "HTTPS (Port 443) →  Encrypted with TLS — safe for passwords & banking",
        "",
        "SSH (Port 22) →  Encrypted remote access (replaces insecure Telnet)"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "These protocols are the workhorses of the internet, and each has significant security implications. Understanding them is essential for both defense and attack analysis.\n\n"
    "DNS (Domain Name System):\n"
    "DNS is often called 'the phonebook of the internet.' When you type www.zamtel.co.zm, your computer queries a DNS server to get the corresponding IP address. Without DNS, we'd have to memorize IP addresses.\n\n"
    "DNS security threats:\n"
    "- DNS Spoofing / Poisoning: An attacker corrupts the DNS cache to redirect users to malicious sites. For example, a user trying to reach their bank's website is instead sent to a phishing site that looks identical.\n"
    "- DNS Tunneling: Malware uses DNS queries to secretly exfiltrate data or receive commands. Since DNS is usually allowed through firewalls, it's an attractive covert channel.\n"
    "- DNS Hijacking: Attackers modify DNS settings on routers or devices to point to malicious DNS servers.\n\n"
    "Defenses: DNSSEC (DNS Security Extensions) adds cryptographic verification to DNS responses. DNS filtering services (like OpenDNS, Cloudflare 1.1.1.1) block known malicious domains. SOC analysts monitor DNS logs for unusual query volumes and suspicious domain names.\n\n"
    "DHCP (Dynamic Host Configuration Protocol):\n"
    "DHCP automatically assigns IP addresses, subnet masks, default gateways, and DNS servers to devices joining a network. It's what happens when you connect to WiFi and instantly get internet access.\n\n"
    "DHCP security threats:\n"
    "- Rogue DHCP Server: An attacker sets up a malicious DHCP server on the network. When devices connect, they receive attacker-controlled IP configurations — typically pointing to a rogue gateway that intercepts all traffic (man-in-the-middle attack).\n"
    "- DHCP Starvation: An attacker floods the DHCP server with requests, exhausting the IP address pool and causing a denial of service.\n\n"
    "Defenses: DHCP snooping on managed switches creates trusted ports for legitimate DHCP servers and drops DHCP offers from untrusted ports.\n\n"
    "HTTP vs HTTPS:\n"
    "This is one of the most important distinctions for end-user security. HTTP sends data in plaintext. Anyone on the same network (coffee shop WiFi, shared office network) can use packet sniffing tools like Wireshark to read usernames, passwords, and session cookies.\n"
    "HTTPS encrypts the connection using TLS (Transport Layer Security). The padlock icon in the browser indicates HTTPS.\n"
    "Real-world impact: In 2010, the Firesheep browser extension demonstrated how easy it was to hijack Facebook and Twitter sessions on public WiFi because those sites hadn't fully deployed HTTPS. Today, major sites use HTTPS by default.\n"
    "Best practice: Never enter passwords or payment information on HTTP sites. Always verify `https://` and the padlock icon.\n\n"
    "SSH (Secure Shell):\n"
    "SSH provides encrypted remote command-line access to servers. It replaced Telnet, which sent everything — including passwords — in plaintext. SSH is the primary way administrators manage Linux servers and network devices.\n"
    "Security considerations:\n"
    "- Brute force attacks constantly target SSH on port 22.\n"
    "- Best practices: Disable password authentication, use SSH key pairs, change the default port, implement fail2ban for rate limiting."
)

# Slide 28: Common Ports & Security Risks
add_content_slide(prs,
    "Common Port Numbers & Their Security Risks",
    [
        "Port    Protocol    Use                              Security Note",
        "────    ────────    ─────────────────────────────    ──────────────────────────────",
        "22      SSH         Secure remote access             Brute-force target — use keys",
        "23      Telnet      Insecure remote access           AVOID — plaintext passwords",
        "25      SMTP        Sending email                    Spam relay target",
        "53      DNS         Domain resolution                DNS tunneling, poisoning",
        "80      HTTP        Web (unencrypted)                Sniffing risk on public WiFi",
        "443     HTTPS       Web (encrypted)                  Standard for secure sites",
        "445     SMB         Windows file sharing             WannaCry, ransomware vector",
        "3389    RDP         Windows remote desktop           Major ransomware entry point",
        "3306    MySQL       Database                         Expose only to app servers",
        "",
        "Attackers scan for open ports to find ways in."
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "Port numbers are the 'doors' through which network services communicate. Knowing common ports is essential for firewall configuration, intrusion detection, and penetration testing. Attackers scan for open ports to identify running services and find exploitable vulnerabilities.\n\n"
    "Port 22 (SSH): The most common remote administration protocol. Because it's so widely used, it's also the most attacked port on the internet. Automated bots constantly attempt to brute-force SSH credentials.\n"
    "- Defense: Disable root login, use key-based authentication, implement fail2ban, consider non-standard port (security through obscurity — not enough alone, but reduces log noise).\n\n"
    "Port 23 (Telnet): An ancient protocol that transmits everything in plaintext. There is absolutely no reason to use Telnet in 2026. If you see port 23 open during a scan, it is a critical finding.\n\n"
    "Port 25 (SMTP): Used for sending email. Misconfigured SMTP servers can be exploited as open relays for spam. Email security protocols like SPF, DKIM, and DMARC protect against email spoofing.\n\n"
    "Port 53 (DNS): As discussed, DNS is critical but vulnerable to poisoning and tunneling. Monitor DNS traffic closely.\n\n"
    "Port 80 (HTTP): Unencrypted web traffic. While many sites redirect to HTTPS, port 80 must remain open for the redirect to work. However, no sensitive data should ever be transmitted over port 80.\n\n"
    "Port 443 (HTTPS): The standard for secure web communication. When students use Airtel Money, Zanaco Internet Banking, or mobile money apps, their data travels over port 443 with TLS encryption.\n\n"
    "Port 445 (SMB): Server Message Block is used for Windows file sharing. The EternalBlue exploit (MS17-010) targeted SMBv1 and was used by WannaCry and NotPetya to spread ransomware globally. Microsoft has strongly recommended disabling SMBv1.\n\n"
    "Port 3389 (RDP): Remote Desktop Protocol is a major attack vector for ransomware groups. In 2019, the FBI reported that RDP was the most common entry point for ransomware attacks. Best practice: Never expose RDP directly to the internet; use a VPN instead.\n\n"
    "Port 3306 (MySQL): Database ports should never be exposed to the internet. They should only be accessible from application servers within the same network segment. Exposed MySQL ports are frequently targeted by automated exploitation scripts.\n\n"
    "Career skill: Nmap port scanning is one of the first skills penetration testers learn. The command `nmap -sS -O target` performs a SYN scan and attempts OS fingerprinting. SOC analysts also use Nmap to verify their network perimeter and ensure no unexpected ports are open."
)

# ======================= CONCLUSION =======================
add_section_slide(prs, "Module 1 Review & What's Next")

# Slide 29: Revision Summary
add_content_slide(prs,
    "Module 1 Revision: Key Takeaways",
    [
        "Computer Fundamentals:",
        "   ✓ IPOS model — every stage is an attack surface",
        "   ✓ CPU, RAM, Storage, NIC — each has security implications",
        "",
        "Operating Systems:",
        "   ✓ Windows & Linux file systems and security directories",
        "   ✓ Command-line navigation (CMD and terminal)",
        "   ✓ Linux permissions (rwx / numeric notation)",
        "",
        "Programming Logic:",
        "   ✓ Variables, if/else, loops, functions",
        "   ✓ Pseudocode and Bash scripting for automation",
        "",
        "Mathematics:",
        "   ✓ Binary, decimal, hexadecimal conversions",
        "   ✓ Boolean algebra and XOR in cryptography",
        "",
        "Networking:",
        "   ✓ OSI & TCP/IP models, IP addressing, subnetting",
        "   ✓ DNS, DHCP, ports, and protocol security risks"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "This revision slide summarizes the foundational knowledge that students must master before advancing to Module 2. Emphasize that each bullet point is directly applicable to real security work.\n\n"
    "Connect the topics together:\n"
    "- A SOC analyst investigating a Windows malware alert uses OS knowledge (Windows file system, CMD commands) to find the malware in AppData.\n"
    "- They use programming logic to understand how the malware's persistence mechanism works (e.g., a loop checking every 10 minutes).\n"
    "- They use networking knowledge to trace the malware's C2 (command and control) connections over TCP port 443 to a malicious domain.\n"
    "- They use binary/hex knowledge to analyze the malware's file signature and hash.\n"
    "- They use Linux skills to run analysis tools on a Kali Linux VM.\n\n"
    "All these skills work together. Module 1 provides the vocabulary and mental models. Module 2 will apply them to actual attacks and defenses.\n\n"
    "Assessment reminder: Students will be assessed on practical labs (OS installation and configuration) and a quiz covering number systems and networking fundamentals. Encourage them to review the lab exercises and practice binary/hex conversions."
)

# Slide 30: Bridge to Module 2
add_content_slide(prs,
    "What's Next? Module 2 Preview",
    [
        "Now that you understand how computers, OSs, and networks work...",
        "",
        "Module 2: Cybersecurity Principles and Defense",
        "",
        "• We will learn HOW attackers exploit everything we just studied:",
        "   → Port 23 Telnet? That's a vulnerability.",
        "   → XOR operations? That's how malware hides and encryption works.",
        "   → /tmp and C:\\Temp? Malware's favorite hiding spots.",
        "",
        "Coming up:",
        "   • CIA Triad, Threats & Attack Types",
        "   • Network Security: Firewalls, VPNs, Wireless Security",
        "   • Cryptography: Encryption, Hashing, Digital Signatures",
        "   • Authentication, Access Control, IDS/IPS"
    ],
    "INSTRUCTOR NOTES — Cybersecurity Connection:\n"
    "This final slide sets the hook for Module 2. The key message: 'You are no longer a beginner computer user. You now understand the machinery. In Module 2, we start looking at how that machinery gets attacked — and how to defend it.'\n\n"
    "Make explicit connections back to Module 1 topics:\n"
    "- 'Remember the OSI model? In Module 2, you'll learn where firewalls (Layer 3-4), IDS/IPS (Layer 4-7), and WAFs (Layer 7) fit.'\n"
    "- 'Remember binary and XOR? In Module 2, you'll see how symmetric encryption (AES) and hashing (SHA-256) use these exact concepts to protect data.'\n"
    "- 'Remember Windows user accounts? In Module 2, you'll learn about Active Directory, RBAC (Role-Based Access Control), and multi-factor authentication.'\n"
    "- 'Remember TCP ports? In Module 2, you'll configure firewall rules to allow/block specific ports and services.'\n\n"
    "Motivational closing: The cybersecurity industry in Zambia and globally is growing rapidly. MTN Zambia, Airtel, Zanaco, ZESCO, ZICTA, NGOs, and international organizations all need skilled cybersecurity professionals. The foundation you built in Module 1 is the first step toward an in-demand, well-paying career. Keep practicing, stay curious, and remember — the best defenders understand systems as well as attackers do.\n\n"
    "Final activity suggestion: End the session with a quick Kahoot quiz or Jeopardy-style review game covering Module 1 topics. Gamification reinforces learning and identifies areas where students may need additional support before the formal assessment."
)

# Save presentation
output_path = "/Users/michael/Documents/Programming/edutrack-lms/Module_1_Foundations_of_Computing_and_Mathematics.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
