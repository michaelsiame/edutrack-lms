#!/usr/bin/env python3
"""
Create PPTX for Python Practical Exercises
Topic 3: Introduction to Programming Logic
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
PRIMARY_BLUE = RGBColor(46, 112, 218)
SECONDARY_AMBER = RGBColor(246, 183, 69)
DARK_BLUE = RGBColor(30, 74, 138)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(51, 51, 51)
LIGHT_GRAY = RGBColor(248, 249, 250)
CODE_BG = RGBColor(30, 30, 30)
CODE_GREEN = RGBColor(80, 250, 123)
CODE_YELLOW = RGBColor(241, 250, 140)
CODE_BLUE = RGBColor(139, 233, 253)
CODE_PINK = RGBColor(255, 121, 198)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_BLUE
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = SECONDARY_AMBER
    p.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Edutrack Computer Training College - Cybersecurity Certificate Program"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, section_num, section_title):
    """Add section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    
    # Section number
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"PART {section_num}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_AMBER
    p.alignment = PP_ALIGN.CENTER
    
    # Section title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullet_points, code_example=None):
    """Add content slide with bullets and optional code"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = PRIMARY_BLUE
    title_bg.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content area
    if code_example:
        # Split layout: bullets on left, code on right
        content_left = Inches(0.5)
        content_width = Inches(6)
        code_left = Inches(7)
        code_width = Inches(5.833)
    else:
        content_left = Inches(0.5)
        content_width = Inches(12.333)
    
    # Bullet points
    bullet_box = slide.shapes.add_textbox(content_left, Inches(1.5), content_width, Inches(5.5))
    tf = bullet_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)
        p.space_after = Pt(6)
    
    # Code example
    if code_example:
        # Code background
        code_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, code_left, Inches(1.5), code_width, Inches(5.5))
        code_bg.fill.solid()
        code_bg.fill.fore_color.rgb = CODE_BG
        code_bg.line.fill.background()
        
        # Code text
        code_box = slide.shapes.add_textbox(code_left + Inches(0.2), Inches(1.6), code_width - Inches(0.4), Inches(5.3))
        tf = code_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = code_example
        p.font.name = "Consolas"
        p.font.size = Pt(14)
        p.font.color.rgb = CODE_GREEN
        
    return slide

def add_code_slide(prs, title, code, explanation=None):
    """Add a slide focused on code example"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = PRIMARY_BLUE
    title_bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Code background
    code_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(12.333), Inches(4.5))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = CODE_BG
    code_bg.line.fill.background()
    
    # Code text
    code_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.8), Inches(4.2))
    tf = code_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = code
    p.font.name = "Consolas"
    p.font.size = Pt(16)
    p.font.color.rgb = CODE_GREEN
    p.line_spacing = 1.2
    
    # Explanation
    if explanation:
        exp_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(12.333), Inches(1.2))
        tf = exp_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = explanation
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.font.italic = True
    
    return slide

# ============== CREATE SLIDES ==============

# Title Slide
add_title_slide(prs, 
    "Python Programming Practical Exercises",
    "Topic 3: Introduction to Programming Logic\nModule 1: Foundations of Computing and Mathematics")

# Learning Objectives
add_content_slide(prs, "Learning Objectives", [
    "Write and execute basic Python commands",
    "Work with variables and different data types (strings, integers, floats, booleans)",
    "Use operators for calculations and comparisons",
    "Implement control flow with conditionals and loops",
    "Create functions to organize code",
    "Apply Python basics to cybersecurity scenarios"
])

# Section 1: Getting Started
add_section_slide(prs, "1", "Getting Started with Python")

add_content_slide(prs, "Your First Python Program", [
    "Python uses the print() function to display output",
    "Strings can be enclosed in single or double quotes",
    "Variables store data for later use",
    "f-strings (formatted strings) make it easy to include variables in output"
], '''#!/usr/bin/env python3
# Hello World for Cybersecurity

print("Welcome to Python!")
print("-" * 40)

# Variables
student_name = "Your Name"
course = "Cybersecurity"

# f-string formatting
print(f"Hello, {student_name}!")
print(f"Course: {course}")''')

add_content_slide(prs, "Python as a Calculator", [
    "Arithmetic operators: +, -, *, /, //, %, **",
    "Order of operations follows standard math rules",
    "Parentheses can be used to group operations",
    "Useful for network calculations and cryptography"
], '''# Network calculations
ips_per_subnet = 256
subnets = 4
total = ips_per_subnet * subnets
print(f"Total IPs: {total}")

# Password cracking time calculation
password_space = 95 ** 8  # 8 chars, 95 printable
attempts_per_sec = 1_000_000
seconds = password_space / attempts_per_sec
print(f"Time to crack: {seconds} seconds")''')

# Section 2: Variables and Data Types
add_section_slide(prs, "2", "Variables and Data Types")

add_content_slide(prs, "Understanding Data Types", [
    "Strings (str): Text data - usernames, emails, hashes",
    "Integers (int): Whole numbers - port numbers, counts",
    "Floats (float): Decimal numbers - percentages, measurements",
    "Booleans (bool): True/False values - access control",
    "None: Represents no value or null"
], '''# Security data examples
username = "admin"              # String
port = 443                     # Integer
risk_score = 7.5               # Float
is_authenticated = True        # Boolean
result = None                  # NoneType

print(f"User: {username}")
print(f"HTTPS Port: {port}")
print(f"Risk Score: {risk_score}")''')

add_content_slide(prs, "Lists and Dictionaries", [
    "Lists: Ordered collections of items (mutable)",
    "Use index [0] to access items (starts at 0)",
    "Dictionaries: Key-value pairs for structured data",
    "Use key ['name'] to access values in dictionaries"
], '''# List of blocked IPs
blocked_ips = ["192.168.1.100", 
               "10.0.0.50"]
blocked_ips.append("172.16.0.25")

# Dictionary for user account
user = {
    "username": "jdoe",
    "role": "analyst",
    "active": True
}

print(f"First blocked IP: {blocked_ips[0]}")
print(f"User role: {user['role']}")''')

# Section 3: Operators
add_section_slide(prs, "3", "Operators in Python")

add_content_slide(prs, "Arithmetic and Comparison Operators", [
    "Arithmetic: +, -, *, /, // (floor division), % (modulo), ** (power)",
    "Comparison: == (equal), != (not equal), >, <, >=, <=",
    "Result of comparisons is always True or False",
    "Modulo (%) is useful for hashing and checksums"
], '''# Packet loss calculation
packets_in = 1500
packets_out = 1420
packet_loss = packets_in - packets_out
loss_percent = (packet_loss / packets_in) * 100
print(f"Loss: {loss_percent:.2f}%")

# Comparison
failed_logins = 5
threshold = 3
is_locked = failed_logins > threshold
print(f"Account locked: {is_locked}")''')

add_content_slide(prs, "Logical Operators", [
    "and: Both conditions must be True",
    "or: At least one condition must be True",
    "not: Inverts the condition (True becomes False)",
    "Used for complex access control decisions"
], '''# Access control logic
is_business_hours = True
is_weekend = False
has_admin_priv = True

# AND - both must be true
allow_access = is_business_hours and not is_weekend

# OR - at least one must be true
needs_review = failed_logins > 3 or not is_business_hours

# NOT - invert
is_blocked = not allow_access

print(f"Allow access: {allow_access}")''')

# Section 4: Control Flow
add_section_slide(prs, "4", "Control Flow - Conditionals")

add_code_slide(prs, "If-Elif-Else Statements", '''# Risk level assessment
cvss_score = 7.5

if cvss_score >= 9.0:
    risk_level = "CRITICAL"
    action = "Patch within 24 hours"
elif cvss_score >= 7.0:
    risk_level = "HIGH"
    action = "Patch within 1 week"
elif cvss_score >= 4.0:
    risk_level = "MEDIUM"
    action = "Patch within 1 month"
else:
    risk_level = "LOW"
    action = "Schedule for next maintenance"

print(f"Risk: {risk_level}")
print(f"Action: {action}")''', 
"The if-elif-else chain checks conditions in order and executes the first matching block")

add_code_slide(prs, "Nested Conditionals", '''# Access control with multiple checks
user_role = "analyst"
is_authenticated = True
security_clearance = 2
required_clearance = 3

if is_authenticated:
    print("User is authenticated")
    
    if user_role == "admin":
        print("Full system access granted")
    elif user_role == "analyst":
        if security_clearance >= required_clearance:
            print("Access granted to classified data")
        else:
            print("Access denied - Insufficient clearance")
    else:
        print("Read-only access granted")
else:
    print("Access denied - Please log in")''',
"Nested if statements allow for multi-level decision making")

# Section 5: Loops
add_section_slide(prs, "5", "Loops - Repeating Actions")

add_code_slide(prs, "For Loops", '''# Check multiple ports
ports_to_check = [22, 80, 443, 3306, 8080]

for port in ports_to_check:
    # Simulate port check
    if port in [80, 443]:
        status = "OPEN"
    else:
        status = "CLOSED"
    print(f"Port {port}: {status}")

# Using enumerate for index
passwords = ["weak", "Password1", "Str0ng!"]
for i, pwd in enumerate(passwords, 1):
    print(f"{i}. Checking password...")''',
"For loops iterate over a sequence (list, string, range, etc.)")

add_code_slide(prs, "While Loops", '''# Login attempt simulation
max_attempts = 3
attempts = 0
password_correct = False

while attempts < max_attempts and not password_correct:
    attempts += 1
    print(f"Attempt {attempts} of {max_attempts}")
    
    # Simulate password check
    if attempts == 2:
        password_correct = True
        print("  ✓ Login successful!")
    else:
        print("  ✗ Invalid password")

if not password_correct:
    print("⚠️  Account locked")''',
"While loops continue as long as a condition is True")

# Section 6: Functions
add_section_slide(prs, "6", "Functions - Organizing Code")

add_code_slide(prs, "Defining Functions", '''# Function with parameters and return value
def check_password_strength(password, min_length=8):
    """Check password meets security requirements"""
    score = 0
    
    if len(password) >= min_length:
        score += 1
    if any(c.isdigit() for c in password):
        score += 1
    if any(c.isupper() for c in password):
        score += 1
    
    is_strong = score >= 3
    return is_strong, f"Score: {score}/4"

# Using the function
strong, msg = check_password_strength("MyP@ssw0rd")
print(f"Strong: {strong}, {msg}")''',
"Functions encapsulate reusable logic with parameters and return values")

add_code_slide(prs, "Function with Multiple Returns", '''# Analyze IP address
def analyze_ip(ip_string):
    """Analyze IP and return properties"""
    octets = [int(o) for o in ip_string.split(".")]
    
    # Determine IP class
    first = octets[0]
    if first < 128:
        ip_class = "A"
    elif first < 192:
        ip_class = "B"
    elif first < 224:
        ip_class = "C"
    else:
        ip_class = "D"
    
    # Check if private
    is_private = (first == 10 or 
                  (first == 172 and 16 <= octets[1] <= 31))
    
    return {
        "ip": ip_string,
        "class": ip_class,
        "is_private": is_private
    }

result = analyze_ip("192.168.1.1")
print(f"IP: {result['ip']}, Class: {result['class']}")''',
"Functions can return dictionaries with multiple related values")

# Section 7: Practical Scenarios
add_section_slide(prs, "7", "Practical Cybersecurity Scenarios")

add_code_slide(prs, "Scenario: Log File Analyzer", '''def analyze_auth_log(log_entries):
    """Analyze logs for brute force attempts"""
    stats = {
        "total": len(log_entries),
        "failed": 0,
        "unique_ips": set()
    }
    failed_by_ip = {}
    
    for entry in log_entries:
        stats["unique_ips"].add(entry["ip"])
        if entry["status"] == "failed":
            stats["failed"] += 1
            failed_by_ip[entry["ip"]] = \
                failed_by_ip.get(entry["ip"], 0) + 1
    
    # Detect brute force (5+ failed attempts)
    for ip, count in failed_by_ip.items():
        if count >= 5:
            print(f"🚨 Brute force from {ip}: {count} attempts")
    
    return stats''',
"Real-world example: Analyzing authentication logs for security threats")

add_code_slide(prs, "Scenario: Password Generator", '''import random
import string

def generate_password(length=12):
    """Generate secure random password"""
    # Character sets
    lowercase = string.ascii_lowercase
    uppercase = string.ascii_uppercase
    digits = string.digits
    special = "!@#$%^&*"
    
    # Ensure at least one of each type
    password = [
        random.choice(lowercase),
        random.choice(uppercase),
        random.choice(digits),
        random.choice(special)
    ]
    
    # Fill remaining with random choices
    all_chars = lowercase + uppercase + digits + special
    for _ in range(length - 4):
        password.append(random.choice(all_chars))
    
    # Shuffle and return
    random.shuffle(password)
    return ''.join(password)

print(f"Generated: {generate_password(16)}")''',
"Security tool: Generating strong passwords with guaranteed character types")

# Summary Slide
add_content_slide(prs, "Summary - What You Learned", [
    "✓ Basic Python syntax: print(), variables, comments",
    "✓ Data types: strings, integers, floats, booleans, lists, dictionaries",
    "✓ Operators: arithmetic, comparison, logical",
    "✓ Control flow: if-elif-else for decision making",
    "✓ Loops: for and while for repeating actions",
    "✓ Functions: creating reusable code blocks",
    "✓ Practical applications: log analysis, password generation",
    "",
    "Next Steps:",
    "• Complete the exercises in the python_exercises folder",
    "• Practice writing your own scripts",
    "• Apply Python to automate security tasks"
])

# Resources Slide
add_content_slide(prs, "Resources and Next Steps", [
    "Exercise Files: course_materials/module1_foundation/topic3_programming/python_exercises/",
    "",
    "Files included:",
    "• 00_cheat_sheet.py - Quick reference for all concepts",
    "• 01_hello_security.py to 07_functions.py - Practice exercises",
    "• 08_log_analyzer.py - Real-world log analysis tool",
    "• 09_password_generator.py - Security tool example",
    "• 10_firewall_simulator.py - Network security simulation",
    "• 99_mini_project.py - Capstone exercise",
    "",
    "Online Resources:",
    "• Python Official Tutorial: docs.python.org/3/tutorial/",
    "• Practice: hackerrank.com/domains/python",
    "• Cybersecurity Python: github.com/topics/cybersecurity-python"
])

# Save presentation
output_path = "Python_Practical_Exercises.pptx"
prs.save(output_path)
print(f"Presentation created: {output_path}")
print(f"Total slides: {len(prs.slides)}")
